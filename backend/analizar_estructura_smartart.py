#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Analiza la estructura XML de los SmartArt para entender cómo acceder a ellos
"""
import os
from pptx import Presentation
from lxml import etree

TEMPLATE_PATH = "../Template/Template-PreCotizacion.pptx"

def analyze_smartart_structure():
    """Analiza la estructura interna de los SmartArt"""
    if not os.path.exists(TEMPLATE_PATH):
        print(f"❌ Template no encontrado: {TEMPLATE_PATH}")
        return
    
    print("=" * 70)
    print("🔬 ANÁLISIS DE ESTRUCTURA SMARTART")
    print("=" * 70)
    
    prs = Presentation(TEMPLATE_PATH)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 DIAPOSITIVA {slide_idx + 1}")
        print("-" * 70)
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_name = getattr(shape, "name", "Sin nombre")
            shape_type = shape.shape_type
            
            # Tipo 3 = GROUP, Tipo 15 = GRAPHIC_FRAME (SmartArt)
            if shape_type in [3, 15]:
                print(f"\n   Shape {shape_idx}: {shape_name}")
                print(f"   Tipo: {shape_type} ({get_shape_type_name(shape_type)})")
                
                # Intentar acceder a shapes internos si es un grupo
                if hasattr(shape, 'shapes'):
                    print(f"   ✅ Tiene shapes internos: {len(shape.shapes)}")
                    for sub_idx, sub_shape in enumerate(shape.shapes):
                        sub_name = getattr(sub_shape, "name", "Sin nombre")
                        print(f"      └─ Sub-shape {sub_idx}: {sub_name}")
                        if sub_shape.has_text_frame:
                            text = sub_shape.text
                            if "{{" in text:
                                print(f"         📝 Texto: {text[:100]}")
                
                # Buscar en el XML
                try:
                    xml_str = etree.tostring(shape.element, encoding='unicode', pretty_print=True)
                    
                    # Buscar placeholders en el XML
                    if "{{" in xml_str:
                        print(f"   ✅ Contiene placeholders en XML")
                        
                        # Extraer fragmento con placeholder
                        lines = xml_str.split('\n')
                        for i, line in enumerate(lines):
                            if "{{" in line:
                                # Mostrar contexto (3 líneas antes y después)
                                start = max(0, i-2)
                                end = min(len(lines), i+3)
                                print(f"\n   📋 Fragmento XML (línea {i}):")
                                for j in range(start, end):
                                    marker = ">>> " if j == i else "    "
                                    print(f"      {marker}{lines[j].strip()[:80]}")
                                
                                # Solo mostrar el primer placeholder de cada shape
                                break
                    else:
                        print(f"   ⚠️  NO contiene placeholders en XML")
                        
                except Exception as e:
                    print(f"   ❌ Error analizando XML: {e}")
    
    print("\n" + "=" * 70)

def get_shape_type_name(shape_type):
    """Convierte el número de tipo a nombre"""
    types = {
        1: "AUTO_SHAPE",
        3: "GROUP",
        13: "PICTURE",
        14: "PLACEHOLDER",
        15: "GRAPHIC_FRAME (SmartArt/Chart)",
        17: "TEXT_BOX",
        19: "TABLE"
    }
    return types.get(shape_type, f"Unknown ({shape_type})")

if __name__ == "__main__":
    analyze_smartart_structure()
