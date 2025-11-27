#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Diagnóstico mejorado del Template PPTX - Incluye SmartArt
"""
import os
import re
from pptx import Presentation

TEMPLATE_PATH = "../Template/Template-PreCotizacion.pptx"

def find_text_in_smartart(shape):
    """Busca texto dentro de SmartArt y otros objetos complejos usando XML"""
    placeholders = []
    try:
        # Buscar TODOS los elementos <a:t> en el shape completo
        # Esto incluye SmartArt, gráficos, formas anidadas, etc.
        namespaces = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        }
        
        for text_elem in shape.element.findall('.//a:t', namespaces):
            if text_elem.text:
                matches = re.findall(r'\{\{[^}]+\}\}', text_elem.text)
                for match in matches:
                    placeholders.append({
                        'placeholder': match,
                        'full_text': text_elem.text,
                        'type': 'SmartArt/XML'
                    })
    except Exception as e:
        pass
    return placeholders

def analyze_template():
    """Analiza el template PPTX incluyendo SmartArt"""
    if not os.path.exists(TEMPLATE_PATH):
        print(f"❌ Template no encontrado: {TEMPLATE_PATH}")
        return
    
    print("=" * 70)
    print("🔍 DIAGNÓSTICO COMPLETO DEL TEMPLATE PPTX (con SmartArt)")
    print("=" * 70)
    
    prs = Presentation(TEMPLATE_PATH)
    
    # Buscar placeholders
    placeholders_found = {}
    smartart_count = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            shape_name = getattr(shape, "name", "Sin nombre")
            shape_type = shape.shape_type
            
            # Contar SmartArt
            if shape_type == 15:  # GRAPHIC_FRAME (incluye SmartArt)
                smartart_count += 1
            
            # Buscar en text_frame
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = run.text
                        if "{{" in text and "}}" in text:
                            matches = re.findall(r'\{\{[^}]+\}\}', text)
                            for ph in matches:
                                if ph not in placeholders_found:
                                    placeholders_found[ph] = []
                                placeholders_found[ph].append({
                                    "slide": slide_idx + 1,
                                    "shape_name": shape_name,
                                    "type": "TextFrame",
                                    "full_text": text
                                })
            
            # Buscar en tablas
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        text = cell.text
                        if "{{" in text and "}}" in text:
                            matches = re.findall(r'\{\{[^}]+\}\}', text)
                            for ph in matches:
                                if ph not in placeholders_found:
                                    placeholders_found[ph] = []
                                placeholders_found[ph].append({
                                    "slide": slide_idx + 1,
                                    "shape_name": shape_name,
                                    "type": f"Table (Fila {row_idx}, Col {col_idx})",
                                    "full_text": text
                                })
            
            # Buscar en SmartArt
            smartart_placeholders = find_text_in_smartart(shape)
            for item in smartart_placeholders:
                ph = item['placeholder']
                if ph not in placeholders_found:
                    placeholders_found[ph] = []
                placeholders_found[ph].append({
                    "slide": slide_idx + 1,
                    "shape_name": shape_name,
                    "type": "SmartArt/XML",
                    "full_text": item['full_text']
                })
    
    print(f"\n📊 ESTADÍSTICAS:")
    print(f"   Diapositivas: {len(prs.slides)}")
    print(f"   SmartArt detectados: {smartart_count}")
    print(f"   Placeholders únicos: {len(placeholders_found)}")
    
    if placeholders_found:
        print(f"\n🏷️  PLACEHOLDERS ENCONTRADOS ({len(placeholders_found)}):\n")
        for ph, locations in sorted(placeholders_found.items()):
            print(f"   {ph}")
            for loc in locations:
                slide = loc['slide']
                shape_name = loc['shape_name']
                loc_type = loc['type']
                print(f"      └─ Diapositiva {slide}, Shape '{shape_name}', Tipo: {loc_type}")
                if len(loc['full_text']) < 80:
                    print(f"         Texto: \"{loc['full_text']}\"")
    else:
        print("   ⚠️  No se encontraron placeholders {{...}}")
    
    print("\n" + "=" * 70)
    print(f"✅ Si ves tus placeholders listados arriba con 'Tipo: SmartArt',")
    print(f"   entonces el código actualizado podrá reemplazarlos correctamente.")
    print("=" * 70)

if __name__ == "__main__":
    analyze_template()
