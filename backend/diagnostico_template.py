#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Diagnóstico del Template PPTX
- Lista todos los placeholders encontrados
- Verifica dimensiones de página
- Busca texto en todos los shapes
"""
import os
from pptx import Presentation
from pptx.util import Inches

TEMPLATE_PATH = "../Template/Template-PreCotizacion.pptx"

def analyze_template():
    """Analiza el template PPTX"""
    if not os.path.exists(TEMPLATE_PATH):
        print(f"❌ Template no encontrado: {TEMPLATE_PATH}")
        return
    
    print("=" * 70)
    print("🔍 DIAGNÓSTICO DEL TEMPLATE PPTX")
    print("=" * 70)
    
    prs = Presentation(TEMPLATE_PATH)
    
    # 1. DIMENSIONES DE PÁGINA
    print(f"\n📐 DIMENSIONES DE PÁGINA:")
    width_inches = prs.slide_width / Inches(1)
    height_inches = prs.slide_height / Inches(1)
    print(f"   Ancho:  {prs.slide_width} EMUs = {width_inches:.2f} pulgadas")
    print(f"   Alto:   {prs.slide_height} EMUs = {height_inches:.2f} pulgadas")
    
    # Tamaños estándar
    if abs(width_inches - 10) < 0.1 and abs(height_inches - 7.5) < 0.1:
        print(f"   ✅ Formato: Presentación estándar (10x7.5 pulgadas)")
    elif abs(width_inches - 8.5) < 0.1 and abs(height_inches - 11) < 0.1:
        print(f"   ✅ Formato: Carta vertical (8.5x11 pulgadas)")
    elif abs(width_inches - 11) < 0.1 and abs(height_inches - 8.5) < 0.1:
        print(f"   ✅ Formato: Carta horizontal (11x8.5 pulgadas)")
    else:
        print(f"   ⚠️  Formato personalizado")
    
    # 2. BUSCAR PLACEHOLDERS
    print(f"\n🏷️  PLACEHOLDERS ENCONTRADOS:")
    placeholders_found = {}
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # Buscar en text_frame
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = run.text
                        if "{{" in text and "}}" in text:
                            # Extraer placeholders
                            import re
                            found = re.findall(r'\{\{[^}]+\}\}', text)
                            for ph in found:
                                if ph not in placeholders_found:
                                    placeholders_found[ph] = []
                                placeholders_found[ph].append({
                                    "slide": slide_idx + 1,
                                    "shape": shape_idx,
                                    "shape_name": getattr(shape, "name", "Sin nombre"),
                                    "full_text": text
                                })
            
            # Buscar en tablas
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        text = cell.text
                        if "{{" in text and "}}" in text:
                            import re
                            found = re.findall(r'\{\{[^}]+\}\}', text)
                            for ph in found:
                                if ph not in placeholders_found:
                                    placeholders_found[ph] = []
                                placeholders_found[ph].append({
                                    "slide": slide_idx + 1,
                                    "shape": shape_idx,
                                    "shape_name": getattr(shape, "name", "Sin nombre"),
                                    "table_cell": f"Fila {row_idx}, Col {col_idx}",
                                    "full_text": text
                                })
    
    if placeholders_found:
        print(f"\n   ✅ Encontrados {len(placeholders_found)} placeholders únicos:\n")
        for ph, locations in sorted(placeholders_found.items()):
            print(f"   {ph}")
            for loc in locations:
                slide = loc['slide']
                shape_name = loc['shape_name']
                if 'table_cell' in loc:
                    print(f"      └─ Diapositiva {slide}, Shape '{shape_name}', {loc['table_cell']}")
                else:
                    print(f"      └─ Diapositiva {slide}, Shape '{shape_name}'")
                    if len(loc['full_text']) < 100:
                        print(f"         Texto: \"{loc['full_text']}\"")
    else:
        print("   ⚠️  No se encontraron placeholders {{...}}")
    
    # 3. BUSCAR TABLA DE AHORROS
    print(f"\n📊 TABLA DE AHORROS:")
    tabla_encontrada = False
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if getattr(shape, "name", "") == "TABLA_AHORROS" and shape.has_table:
                print(f"   ✅ Encontrada tabla 'TABLA_AHORROS' en diapositiva {slide_idx + 1}")
                print(f"      Filas: {len(shape.table.rows)}, Columnas: {len(shape.table.columns)}")
                # Mostrar encabezados
                if shape.table.rows:
                    headers = [cell.text for cell in shape.table.rows[0].cells]
                    print(f"      Encabezados: {headers}")
                tabla_encontrada = True
    
    if not tabla_encontrada:
        print("   ⚠️  No se encontró tabla con nombre 'TABLA_AHORROS'")
        print("   🔍 Buscando tablas con encabezados 'año' y 'ahorro'...")
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    if table.rows:
                        headers = [cell.text.lower() for cell in table.rows[0].cells]
                        if any("año" in h or "ano" in h for h in headers) and any("ahorro" in h for h in headers):
                            print(f"   ✅ Posible tabla encontrada en diapositiva {slide_idx + 1}")
                            print(f"      Shape: {getattr(shape, 'name', 'Sin nombre')}")
                            print(f"      Filas: {len(table.rows)}, Columnas: {len(table.columns)}")
                            print(f"      Encabezados: {[cell.text for cell in table.rows[0].cells]}")
    
    # 4. TOTAL DE DIAPOSITIVAS
    print(f"\n📄 TOTAL DE DIAPOSITIVAS: {len(prs.slides)}")
    
    print("\n" + "=" * 70)

if __name__ == "__main__":
    analyze_template()
