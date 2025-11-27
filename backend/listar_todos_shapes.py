#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Lista TODOS los shapes del template con sus tipos y texto
"""
import os
from pptx import Presentation

TEMPLATE_PATH = "../Template/Template-PreCotizacion.pptx"

def list_all_shapes():
    """Lista todos los shapes con detalles"""
    if not os.path.exists(TEMPLATE_PATH):
        print(f"❌ Template no encontrado: {TEMPLATE_PATH}")
        return
    
    print("=" * 70)
    print("📋 LISTA COMPLETA DE SHAPES")
    print("=" * 70)
    
    prs = Presentation(TEMPLATE_PATH)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 DIAPOSITIVA {slide_idx + 1}")
        print("-" * 70)
        
        if not slide.shapes:
            print("   (Sin shapes)")
            continue
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_name = getattr(shape, "name", "Sin nombre")
            shape_type = shape.shape_type
            
            print(f"\n   [{shape_idx}] {shape_name}")
            print(f"   Tipo: {shape_type}")
            
            # Ver si tiene texto
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    # Mostrar primeros 100 caracteres
                    display_text = text[:100] + "..." if len(text) > 100 else text
                    print(f"   Texto: {display_text}")
                    
                    # Marcar si tiene placeholders
                    if "{{" in text:
                        print(f"   ✅ TIENE PLACEHOLDERS")
            
            # Ver si tiene tabla
            if shape.has_table:
                table = shape.table
                print(f"   Tabla: {len(table.rows)} filas x {len(table.columns)} columnas")
                
                # Ver si la tabla tiene placeholders
                has_placeholder = False
                for row in table.rows:
                    for cell in row.cells:
                        if "{{" in cell.text:
                            has_placeholder = True
                            break
                    if has_placeholder:
                        break
                
                if has_placeholder:
                    print(f"   ✅ TIENE PLACEHOLDERS EN TABLA")
            
            # Ver si es un grupo
            if hasattr(shape, 'shapes'):
                print(f"   Grupo: {len(shape.shapes)} sub-shapes")
                for sub_idx, sub_shape in enumerate(shape.shapes[:5]):  # Máximo 5
                    sub_name = getattr(sub_shape, "name", "Sin nombre")
                    sub_text = ""
                    if sub_shape.has_text_frame:
                        sub_text = sub_shape.text.strip()[:50]
                    print(f"      └─ [{sub_idx}] {sub_name}: {sub_text}")
            
            # Buscar en XML para SmartArt
            try:
                # Buscar elementos <a:t> en el XML
                namespaces = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                text_elems = shape.element.findall('.//a:t', namespaces)
                
                if text_elems:
                    print(f"   Elementos <a:t> en XML: {len(text_elems)}")
                    
                    # Ver si alguno tiene placeholders
                    for elem in text_elems[:3]:  # Máximo 3
                        if elem.text and "{{" in elem.text:
                            print(f"      ✅ XML con placeholder: {elem.text[:60]}")
            except:
                pass
    
    print("\n" + "=" * 70)

if __name__ == "__main__":
    list_all_shapes()
