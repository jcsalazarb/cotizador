import os
import json
import time
import socket
import tempfile
import shutil
import subprocess
import smtplib
import secrets
import unicodedata
from math import ceil
from datetime import datetime, timezone, timedelta
from typing import Optional, Dict, Any
from fastapi import FastAPI, HTTPException, Depends, Request
from fastapi.responses import JSONResponse, FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBasic, HTTPBasicCredentials
from pydantic import BaseModel, EmailStr, Field, field_validator, ValidationError
from fastapi.exceptions import RequestValidationError
from email.message import EmailMessage
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

# Importar modelos de SQLAlchemy
try:
    from models import get_db_session, Panel, Inversor, Bateria, Ciudad, Parametro, Consecutivo, Estadistica
    POSTGRES_AVAILABLE = True
except ImportError as e:
    print(f"⚠️ Warning: No se pudieron importar modelos de PostgreSQL: {e}")
    POSTGRES_AVAILABLE = False

load_dotenv()

# ========================================
# ⏰ ZONA HORARIA COLOMBIA (UTC-5)
# ========================================
COLOMBIA_TZ = timezone(timedelta(hours=-5))

def now_colombia():
    """Retorna la hora actual en Colombia (UTC-5)"""
    return datetime.now(COLOMBIA_TZ)

# ========================================
# 📁 CONFIGURACIÓN DE RUTAS
# ========================================
APP_DIR = os.path.dirname(os.path.abspath(__file__))
CONFIG_DIR = os.path.join(APP_DIR, "config")
EQUIPOS_FILE = os.path.join(CONFIG_DIR, "equipos.json")
CIUDADES_FILE = os.path.join(CONFIG_DIR, "ciudades.json")
PARAMETROS_FILE = os.path.join(CONFIG_DIR, "parametros.json")
ESTADISTICAS_FILE = os.path.join(CONFIG_DIR, "estadisticas.json")
CONSECUTIVO_FILE = os.path.join(CONFIG_DIR, "consecutivo.json")
TEMPLATE_DIR = os.path.join(APP_DIR, "..", "Template")
TEMPLATE_PPTX = os.path.join(TEMPLATE_DIR, "Template-PreCotizacion.pptx")
TEMPLATE_PPTX_OP2 = os.path.join(TEMPLATE_DIR, "Template-PreCotizacion2.pptx")

# ========================================
# 🔒 CONFIGURACIÓN DE SEGURIDAD
# ========================================
app = FastAPI(
    title="NASSA Solar API",
    description="API segura para cotizaciones de sistemas fotovoltaicos",
    version="1.0.0",
    docs_url="/docs",
    redoc_url="/redoc"
)

ALLOWED_ORIGINS = os.getenv("ALLOWED_ORIGINS", "*").split(",")
app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["*"],
)

# ========================================
# �️ MANEJADOR DE ERRORES DE VALIDACIÓN
# ========================================
@app.exception_handler(RequestValidationError)
async def validation_exception_handler(request: Request, exc: RequestValidationError):
    """
    Captura errores de validación de Pydantic y devuelve mensajes amigables en español
    """
    errores_amigables = []
    
    # Diccionario de traducción de campos
    campo_nombres = {
        "nombre": "Nombre",
        "telefono": "Teléfono",
        "email": "Correo electrónico",
        "ciudad": "Ciudad",
        "direccion": "Dirección",
        "consumoMensual": "Consumo mensual",
        "valorFactura": "Valor de la factura",
        "valorKwh": "Valor kWh",
        "nic": "NIC/Número de instalación",
        "tipoVivienda": "Tipo de vivienda",
        "areaDisponible": "Área disponible",
        "numeroPisos": "Número de pisos",
        "sistemaElectrico": "Sistema eléctrico",
        "porcentajeConsumodia": "Porcentaje consumo día",
        "porcentajeAhorroEnergia": "Porcentaje de ahorro energético",
        "tipoSistemaFV": "Tipo de sistema fotovoltaico",
        "legalizacion": "Legalización",
        "seleccionManual": "Selección manual",
        "panel": "Panel",
        "inversor": "Inversor",
        "bateria": "Batería",
        "identificacion": "Identificación"
    }
    
    for error in exc.errors():
        campo = error["loc"][-1] if error["loc"] else "campo"
        tipo_error = error["type"]
        mensaje_original = error["msg"]
        
        # Obtener nombre amigable del campo
        nombre_campo = campo_nombres.get(campo, campo)
        
        # Generar mensaje amigable según el tipo de error
        if tipo_error == "string_too_short":
            min_length = error.get("ctx", {}).get("min_length", "")
            valor_actual = error.get("input", "")
            mensaje = f"{nombre_campo}: debe tener al menos {min_length} caracteres (actualmente tiene {len(str(valor_actual))})"
        
        elif tipo_error == "string_too_long":
            max_length = error.get("ctx", {}).get("max_length", "")
            mensaje = f"{nombre_campo}: no puede exceder {max_length} caracteres"
        
        elif tipo_error == "string_pattern_mismatch":
            if campo == "telefono":
                mensaje = f"{nombre_campo}: debe tener un formato válido (7-20 dígitos, puede incluir +, espacios, guiones o paréntesis)"
            elif campo == "tipoVivienda":
                mensaje = f"{nombre_campo}: debe ser 'casa', 'apartamento', 'local' o 'empresa'"
            elif campo == "sistemaElectrico":
                mensaje = f"{nombre_campo}: debe ser 'monofasico', 'bifasico' o 'trifasico'"
            elif campo == "tipoSistemaFV":
                mensaje = f"{nombre_campo}: debe ser 'ongrid', 'offgrid', 'hibrido_incluido' o 'hibrido_opcional'"
            elif campo in ["legalizacion", "seleccionManual"]:
                mensaje = f"{nombre_campo}: debe ser 'SI' o 'NO'"
            else:
                mensaje = f"{nombre_campo}: formato inválido"
        
        elif tipo_error == "value_error.email":
            mensaje = f"{nombre_campo}: debe ser un correo electrónico válido"
        
        elif tipo_error == "greater_than":
            limite = error.get("ctx", {}).get("gt", "")
            mensaje = f"{nombre_campo}: debe ser mayor a {limite}"
        
        elif tipo_error == "less_than":
            limite = error.get("ctx", {}).get("lt", "")
            mensaje = f"{nombre_campo}: debe ser menor a {limite}"
        
        elif tipo_error == "greater_than_equal":
            limite = error.get("ctx", {}).get("ge", "")
            mensaje = f"{nombre_campo}: debe ser mayor o igual a {limite}"
        
        elif tipo_error == "less_than_equal":
            limite = error.get("ctx", {}).get("le", "")
            mensaje = f"{nombre_campo}: debe ser menor o igual a {limite}"
        
        elif tipo_error == "missing":
            mensaje = f"{nombre_campo}: es obligatorio"
        
        elif tipo_error == "value_error":
            # Errores personalizados de validadores
            mensaje = f"{nombre_campo}: {mensaje_original}"
        
        else:
            # Fallback para otros tipos de error
            mensaje = f"{nombre_campo}: {mensaje_original}"
        
        errores_amigables.append({
            "campo": campo,
            "mensaje": mensaje
        })
    
    return JSONResponse(
        status_code=422,
        content={
            "error": "Errores de validación",
            "mensaje": "Por favor, corrige los siguientes campos:",
            "errores": errores_amigables
        }
    )

# ========================================
# �📁 CONFIGURACIÓN DE ARCHIVOS ESTÁTICOS
# ========================================
STATIC_DIR = os.path.join(APP_DIR, "static")
if os.path.exists(STATIC_DIR):
    app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")

# Servir index.html desde la raíz
@app.get("/")
async def root():
    """Redirige a la página principal del cotizador"""
    index_path = os.path.join(STATIC_DIR, "index.html")
    if os.path.exists(index_path):
        return FileResponse(index_path)
    return {"message": "NASSA Solar - API de Cotización", "docs": "/docs"}

security = HTTPBasic()
RATE_LIMIT = int(os.getenv("RATE_LIMIT", "10"))
RATE_WINDOW = 60
request_counts: Dict[str, list] = {}

def rate_limit(request: Request):
    """Rate limiting por IP"""
    ip = request.client.host
    now = time.time()
    arr = request_counts.setdefault(ip, [])
    request_counts[ip] = [t for t in arr if now - t < RATE_WINDOW]
    if len(request_counts[ip]) >= RATE_LIMIT:
        raise HTTPException(429, f"Límite {RATE_LIMIT} req/min excedido")
    request_counts[ip].append(now)

def auth_admin(credentials: HTTPBasicCredentials = Depends(security)):
    """Verificar credenciales para endpoints administrativos"""
    u_ok = secrets.compare_digest(credentials.username, os.getenv("ADMIN_USER", "admin"))
    p_ok = secrets.compare_digest(credentials.password, os.getenv("ADMIN_PASS", "changeme"))
    if not (u_ok and p_ok):
        raise HTTPException(
            status_code=401,
            detail="Credenciales inválidas",
            headers={"WWW-Authenticate": "Basic"}
        )
    return credentials.username

# ========================================
# 🔐 MODELOS CON VALIDACIÓN ROBUSTA
# ========================================
class CotizarRequest(BaseModel):
    nombre: str = Field(..., min_length=3, max_length=100)
    telefono: str = Field(..., pattern=r'^\+?[0-9\s\-()]{7,20}$')
    email: EmailStr
    ciudad: str = Field(..., min_length=2, max_length=80)
    direccion: str = Field(..., min_length=5, max_length=200)
    consumoMensual: float = Field(..., gt=50, lt=50000)
    valorFactura: float = Field(..., gt=10000, lt=100000000)
    valorKwh: float = Field(..., gt=100, lt=5000)
    nic: str = Field(..., min_length=3, max_length=25)
    tipoVivienda: str = Field(..., pattern=r'^(casa|apartamento|local|empresa)$')
    areaDisponible: Optional[float] = Field(0, ge=0, lt=20000)
    numeroPisos: Optional[str] = Field("1", pattern=r'^[1-6]$')
    sistemaElectrico: str = Field(..., pattern=r'^(monofasico|bifasico|trifasico)$')
    porcentajeConsumodia: float = Field(..., ge=0, le=100)
    porcentajeAhorroEnergia: float = Field(100.0, ge=10, le=100, description="Porcentaje del consumo que desea cubrir con energía solar")
    tipoSistemaFV: str = Field(..., pattern=r'^(ongrid|offgrid|hibrido_incluido|hibrido_opcional)$')
    hspCalculado: Optional[float] = Field(None, gt=0, lt=9)
    legalizacion: str = Field(..., pattern=r'^(SI|NO)$')
    seleccionManual: str = Field(..., pattern=r'^(SI|NO)$')
    panel: Optional[str] = Field(None, pattern=r'^panel[1-9]\d?$')
    inversor: Optional[str] = Field(None, pattern=r'^inv[1-9]\d?$')
    bateria: Optional[str] = Field(None, pattern=r'^bat[1-9]\d?$')
    identificacion: Optional[str] = Field(None, max_length=20)

    @field_validator("bateria")
    @classmethod
    def validar_bateria(cls, v, info):
        if info.data.get("tipoSistemaFV") in ("offgrid", "hibrido_incluido") and not v:
            raise ValueError("Batería requerida para el tipo de sistema seleccionado")
        return v
    
    @field_validator("panel", "inversor")
    @classmethod
    def validar_equipos_manuales(cls, v, info):
        # Si seleccionManual es SI, panel e inversor son obligatorios
        if info.data.get("seleccionManual") == "SI" and not v:
            raise ValueError("Panel e inversor son requeridos cuando se selecciona manualmente")
        return v

# ========================================
# 📁 FUNCIONES AUXILIARES
# ========================================
def load_json(path: str) -> dict:
    """Cargar archivo JSON con manejo de errores"""
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except FileNotFoundError:
        raise HTTPException(404, f"Archivo no encontrado: {os.path.basename(path)}")
    except json.JSONDecodeError:
        raise HTTPException(500, f"JSON inválido: {os.path.basename(path)}")

def obtener_siguiente_consecutivo() -> str:
    """
    Genera el siguiente número de cotización con formato NASSA-YYYY-#### 
    Usa PostgreSQL para evitar duplicados en concurrencia
    """
    ano_actual = now_colombia().year
    
    if not POSTGRES_AVAILABLE:
        # Fallback a archivo JSON si PostgreSQL no disponible
        print("⚠️ PostgreSQL no disponible, usando archivo JSON como fallback")
        import fcntl
        consecutivo_file = os.path.join(CONFIG_DIR, "consecutivo.json")
        
        if not os.path.exists(consecutivo_file):
            data_inicial = {"ano_actual": ano_actual, "ultimo_consecutivo": 0}
            with open(consecutivo_file, "w", encoding="utf-8") as f:
                json.dump(data_inicial, f, ensure_ascii=False, indent=2)
        
        with open(consecutivo_file, "r+", encoding="utf-8") as f:
            fcntl.flock(f.fileno(), fcntl.LOCK_EX)
            try:
                data = json.load(f)
                if data.get("ano_actual") != ano_actual:
                    data["ano_actual"] = ano_actual
                    data["ultimo_consecutivo"] = 0
                data["ultimo_consecutivo"] += 1
                consecutivo = data["ultimo_consecutivo"]
                f.seek(0)
                json.dump(data, f, ensure_ascii=False, indent=2)
                f.truncate()
                cotizacion_id = f"NASSA-{ano_actual}-{consecutivo:04d}"
                return cotizacion_id
            finally:
                fcntl.flock(f.fileno(), fcntl.LOCK_UN)
    
    # Usar PostgreSQL (RECOMENDADO)
    try:
        with get_db_session() as session:
            # Buscar consecutivo del año actual con FOR UPDATE (lock de fila)
            consecutivo_obj = session.query(Consecutivo).filter(
                Consecutivo.ano_actual == ano_actual
            ).with_for_update().first()
            
            if consecutivo_obj is None:
                # Crear nuevo registro para el año
                print(f"📅 Creando nuevo consecutivo para el año {ano_actual}")
                consecutivo_obj = Consecutivo(
                    ano_actual=ano_actual,
                    ultimo_consecutivo=1
                )
                session.add(consecutivo_obj)
                session.commit()
                consecutivo = 1
            else:
                # Incrementar consecutivo existente
                consecutivo_obj.ultimo_consecutivo += 1
                consecutivo = consecutivo_obj.ultimo_consecutivo
                session.commit()
            
            cotizacion_id = f"NASSA-{ano_actual}-{consecutivo:04d}"
            print(f"✅ Cotización ID generado: {cotizacion_id} (consecutivo {consecutivo} del año {ano_actual})")
            return cotizacion_id
            
    except Exception as e:
        print(f"❌ Error obteniendo consecutivo de PostgreSQL: {e}")
        # Fallback a archivo JSON si falla PostgreSQL
        print("⚠️ Usando fallback a archivo JSON")
        import fcntl
        consecutivo_file = os.path.join(CONFIG_DIR, "consecutivo.json")
        
        if not os.path.exists(consecutivo_file):
            data_inicial = {"ano_actual": ano_actual, "ultimo_consecutivo": 0}
            with open(consecutivo_file, "w", encoding="utf-8") as f:
                json.dump(data_inicial, f, ensure_ascii=False, indent=2)
        
        with open(consecutivo_file, "r+", encoding="utf-8") as f:
            fcntl.flock(f.fileno(), fcntl.LOCK_EX)
            try:
                data = json.load(f)
                if data.get("ano_actual") != ano_actual:
                    data["ano_actual"] = ano_actual
                    data["ultimo_consecutivo"] = 0
                data["ultimo_consecutivo"] += 1
                consecutivo = data["ultimo_consecutivo"]
                f.seek(0)
                json.dump(data, f, ensure_ascii=False, indent=2)
                f.truncate()
                cotizacion_id = f"NASSA-{ano_actual}-{consecutivo:04d}"
                return cotizacion_id
            finally:
                fcntl.flock(f.fileno(), fcntl.LOCK_UN)

def calcular_valor_legalizacion(capacidad_instalada_w: float, parametros: dict) -> float:
    """
    Calcula el costo de legalización según tabla de rangos
    
    Args:
        capacidad_instalada_w: Capacidad en Watts
        parametros: Dict con tabla_legalizacion.rangos
    
    Returns:
        Valor de legalización (no exento de IVA)
    """
    tabla = parametros.get("tabla_legalizacion", {}).get("rangos", [])
    
    for rango in tabla:
        if rango["min"] <= capacidad_instalada_w <= rango["max"]:
            return rango["valor"]
    
    # Fallback: último rango si no encuentra
    return tabla[-1]["valor"] if tabla else 0

def obtener_equipos_defaults(equipos: dict, sistema_electrico: str = None) -> dict:
    """
    Obtener equipos marcados como default en equipos.json
    
    Args:
        equipos: Diccionario con paneles, inversores y baterías
        sistema_electrico: Tipo de sistema eléctrico (monofasico, bifasico, trifasico)
                          Si se proporciona, filtra inversores compatibles
    
    Returns:
        Dict con IDs de panel, inversor y batería por defecto
        
    Algoritmo de selección (4 niveles de prioridad):
        1. Equipo marcado como default=True y compatible con sistema_electrico
        2. Primer equipo compatible (si no hay default para ese tipo)
        3. Equipo marcado como default=True (ignorando compatibilidad)
        4. Primer equipo disponible (fallback final)
    """
    # Panel: Buscar default o primer disponible
    panel_default = next((p for p in equipos["paneles"] if p.get("default", False)), None)
    if not panel_default:
        panel_default = equipos["paneles"][0] if equipos["paneles"] else None
        if panel_default:
            print(f"⚠️ No hay panel default, usando primer disponible: {panel_default.get('id')}")
    
    # Batería: Buscar default o primera disponible
    bateria_default = next((b for b in equipos["baterias"] if b.get("default", False)), None)
    if not bateria_default:
        bateria_default = equipos["baterias"][0] if equipos["baterias"] else None
        if bateria_default:
            print(f"⚠️ No hay batería default, usando primera disponible: {bateria_default.get('id')}")
    
    # FIX #3 y #4: Selección inteligente de inversor según sistema eléctrico
    inversor_default = None
    
    if sistema_electrico:
        # NIVEL 1: Buscar inversor default=True compatible con sistema_electrico
        inversor_default = next(
            (i for i in equipos["inversores"] 
             if i.get("default", False) and i.get("tipo_sistema") == sistema_electrico),
            None
        )
        
        # NIVEL 2: Si no hay default compatible, buscar PRIMER inversor compatible
        if not inversor_default:
            print(f"⚠️ No hay inversor default para {sistema_electrico}, buscando primer compatible...")
            inversor_default = next(
                (i for i in equipos["inversores"] if i.get("tipo_sistema") == sistema_electrico),
                None
            )
            
        # NIVEL 3: Si aún no hay compatible, usar cualquier default (sin importar tipo)
        if not inversor_default:
            print(f"⚠️⚠️ No hay inversores compatibles con {sistema_electrico}, usando default general...")
            inversor_default = next(
                (i for i in equipos["inversores"] if i.get("default", False)),
                None
            )
    else:
        # Sin sistema eléctrico especificado, usar default general
        inversor_default = next((i for i in equipos["inversores"] if i.get("default", False)), None)
    
    # NIVEL 4: Fallback final - primer inversor disponible
    if not inversor_default:
        inversor_default = equipos["inversores"][0] if equipos["inversores"] else None
        if inversor_default:
            print(f"⚠️⚠️⚠️ Usando primer inversor disponible como fallback: {inversor_default.get('id')}")
    
    return {
        "panel": panel_default["id"] if panel_default else None,
        "inversor": inversor_default["id"] if inversor_default else None,
        "bateria": bateria_default["id"] if bateria_default else None
    }

# ========================================
# 📊 CARGA DE DATOS DESDE POSTGRESQL
# ========================================
def cargar_datos_desde_postgres():
    """
    Carga equipos, ciudades y parámetros desde PostgreSQL.
    Retorna estructura compatible con formato JSON legacy.
    FALLBACK: Si PostgreSQL no está disponible, carga desde archivos JSON.
    """
    # Intentar cargar desde PostgreSQL
    if POSTGRES_AVAILABLE:
        try:
            session = get_db_session()
            try:
                # Cargar paneles
                paneles_db = session.query(Panel).all()
                paneles = [{
                    "id": p.id,
                    "nombre": p.nombre,
                    "capacidad": p.capacidad,
                    "precio": p.precio,
                    "area": p.area,
                    "eficienciaPanel": p.eficienciaPanel,
                    "default": p.default
                } for p in paneles_db]
                
                # Cargar inversores
                inversores_db = session.query(Inversor).all()
                inversores = [{
                    "id": i.id,
                    "nombre": i.nombre,
                    "capacidad": i.capacidad,
                    "precio": i.precio,
                    "eficiencia": i.eficiencia,
                    "tipo": i.tipo,
                    "paneles_por_inversor": i.paneles_por_inversor,
                    "sobredimensionamiento": i.sobredimensionamiento,
                    "sistemaElectrico": i.sistemaElectrico,
                    "default": i.default
                } for i in inversores_db]
                
                # Cargar baterías
                baterias_db = session.query(Bateria).all()
                baterias = [{
                    "id": b.id,
                    "nombre": b.nombre,
                    "capacidad": b.capacidad,
                    "precio": b.precio,
                    "default": b.default
                } for b in baterias_db]
                
                # Cargar ciudades
                ciudades_db = session.query(Ciudad).all()
                ciudades = {
                    c.key: {
                        "hsp": c.hsp, 
                        "nombre": c.nombre,
                        "factorTemperatura": getattr(c, 'factorTemperatura', 0.90)  # Default 0.90 si no existe
                    } 
                    for c in ciudades_db
                }
                
                # Cargar parámetros (reconstruir dict anidado)
                parametros_db = session.query(Parametro).all()
                parametros = {}
                for p in parametros_db:
                    parametros[p.seccion] = p.data
                
                equipos = {
                    "paneles": paneles,
                    "inversores": inversores,
                    "baterias": baterias
                }
                
                print(f"✅ Datos cargados desde PostgreSQL: {len(paneles)} paneles, {len(inversores)} inversores, {len(baterias)} baterías, {len(ciudades)} ciudades")
                return equipos, ciudades, parametros
                
            finally:
                session.close()
        except Exception as e:
            print(f"⚠️ Error conectando a PostgreSQL ({type(e).__name__}: {str(e)[:100]})")
            print(f"   → Usando archivos JSON como fallback")
    else:
        print(f"⚠️ PostgreSQL no disponible (models no importados)")
        print(f"   → Usando archivos JSON como fallback")
    
    # FALLBACK: Cargar desde archivos JSON
    try:
        equipos = load_json(EQUIPOS_FILE)
        ciudades = load_json(CIUDADES_FILE)
        
        # Parametros tiene estructura diferente en JSON
        parametros_path = os.path.join(CONFIG_DIR, "parametros.json")
        if os.path.exists(parametros_path):
            parametros = load_json(parametros_path)
        else:
            # Valores por defecto si no existe el archivo
            parametros = {
                "depreciation": {
                    "enabled": True,
                    "years": 3,
                    "percentage": 0.35
                },
                "rent_deduction": {
                    "enabled": True,
                    "years": 5,
                    "percentage_base": 0.50,
                    "percentage_effective": 0.35
                }
            }
        
        print(f"✅ Datos cargados desde JSON: {len(equipos.get('paneles', []))} paneles, {len(equipos.get('inversores', []))} inversores, {len(equipos.get('baterias', []))} baterías, {len(ciudades)} ciudades")
        return equipos, ciudades, parametros
        
    except Exception as e:
        print(f"❌ Error crítico cargando datos desde JSON: {type(e).__name__}: {e}")
        raise

# ========================================
# 🧮 FUNCIÓN DE CÁLCULO
# ========================================
def calcular_cotizacion(data: dict, equipos: dict, ciudades: dict, parametros: dict = None) -> dict:
    """
    Lógica de cálculo de cotización completa con:
    - Porcentaje de ahorro de energía
    - Lógica MICRO vs STRING para inversores
    - Costo de legalización por rangos
    - Consecutivo controlado
    """
    # Cargar parámetros de configuración
    if parametros is None:
        # Fallback: cargar desde PostgreSQL si no se pasan
        _, _, parametros = cargar_datos_desde_postgres()
    
    # Extraer parámetros
    costos = parametros["costos_instalacion"]
    fiscales = parametros["parametros_fiscales"]
    proyeccion = parametros["parametros_proyeccion"]
    parametros_sistema = parametros.get("parametros_sistema", {})
    
    # PUNTO 6: Porcentaje de ahorro de energía (default 100%)
    porcentajeAhorroEnergia = float(data.get("porcentajeAhorroEnergia", parametros_sistema.get("porcentaje_ahorro_default", 100)))
    
    consumoMensual = float(data["consumoMensual"])
    # MODIFICACIÓN CRÍTICA: Consumo objetivo basado en % ahorro
    consumoObjetivo = consumoMensual * (porcentajeAhorroEnergia / 100.0)
    
    valorFactura = float(data["valorFactura"])
    valorKwh = float(data["valorKwh"])
    ciudad_key = data["ciudad"].lower().strip().replace(" ", "_")
    
    # FIX: Compatibilidad con nuevo formato de ciudades.json (objeto con .hsp)
    ciudad_data = ciudades.get(ciudad_key, ciudades.get("default", 4.5))
    hsp_value = ciudad_data.get("hsp") if isinstance(ciudad_data, dict) else ciudad_data
    hsp = float(data.get("hspCalculado") or hsp_value)
    
    # NUEVO: Factor de temperatura por ciudad
    factorTemperatura = ciudad_data.get("factorTemperatura", 0.90) if isinstance(ciudad_data, dict) else 0.90
    
    # VALIDACIÓN CRÍTICA: Si factorTemperatura está fuera del rango 0.5-1.0, posiblemente está en formato porcentual
    if factorTemperatura > 1.0:
        print(f"⚠️⚠️⚠️ ERROR DETECTADO en calcular_cotizacion(): factorTemperatura = {factorTemperatura}")
        print(f"⚠️ Valor fuera de rango [0.5-1.0]. Posiblemente está en formato porcentual.")
        print(f"⚠️ CORRECCIÓN AUTOMÁTICA: Dividiendo entre 100")
        factorTemperatura = factorTemperatura / 100
        print(f"✅ Nuevo valor: {factorTemperatura}")
    elif factorTemperatura < 0.5:
        print(f"⚠️⚠️⚠️ ERROR DETECTADO en calcular_cotizacion(): factorTemperatura = {factorTemperatura}")
        print(f"⚠️ Valor demasiado bajo [<0.5]. Usando default 0.90")
        factorTemperatura = 0.90

    panel = next((x for x in equipos["paneles"] if x["id"] == data["panel"]), None)
    inversor = next((x for x in equipos["inversores"] if x["id"] == data["inversor"]), None)
    bateria = next((x for x in equipos["baterias"] if x["id"] == data.get("bateria")), None) if data.get("bateria") else None
    if not panel or not inversor:
        raise ValueError("Panel o inversor no encontrado")
    
    factorAreaEfectiva = parametros_sistema.get("factor_area_efectiva", 1.2)
    areaPanel = panel.get("area", 2.0)
    eficiencia_panel = panel.get("eficienciaPanel", parametros_sistema.get("eficiencia_panel_default", 1.0))
    eficiencia_inversor = inversor.get("eficiencia", parametros_sistema.get("eficiencia_inversor_default", 1.0))
    
    # CÁLCULO INICIAL de paneles basado en consumo objetivo
    # Fórmula: energiaPanelDia = (capacidad * eficiencia_panel * hsp * factorTemperatura) / 1000
    consumoDiario = consumoObjetivo / 30
    energiaPanelDia = (panel["capacidad"] * eficiencia_panel * hsp * factorTemperatura) / 1000
    numeroPaneles_inicial = int(ceil((consumoDiario * 1.2) / energiaPanelDia))
    
    # DEBUGGING: Logging detallado del cálculo
    print(f"\n{'='*80}")
    print(f"🔍 CÁLCULO DE PANELES - DEBUGGING")
    print(f"{'='*80}")
    print(f"📊 DATOS DE ENTRADA:")
    print(f"   consumoMensual: {consumoMensual} kWh")
    print(f"   consumoObjetivo: {consumoObjetivo} kWh (porcentaje: {porcentajeAhorroEnergia}%)")
    print(f"   consumoDiario: {consumoDiario:.2f} kWh/día")
    print(f"\n🔧 EQUIPOS SELECCIONADOS:")
    print(f"   Panel: {panel['id']} - {panel['nombre']}")
    print(f"   Capacidad panel: {panel['capacidad']} W")
    print(f"   Eficiencia panel: {eficiencia_panel}")
    print(f"   Inversor: {inversor['id']} - {inversor['nombre']}")
    print(f"   Eficiencia inversor: {eficiencia_inversor}")
    print(f"\n🌡️ PARÁMETROS CLIMÁTICOS:")
    print(f"   Ciudad: {data.get('ciudad')}")
    print(f"   HSP: {hsp}")
    print(f"   Factor Temperatura: {factorTemperatura} {'⚠️ ERROR: Debe estar entre 0.5 y 1.0' if factorTemperatura > 1.0 else '✅'}")
    print(f"\n⚡ CÁLCULO DE ENERGÍA POR PANEL:")
    print(f"   Fórmula: (capacidad * efic_panel * HSP * factorTemp) / 1000")
    print(f"   energiaPanelDia = ({panel['capacidad']} * {eficiencia_panel} * {hsp} * {factorTemperatura}) / 1000")
    print(f"   energiaPanelDia = {energiaPanelDia:.4f} kWh/día/panel")
    print(f"\n🔢 NÚMERO DE PANELES:")
    print(f"   Fórmula: ceil((consumoDiario * 1.2) / energiaPanelDia)")
    print(f"   numeroPaneles = ceil(({consumoDiario:.2f} * 1.2) / {energiaPanelDia:.4f})")
    print(f"   numeroPaneles_inicial = {numeroPaneles_inicial}")
    print(f"{'='*80}\n")
    
    # PUNTO 5: Lógica de inversores MICRO vs STRING
    tipo_inversor = inversor.get("tipo", "STRING")
    
    if tipo_inversor == "MICRO":
        # MICRO: Basado en paneles por inversor
        paneles_por_inversor = inversor.get("paneles_por_inversor", 4)
        numeroInversores_raw = numeroPaneles_inicial / paneles_por_inversor
        decimal = numeroInversores_raw - int(numeroInversores_raw)
        
        if decimal < 0.49:
            numeroInversores = int(numeroInversores_raw)
        else:
            numeroInversores = int(numeroInversores_raw) + 1
        
        # Recalcular paneles si se redondeó hacia abajo
        numeroPaneles = numeroInversores * paneles_por_inversor
        
    else:  # STRING
        # STRING: Basado en capacidad con sobredimensionamiento
        sobredimensionamiento = inversor.get("sobredimensionamiento", 0.40)
        capacidad_inversor_w = inversor["capacidad"]
        capacidad_efectiva_w = capacidad_inversor_w * (1 + sobredimensionamiento)
        
        # Calcular capacidad instalada inicial
        capacidadInstalada_inicial = (numeroPaneles_inicial * panel["capacidad"]) / 1000  # kW
        
        # Número de inversores necesarios
        numeroInversores = int(ceil((capacidadInstalada_inicial * 1000) / capacidad_efectiva_w))
        
        # Capacidad máxima que pueden manejar los inversores
        capacidad_maxima_sistema = (numeroInversores * capacidad_efectiva_w) / 1000  # kW
        
        # Si la capacidad inicial excede la máxima, ajustar paneles
        if capacidadInstalada_inicial > capacidad_maxima_sistema:
            numeroPaneles = int((capacidad_maxima_sistema * 1000) / panel["capacidad"])
        else:
            numeroPaneles = numeroPaneles_inicial
    
    # RECALCULAR todo con paneles y/o inversores ajustados
    capacidadInstalada = (numeroPaneles * panel["capacidad"]) / 1000  # kW
    generacionMensual = numeroPaneles * energiaPanelDia * 30 * eficiencia_inversor
    generacionAnual = generacionMensual * 12
    areaRequerida = round(numeroPaneles * areaPanel * factorAreaEfectiva, 2)
    
    # DEBUGGING: Logging del resultado final
    print(f"\n{'='*80}")
    print(f"📈 RESULTADO FINAL DEL CÁLCULO")
    print(f"{'='*80}")
    print(f"   Paneles finales: {numeroPaneles}")
    print(f"   Inversores: {numeroInversores}")
    print(f"   Capacidad instalada: {capacidadInstalada:.2f} kW")
    print(f"   Generación MENSUAL: {generacionMensual:.2f} kWh {'⚠️ POSIBLE ERROR' if generacionMensual > consumoMensual * 10 else '✅'}")
    print(f"   Generación ANUAL: {generacionAnual:.2f} kWh")
    print(f"   Área requerida: {areaRequerida} m²")
    print(f"   Ratio generación/consumo: {(generacionMensual/consumoMensual*100):.1f}%")
    print(f"{'='*80}\n")

    # Costos básicos desde parámetros configurables
    soporteria = costos["soporteria_por_panel"]
    instalacion = costos["instalacion_por_panel"]
    materialesAdicionales = costos["materiales_por_panel"]
    mantenimientoAnual = costos["mantenimiento_anual_por_kw"]

    costoPaneles = numeroPaneles * panel.get("precio", 0)
    costoInversores = numeroInversores * inversor.get("precio", 0)
    costoBaterias = bateria.get("precio", 0) if bateria else 0
    costoSoporteria = numeroPaneles * soporteria
    costoInstalacion = numeroPaneles * instalacion
    costoMateriales = numeroPaneles * materialesAdicionales

    # PUNTO 1: Valor de legalización (NO exento de IVA)
    capacidadInstalada_w = capacidadInstalada * 1000
    valorLegalizacion = calcular_valor_legalizacion(capacidadInstalada_w, parametros)
    incluir_legalizacion = data.get("legalizacion", "NO") == "SI"
    costoLegalizacion = valorLegalizacion if incluir_legalizacion else 0
    ivaLegalizacion = costoLegalizacion * fiscales["iva_porcentaje"]

    # Cálculo de costos totales
    subtotalAntesIVA = costoPaneles + costoInversores + costoBaterias + costoSoporteria + costoInstalacion + costoMateriales
    subtotalConLegalizacion = subtotalAntesIVA + costoLegalizacion
    
    # IVA: Equipos + Legalización
    ivaEquipos = (costoBaterias + costoSoporteria + costoInstalacion + costoMateriales) * fiscales["iva_porcentaje"]
    ivaTotal = ivaEquipos + ivaLegalizacion
    
    valorTotalSistema = subtotalConLegalizacion + ivaTotal

    porcentajeProduccionMensual = ((generacionMensual / consumoMensual) * 100) if consumoMensual > 0 else 0
    ahorroMensualEnergia = generacionMensual * valorKwh
    ahorroAnualEnergia = ahorroMensualEnergia * 12

    deduccionRentaBase = subtotalAntesIVA * fiscales["deduccion_renta_base_porcentaje"]
    deduccionRentaEfectiva = deduccionRentaBase * fiscales["impuesto_renta_porcentaje"]
    ahorroAnualDeduccion = (deduccionRentaEfectiva / fiscales["anos_deduccion"])
    depreciacionAnual = subtotalAntesIVA / fiscales["anos_depreciacion"]
    ahorroAnualDepreciacion = depreciacionAnual * fiscales["impuesto_renta_porcentaje"]

    tabla = []
    ahorroAcum = 0
    payback = 0
    alcanzado = False
    costoMantBase = capacidadInstalada * mantenimientoAnual
    acumxgen = acumxdeduc = acumxdepre = 0

    # Usar parámetros configurables para la proyección
    anos_proyeccion = proyeccion["anos_proyeccion"]
    degradacion_anual = proyeccion["degradacion_anual_panel"]
    factor_primer_ano = proyeccion["factor_primer_ano"]
    incremento_anual = proyeccion["incremento_anual_kwh"]
    anos_depreciacion = fiscales["anos_depreciacion"]
    anos_deduccion = fiscales["anos_deduccion"]

    for year in range(1, anos_proyeccion + 1):
        valorKwhAño = valorKwh * ((1 + incremento_anual) ** (year - 1))
        degradacion = ((1 - degradacion_anual) ** (year - 1))
        factorInicio = factor_primer_ano if year == 1 else 1
        prodAnual = generacionAnual * degradacion * factorInicio
        ahorroGeneracion = min(prodAnual * valorKwhAño, valorFactura * 12 * ((1 + incremento_anual) ** (year - 1)))
        ahorroDep = ahorroAnualDepreciacion if year <= anos_depreciacion else 0
        ahorroDed = ahorroAnualDeduccion if year <= anos_deduccion else 0
        costoMant = costoMantBase * ((1 + incremento_anual) ** (year - 1))
        ahorroTotal = ahorroGeneracion + ahorroDep + ahorroDed - costoMant
        ahorroAcum += ahorroTotal
        roi = ((ahorroAcum - valorTotalSistema) / valorTotalSistema) * 100 if valorTotalSistema else 0

        if not alcanzado and ahorroAcum >= valorTotalSistema:
            alcanzado = True
            payback = year

        if not alcanzado or year <= payback:
            acumxgen += ahorroGeneracion
            acumxdeduc += ahorroDed
            acumxdepre += ahorroDep

        tabla.append({
            "año": year,
            "valorKwhAño": round(valorKwhAño),
            "produccionAnual": round(prodAnual),
            "ahorroGeneracion": round(ahorroGeneracion),
            "ahorroDep": round(ahorroDep),
            "ahorroDed": round(ahorroDed),
            "costoMant": round(costoMant),
            "ahorroTotalAño": round(ahorroTotal),
            "ahorroAcumulado": round(ahorroAcum),
            "roi": round(roi, 2)
        })

    TotalAcum = acumxgen + acumxdeduc + acumxdepre
    TotalAcum = acumxgen + acumxdeduc + acumxdepre
    tiempoRetorno = payback or (valorTotalSistema / (ahorroAnualEnergia + ahorroAnualDeduccion + ahorroAnualDepreciacion + 1e-9))

    # PUNTO 3: Generar consecutivo controlado
    cotizacion_id = obtener_siguiente_consecutivo()

    return {
        "fecha": now_colombia().isoformat(),
        "cotizacionId": cotizacion_id,
        "panel": {"id": panel["id"], "nombre": panel["nombre"], "capacidad": panel["capacidad"], "area": areaPanel},
        "inversor": {
            "id": inversor["id"], 
            "nombre": inversor["nombre"], 
            "capacidad": inversor["capacidad"],
            "tipo": tipo_inversor,
            "paneles_por_inversor": inversor.get("paneles_por_inversor") if tipo_inversor == "MICRO" else None,
            "sobredimensionamiento": inversor.get("sobredimensionamiento") if tipo_inversor == "STRING" else None
        },
        "bateria": {"id": bateria["id"], "nombre": bateria["nombre"]} if bateria else None,
        "numeroPaneles": numeroPaneles,
        "numeroInversores": numeroInversores,
        "capacidadInstalada": round(capacidadInstalada, 2),
        "areaRequerida": areaRequerida,
        "areaDisponibleCliente": float(data.get("areaDisponible", 0)),
        "generacionMensual": round(generacionMensual),
        "generacionAnual": round(generacionAnual),
        "porcentajeAhorroEnergia": porcentajeAhorroEnergia,
        "consumoObjetivo": round(consumoObjetivo),
        # PUNTO 2: Desglose detallado de costos para preview
        "desgloseCostos": {
            "costoPaneles": round(costoPaneles),
            "costoInversores": round(costoInversores),
            "costoBaterias": round(costoBaterias),
            "costoSoporteria": round(costoSoporteria),
            "costoInstalacion": round(costoInstalacion),
            "costoMateriales": round(costoMateriales),
            "costoLegalizacion": round(costoLegalizacion),
            "subtotalAntesIVA": round(subtotalAntesIVA),
            "ivaEquipos": round(ivaEquipos),
            "ivaLegalizacion": round(ivaLegalizacion),
            "ivaTotal": round(ivaTotal)
        },
        "valorTotalSistema": round(valorTotalSistema),
        "ahorroMensualEnergia": round(ahorroMensualEnergia),
        "ahorroAnualEnergia": round(ahorroAnualEnergia),
        "tiempoRetorno": round(tiempoRetorno, 1),
        "tablaAhorros": tabla,
        "porcentajeProduccionMensual": round(porcentajeProduccionMensual, 1),
        "subtotalAntesIVA": round(subtotalAntesIVA),
        "deduccionRentaBase": round(deduccionRentaBase),
        "deduccionRentaEfectiva": round(deduccionRentaEfectiva),
        "depreciacionAnual": round(depreciacionAnual),
        "ahorroAnualDepreciacion": round(ahorroAnualDepreciacion),
        "ahorroTotalDepreciacion": round(ahorroAnualDepreciacion * 3),
        "ahorroTotalDeduccion": round(deduccionRentaEfectiva),
        "acumxgen": round(acumxgen),
        "acumxdeduc": round(acumxdeduc),
        "acumxdepre": round(acumxdepre),
        "TotalAcum": round(TotalAcum)
    }

def calcular_segunda_opcion(data: dict, equipos: dict, ciudades: dict, areaDisponible: float, cotizacion_id_base: str, parametros: dict = None) -> dict:
    """
    Calcula cotización ajustada al área disponible del cliente.
    Reduce número de paneles para que quepan en el espacio real.
    Usa la misma lógica MICRO/STRING y legalización que calcular_cotizacion.
    
    Args:
        cotizacion_id_base: ID base (ej: "NASSA-2025-0001") para mantener consistencia
    """
    # Cargar parámetros si no se pasan
    if parametros is None:
        _, _, parametros = cargar_datos_desde_postgres()
    
    costos = parametros["costos_instalacion"]
    fiscales = parametros["parametros_fiscales"]
    proyeccion = parametros["parametros_proyeccion"]
    parametros_sistema = parametros.get("parametros_sistema", {})
    
    # Porcentaje de ahorro
    porcentajeAhorroEnergia = float(data.get("porcentajeAhorroEnergia", parametros_sistema.get("porcentaje_ahorro_default", 100)))
    
    consumoMensual = float(data["consumoMensual"])
    consumoObjetivo = consumoMensual * (porcentajeAhorroEnergia / 100.0)
    valorFactura = float(data["valorFactura"])
    valorKwh = float(data["valorKwh"])
    ciudad_key = data["ciudad"].lower().strip().replace(" ", "_")
    
    ciudad_data = ciudades.get(ciudad_key, ciudades.get("default", 4.5))
    hsp_value = ciudad_data.get("hsp") if isinstance(ciudad_data, dict) else ciudad_data
    hsp = float(data.get("hspCalculado") or hsp_value)
    
    # NUEVO: Factor de temperatura por ciudad
    factorTemperatura = ciudad_data.get("factorTemperatura", 0.90) if isinstance(ciudad_data, dict) else 0.90
    
    # VALIDACIÓN CRÍTICA: Si factorTemperatura está fuera del rango 0.5-1.0, posiblemente está en formato porcentual
    if factorTemperatura > 1.0:
        print(f"⚠️⚠️⚠️ ERROR DETECTADO en calcular_segunda_opcion(): factorTemperatura = {factorTemperatura}")
        print(f"⚠️ Valor fuera de rango [0.5-1.0]. Posiblemente está en formato porcentual.")
        print(f"⚠️ CORRECCIÓN AUTOMÁTICA: Dividiendo entre 100")
        factorTemperatura = factorTemperatura / 100
        print(f"✅ Nuevo valor: {factorTemperatura}")
    elif factorTemperatura < 0.5:
        print(f"⚠️⚠️⚠️ ERROR DETECTADO en calcular_segunda_opcion(): factorTemperatura = {factorTemperatura}")
        print(f"⚠️ Valor demasiado bajo [<0.5]. Usando default 0.90")
        factorTemperatura = 0.90

    panel = next((x for x in equipos["paneles"] if x["id"] == data["panel"]), None)
    inversor = next((x for x in equipos["inversores"] if x["id"] == data["inversor"]), None)
    bateria = next((x for x in equipos["baterias"] if x["id"] == data.get("bateria")), None) if data.get("bateria") else None
    
    factorAreaEfectiva = parametros_sistema.get("factor_area_efectiva", 1.2)
    areaPanel = panel.get("area", 2.0)
    
    # CALCULAR NÚMERO MÁXIMO DE PANELES QUE CABEN EN EL ÁREA
    numeroPaneles_max = max(1, int(areaDisponible / (areaPanel * factorAreaEfectiva)))
    
    # Aplicar lógica MICRO/STRING con restricción de área
    eficiencia_panel = panel.get("eficienciaPanel", parametros_sistema.get("eficiencia_panel_default", 1.0))
    eficiencia_inversor = inversor.get("eficiencia", parametros_sistema.get("eficiencia_inversor_default", 1.0))
    energiaPanelDia = (panel["capacidad"] * eficiencia_panel * hsp * factorTemperatura) / 1000
    
    tipo_inversor = inversor.get("tipo", "STRING")
    
    if tipo_inversor == "MICRO":
        paneles_por_inversor = inversor.get("paneles_por_inversor", 4)
        # Ajustar a múltiplo de paneles_por_inversor que quepa en área
        numeroInversores = numeroPaneles_max // paneles_por_inversor
        if numeroInversores == 0:
            numeroInversores = 1
        numeroPaneles = numeroInversores * paneles_por_inversor
        # Si excede área, reducir un inversor
        if numeroPaneles > numeroPaneles_max:
            numeroInversores = max(1, numeroInversores - 1)
            numeroPaneles = numeroInversores * paneles_por_inversor
    else:  # STRING
        # Usar todos los paneles que quepan
        numeroPaneles = numeroPaneles_max
        sobredimensionamiento = inversor.get("sobredimensionamiento", 0.40)
        capacidad_inversor_w = inversor["capacidad"]
        capacidad_efectiva_w = capacidad_inversor_w * (1 + sobredimensionamiento)
        capacidadInstalada_temp = (numeroPaneles * panel["capacidad"]) / 1000
        numeroInversores = int(ceil((capacidadInstalada_temp * 1000) / capacidad_efectiva_w))
    
    # RECALCULAR con paneles/inversores ajustados
    capacidadInstalada = (numeroPaneles * panel["capacidad"]) / 1000
    generacionMensual = numeroPaneles * energiaPanelDia * 30 * eficiencia_inversor
    generacionAnual = generacionMensual * 12
    areaRequerida = round(numeroPaneles * areaPanel * factorAreaEfectiva, 2)
    
    # Costos con legalización
    soporteria = costos["soporteria_por_panel"]
    instalacion = costos["instalacion_por_panel"]
    materialesAdicionales = costos["materiales_por_panel"]
    mantenimientoAnual = costos["mantenimiento_anual_por_kw"]
    
    costoPaneles = numeroPaneles * panel.get("precio", 0)
    costoInversores = numeroInversores * inversor.get("precio", 0)
    costoBaterias = bateria.get("precio", 0) if bateria else 0
    costoSoporteria = numeroPaneles * soporteria
    costoInstalacion = numeroPaneles * instalacion
    costoMateriales = numeroPaneles * materialesAdicionales
    
    # Legalización (si aplica)
    capacidadInstalada_w = capacidadInstalada * 1000
    valorLegalizacion = calcular_valor_legalizacion(capacidadInstalada_w, parametros)
    incluir_legalizacion = data.get("legalizacion", "NO") == "SI"
    costoLegalizacion = valorLegalizacion if incluir_legalizacion else 0
    ivaLegalizacion = costoLegalizacion * fiscales["iva_porcentaje"]
    
    subtotalAntesIVA = costoPaneles + costoInversores + costoBaterias + costoSoporteria + costoInstalacion + costoMateriales
    subtotalConLegalizacion = subtotalAntesIVA + costoLegalizacion
    ivaEquipos = (costoBaterias + costoSoporteria + costoInstalacion + costoMateriales) * fiscales["iva_porcentaje"]
    ivaTotal = ivaEquipos + ivaLegalizacion
    valorTotalSistema = subtotalConLegalizacion + ivaTotal
    
    porcentajeProduccionMensual = ((generacionMensual / consumoMensual) * 100) if consumoMensual > 0 else 0
    ahorroMensualEnergia = generacionMensual * valorKwh
    ahorroAnualEnergia = ahorroMensualEnergia * 12
    
    # Fiscales
    deduccionRentaBase = subtotalAntesIVA * fiscales["deduccion_renta_base_porcentaje"]
    deduccionRentaEfectiva = deduccionRentaBase * fiscales["impuesto_renta_porcentaje"]
    ahorroAnualDeduccion = (deduccionRentaEfectiva / fiscales["anos_deduccion"])
    depreciacionAnual = subtotalAntesIVA / fiscales["anos_depreciacion"]
    ahorroAnualDepreciacion = depreciacionAnual * fiscales["impuesto_renta_porcentaje"]
    
    # Proyección 30 años - usar nombres correctos
    anos_proyeccion = proyeccion["anos_proyeccion"]
    degradacion_anual = proyeccion["degradacion_anual_panel"]
    factor_primer_ano = proyeccion["factor_primer_ano"]
    incremento_anual = proyeccion["incremento_anual_kwh"]
    anos_depreciacion = fiscales["anos_depreciacion"]
    anos_deduccion = fiscales["anos_deduccion"]
    
    tabla = []
    acumxgen = acumxdeduc = acumxdepre = 0
    ahorroAcum = 0
    payback = None
    costoMantBase = capacidadInstalada * mantenimientoAnual
    alcanzado = False
    
    for year in range(1, anos_proyeccion + 1):
        valorKwhAño = valorKwh * ((1 + incremento_anual) ** (year - 1))
        degradacion = ((1 - degradacion_anual) ** (year - 1))
        factorInicio = factor_primer_ano if year == 1 else 1
        prodAnual = generacionAnual * degradacion * factorInicio
        ahorroGeneracion = min(prodAnual * valorKwhAño, valorFactura * 12 * ((1 + incremento_anual) ** (year - 1)))
        
        ahorroDep = ahorroAnualDepreciacion if year <= anos_depreciacion else 0
        ahorroDed = ahorroAnualDeduccion if year <= anos_deduccion else 0
        costoMant = costoMantBase * ((1 + incremento_anual) ** (year - 1))
        
        ahorroTotal = ahorroGeneracion + ahorroDep + ahorroDed - costoMant
        ahorroAcum += ahorroTotal
        roi = ((ahorroAcum - valorTotalSistema) / valorTotalSistema) * 100 if valorTotalSistema else 0
        
        if not alcanzado and ahorroAcum >= valorTotalSistema:
            alcanzado = True
            payback = year
        
        if not alcanzado or year <= (payback or anos_proyeccion):
            acumxgen += ahorroGeneracion
            acumxdeduc += ahorroDed
            acumxdepre += ahorroDep
        
        tabla.append({
            "año": year,
            "valorKwhAño": round(valorKwhAño),
            "produccionAnual": round(prodAnual),
            "ahorroGeneracion": round(ahorroGeneracion),
            "ahorroDep": round(ahorroDep),
            "ahorroDed": round(ahorroDed),
            "costoMant": round(costoMant),
            "ahorroTotalAño": round(ahorroTotal),
            "ahorroAcumulado": round(ahorroAcum),
            "roi": round(roi, 2)
        })
    
    TotalAcum = acumxgen + acumxdeduc + acumxdepre
    tiempoRetorno = payback or (valorTotalSistema / (ahorroAnualEnergia + ahorroAnualDeduccion + ahorroAnualDepreciacion + 1e-9))
    
    return {
        "fecha": now_colombia().isoformat(),
        "cotizacionId": cotizacion_id_base + "-OP2",  # Usar el mismo ID base con sufijo
        "panel": {"id": panel["id"], "nombre": panel["nombre"], "capacidad": panel["capacidad"], "area": areaPanel},
        "inversor": {
            "id": inversor["id"], 
            "nombre": inversor["nombre"], 
            "capacidad": inversor["capacidad"],
            "tipo": tipo_inversor,
            "paneles_por_inversor": inversor.get("paneles_por_inversor") if tipo_inversor == "MICRO" else None,
            "sobredimensionamiento": inversor.get("sobredimensionamiento") if tipo_inversor == "STRING" else None
        },
        "bateria": {"id": bateria["id"], "nombre": bateria["nombre"]} if bateria else None,
        "numeroPaneles": numeroPaneles,
        "numeroInversores": numeroInversores,
        "capacidadInstalada": round(capacidadInstalada, 2),
        "areaRequerida": areaRequerida,
        "areaDisponibleCliente": areaDisponible,
        "generacionMensual": round(generacionMensual),
        "generacionAnual": round(generacionAnual),
        "porcentajeAhorroEnergia": porcentajeAhorroEnergia,
        "consumoObjetivo": round(consumoObjetivo),
        "desgloseCostos": {
            "costoPaneles": round(costoPaneles),
            "costoInversores": round(costoInversores),
            "costoBaterias": round(costoBaterias),
            "costoSoporteria": round(costoSoporteria),
            "costoInstalacion": round(costoInstalacion),
            "costoMateriales": round(costoMateriales),
            "costoLegalizacion": round(costoLegalizacion),
            "subtotalAntesIVA": round(subtotalAntesIVA),
            "subtotalConLegalizacion": round(subtotalConLegalizacion),
            "ivaEquipos": round(ivaEquipos),
            "ivaLegalizacion": round(ivaLegalizacion),
            "ivaTotal": round(ivaTotal)
        },
        "valorTotalSistema": round(valorTotalSistema),
        "ahorroMensualEnergia": round(ahorroMensualEnergia),
        "ahorroAnualEnergia": round(ahorroAnualEnergia),
        "tiempoRetorno": round(tiempoRetorno, 1),
        "tablaAhorros": tabla,
        "porcentajeProduccionMensual": round(porcentajeProduccionMensual, 1),
        "subtotalAntesIVA": round(subtotalAntesIVA),
        "deduccionRentaBase": round(deduccionRentaBase),
        "deduccionRentaEfectiva": round(deduccionRentaEfectiva),
        "depreciacionAnual": round(depreciacionAnual),
        "ahorroAnualDepreciacion": round(ahorroAnualDepreciacion),
        "ahorroTotalDepreciacion": round(ahorroAnualDepreciacion * 3),
        "ahorroTotalDeduccion": round(deduccionRentaEfectiva),
        "acumxgen": round(acumxgen),
        "acumxdeduc": round(acumxdeduc),
        "acumxdepre": round(acumxdepre),
        "TotalAcum": round(TotalAcum)
    }

# ========================================
# 📄 PROCESAMIENTO DE TEMPLATE PPTX
# ========================================
def _normalize(t: str) -> str:
    """Normaliza texto para comparación de encabezados"""
    if t is None:
        return ""
    t = unicodedata.normalize("NFKD", t).encode("ascii", "ignore").decode("ascii")
    t = t.strip().lower().replace("_", "").replace("-", "").replace(" ", "")
    return t

def _find_table_by_name_or_headers(prs, preferred_name: str = "TABLA_AHORROS"):
    """Busca tabla por nombre o por encabezados"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "name", "") == preferred_name and shape.has_table:
                return shape.table
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                if table.rows and table.columns:
                    headers = [_normalize(cell.text) for cell in table.rows[0].cells]
                    if any("AÑO" in h or "año" in h for h in headers) and any("Ahorro" in h for h in headers):
                        return table
    return None

def _map_columns_by_header(table):
    """Mapea índices de columnas según encabezados"""
    header_map = {}
    headers = [cell.text.strip() for cell in table.rows[0].cells]  # Sin normalizar para preservar mayúsculas
    
    for idx, h in enumerate(headers):
        h_lower = h.lower().replace(" ", "").replace("_", "").replace("-", "")
        
        if "año" in h_lower:
            header_map["año"] = idx
        elif "valorkwh" in h_lower:
            header_map["valorKwh"] = idx
        elif "producción" in h_lower or "produccion" in h_lower:
            header_map["produccionAnual"] = idx
        elif "generacion" in h_lower or "generación" in h_lower:
            header_map["generacion"] = idx
        elif "depreciacion" in h_lower or "depreciación" in h_lower:
            header_map["depreciacion"] = idx
        elif "deduccion" in h_lower or "deducción" in h_lower:
            header_map["deduccion"] = idx
        elif "costomtto" in h_lower or "costomant" in h_lower:
            header_map["costoMtto"] = idx
        elif "totahorro" in h_lower or "ahorrototal" in h_lower:
            header_map["ahorroAño"] = idx
        elif "acumulado" in h_lower:
            header_map["acumulado"] = idx
        elif h.upper() == "ROI":
            header_map["roi"] = idx
    return header_map

def fill_ahorros_table_in_ppt(prs, tabla_ahorros: list, max_years: int = None):
    """Llena la tabla de ahorros en el PPTX con formato y alineación
    
    Args:
        prs: Presentación de PowerPoint
        tabla_ahorros: Lista con datos de ahorros por año
        max_years: Número máximo de años a llenar (None = detectar automáticamente desde tabla)
    """
    table = _find_table_by_name_or_headers(prs, "TABLA_AHORROS")
    if table is None:
        print("⚠️ No se encontró TABLA_AHORROS en el Template - continuando sin llenar tabla")
        return

    col_map = _map_columns_by_header(table)
    print(f"📊 Columnas detectadas en tabla: {list(col_map.keys())}")
    
    required = {"año", "roi"}
    if not required.issubset(col_map.keys()):
        print(f"⚠️ Encabezados insuficientes. Encontrados: {list(col_map.keys())} - continuando sin llenar tabla")
        return

    # Detectar automáticamente el número de filas disponibles (excluyendo el encabezado)
    filas_disponibles = len(table.rows) - 1  # -1 porque la primera fila es el encabezado
    print(f"📋 Filas disponibles en tabla: {filas_disponibles}")
    
    # Determinar cuántos años llenar: el mínimo entre filas disponibles, datos disponibles y max_years (si se especifica)
    if max_years is None:
        n = min(filas_disponibles, len(tabla_ahorros))
    else:
        n = min(max_years, filas_disponibles, len(tabla_ahorros))
    
    print(f"✍️  Llenando {n} años en la tabla")
    
    for i in range(n):
        row_idx = i + 1
        if row_idx >= len(table.rows):
            break
        
        año_data = tabla_ahorros[i]
        año = año_data.get("año", i + 1)
        
        # Función auxiliar para formatear celdas con formato ARIAL 9pt consistente
        def set_cell_value(col_key, value_text):
            if col_key in col_map:
                cell = table.cell(row_idx, col_map[col_key])
                
                # Aplicar formato Arial 9pt de forma explícita
                if cell.text_frame and cell.text_frame.paragraphs:
                    for paragraph in cell.text_frame.paragraphs:
                        # Si hay runs, usar el primer run
                        if paragraph.runs:
                            # Limpiar texto de todos los runs
                            for run in paragraph.runs:
                                run.text = ""
                            # Poner el nuevo valor en el primer run y aplicar formato
                            first_run = paragraph.runs[0]
                            first_run.text = value_text
                            first_run.font.name = "Arial"
                            first_run.font.size = Pt(7)
                        else:
                            # Si no hay runs, crear uno con formato Arial 7pt
                            run = paragraph.add_run()
                            run.text = value_text
                            run.font.name = "Arial"
                            run.font.size = Pt(7)
                        
                        # Cambiar alineación para columnas numéricas
                        if col_key in ["valorKwh", "produccionAnual", "generacion", "depreciacion", 
                                     "deduccion", "costoMtto", "ahorroAño", "acumulado", "roi"]:
                            paragraph.alignment = PP_ALIGN.RIGHT
                        else:
                            paragraph.alignment = PP_ALIGN.CENTER
        
        # Llenar columnas con valores formateados
        set_cell_value("año", str(año))
        set_cell_value("valorKwh", f"${año_data.get('valorKwhAño', 0):,.0f}")
        set_cell_value("produccionAnual", f"{año_data.get('produccionAnual', 0):,.0f}")  # CORREGIDO: era 'produccionAnual'
        set_cell_value("generacion", f"${año_data.get('ahorroGeneracion', 0):,.0f}")
        set_cell_value("depreciacion", f"${año_data.get('ahorroDep', 0):,.0f}")
        set_cell_value("deduccion", f"${año_data.get('ahorroDed', 0):,.0f}")
        set_cell_value("costoMtto", f"${año_data.get('costoMant', 0):,.0f}")
        set_cell_value("ahorroAño", f"${año_data.get('ahorroTotalAño', 0):,.0f}")
        set_cell_value("acumulado", f"${año_data.get('ahorroAcumulado', 0):,.0f}")
        set_cell_value("roi", f"{año_data.get('roi', 0):.2f}%")

    # Limpiar filas vacías
    for j in range(n + 1, len(table.rows)):
        for c in range(len(table.columns)):
            table.cell(j, c).text = ""

def build_placeholders(req: dict, resultado: dict, opcion: str = "") -> dict:
    """
    Construye diccionario de placeholders - TODOS máximo 8 letras
    opcion: "" para una sola opción, "OPCIÓN 1" o "OPCIÓN 2 - Ajustada a área disponible"
    """
    print(f"   🏗️  Construyendo placeholders para: {opcion if opcion else '(única)'}")
    print(f"      N_PANEL: {resultado['numeroPaneles']}, CAP_KW: {resultado['capacidadInstalada']} kW")
    print(f"      AREA_REQ: {resultado['areaRequerida']} m², INVER: ${resultado['valorTotalSistema']:,.0f}")
    
    # Manejar baterías: si no hay batería, dejar campos en blanco
    num_baterias = "1" if resultado.get("bateria") else " "
    bateria_modelo = resultado["bateria"]["nombre"] if resultado.get("bateria") else " "
    
    # Mapeo de nombres de sistemas a formato correcto
    tipo_sistema_map = {
        "ongrid": "ON GRID",
        "offgrid": "OFF GRID",
        "hibrido_incluido": "HIBRIDO CON BATERIA",
        "hibrido_opcional": "HIBRIDO OPCIONAL BATERIA"
    }
    tipo_sistema_texto = tipo_sistema_map.get(req["tipoSistemaFV"], req["tipoSistemaFV"].upper())
    
    return {
        # Información general (8 letras máx)
        "{{COT_ID}}": resultado["cotizacionId"],
        "{{FECHA}}": resultado["fecha"][:10],
        "{{OPCION}}": opcion,
        
        # Datos del cliente (8 letras máx)
        "{{NOMBRE}}": req["nombre"],
        "{{EMAIL}}": req["email"],
        "{{TELEFONO}}": req["telefono"],
        "{{CIUDAD}}": req["ciudad"],
        "{{DIRECC}}": req["direccion"],
        "{{NIC}}": req.get("nic", "N/A"),
        
        # Consumo energético (8 letras máx)
        "{{CONSUMO}}": f"{req['consumoMensual']:.0f} kWh",
        "{{FACTURA}}": f"${req['valorFactura']:,.0f}",
        "{{VAL_KWH}}": f"${req['valorKwh']:,.0f}",
        
        # Características del inmueble (8 letras máx)
        "{{VIVIENDA}}": req["tipoVivienda"],
        "{{SIS_ELEC}}": req["sistemaElectrico"],
        "{{TIPO_FV}}": tipo_sistema_texto,
        
        # Equipamiento (8 letras máx)
        "{{N_PANEL}}": str(resultado["numeroPaneles"]),
        "{{M_PANEL}}": resultado["panel"]["nombre"],
        "{{N_INVER}}": str(resultado["numeroInversores"]),
        "{{M_INVER}}": resultado["inversor"]["nombre"],
        "{{N_BATER}}": num_baterias,
        "{{M_BATER}}": bateria_modelo,
        
        # Especificaciones técnicas (8 letras máx)
        "{{CAP_KW}}": f"{resultado['capacidadInstalada']} kW",
        "{{GEN_MES}}": f"{resultado['generacionMensual']} kWh",
        "{{AREA_REQ}}": f"{resultado['areaRequerida']} m²",
        
        # Análisis financiero (8 letras máx)
        "{{INVER}}": f"${resultado['valorTotalSistema']:,.0f}",
        "{{SUBTOT}}": f"${resultado['subtotalAntesIVA']:,.0f}",
        "{{AHO_MES}}": f"${resultado['ahorroMensualEnergia']:,.0f}",
        "{{RETORNO}}": f"{resultado['tiempoRetorno']} años",
        "{{PORC_PR}}": f"{resultado['porcentajeProduccionMensual']}%",
        # Condiciones comerciales (NEW)
        "{{COND_COM}}": resultado.get("condicionesComerciales", ""),
        # Campos adicionales solicitados
        "{{NPISOS}}": str(req.get('numeroPisos', '1')),
        "{{HSPC}}": f"{req.get('hspCalculado') if req.get('hspCalculado') is not None else ''}",
        "{{AREA}}": f"{req.get('areaDisponible', 0) if req.get('areaDisponible') else 0} m²" if req.get('areaDisponible') else "N/A",
        "{{PCTDIA}}": f"{req.get('porcentajeConsumodia', '')}%",
        # Placeholders de legalización
        "{{NO_LEGALIZA}}": "Gestión de legalización ante operador del sistema instalado" if req.get("legalizacion", "NO") == "NO" else "",
        "{{SI_LEGALIZA}}": "Gestión de legalización ante operador de red, contador, documentos eléctricos y acompañamiento técnico en las 3 visitas." if req.get("legalizacion", "NO") == "SI" else ""
    }

def replace_text_in_shape(shape, mapping: dict):
    """Reemplaza texto en un shape individual manejando placeholders divididos entre runs"""
    if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
        return
    
    for paragraph in shape.text_frame.paragraphs:
        # PASO 1: Intentar reemplazo simple run por run (caso ideal)
        for run in paragraph.runs:
            if run.text:
                original_text = run.text
                new_text = original_text
                
                # Reemplazar placeholders completos dentro de este run
                for placeholder, value in mapping.items():
                    if placeholder in new_text:
                        new_text = new_text.replace(placeholder, value)
                
                if new_text != original_text:
                    run.text = new_text
        
        # PASO 2: Verificar si quedan placeholders divididos entre runs
        full_text = "".join(r.text for r in paragraph.runs)
        
        # Si hay placeholders en el texto completo pero no fueron reemplazados
        has_unreplaced = any(placeholder in full_text for placeholder in mapping.keys())
        
        if has_unreplaced:
            # Placeholder dividido entre runs - consolidar
            replaced_text = full_text
            for placeholder, value in mapping.items():
                replaced_text = replaced_text.replace(placeholder, value)
            
            if replaced_text != full_text:
                # Necesitamos consolidar: guardar formato del primer run
                if paragraph.runs:
                    first_run = paragraph.runs[0]
                    
                    # Limpiar todos los runs
                    for run in paragraph.runs:
                        run.text = ""
                    
                    # Poner todo el texto en el primer run
                    first_run.text = replaced_text

def replace_shape_text(shape, mapping: dict):
    """Reemplaza texto en shapes con estrategia simple de buscar-y-reemplazar"""
    from pptx.enum.text import PP_ALIGN
    
    replaced_count = 0
    
    # 1. Si tiene text_frame, procesar
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        replace_text_in_shape(shape, mapping)
    
    # 2. Si es tabla, procesar cada celda con enfoque run-by-run
    if hasattr(shape, 'has_table') and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                # Procesar cada párrafo y run de la celda
                for paragraph in cell.text_frame.paragraphs:
                    # PASO 1: Intentar reemplazo simple run por run
                    for run in paragraph.runs:
                        if run.text:
                            original_text = run.text
                            new_text = original_text
                            
                            # Reemplazar placeholders
                            for k, v in mapping.items():
                                if k in new_text:
                                    new_text = new_text.replace(k, v)
                                    replaced_count += 1
                            
                            # Actualizar solo si hubo cambios
                            if new_text != original_text:
                                run.text = new_text
                    
                    # PASO 2: Verificar placeholders divididos
                    full_text = "".join(r.text for r in paragraph.runs)
                    has_unreplaced = any(placeholder in full_text for placeholder in mapping.keys())
                    
                    if has_unreplaced:
                        # Placeholder dividido - consolidar
                        replaced_text = full_text
                        for k, v in mapping.items():
                            replaced_text = replaced_text.replace(k, v)
                            if k in full_text:
                                replaced_count += 1
                        
                        if replaced_text != full_text and paragraph.runs:
                            # Limpiar todos los runs
                            for run in paragraph.runs:
                                run.text = ""
                            # Poner texto en el primer run
                            paragraph.runs[0].text = replaced_text
                    
                    # Ajustar alineación para placeholders de totales
                    cell_text = "".join(r.text for r in paragraph.runs)
                    # Alinear celdas con valores monetarios a la derecha
                    if '$' in cell_text:
                        paragraph.alignment = PP_ALIGN.RIGHT
    
    return replaced_count

def fill_template_and_convert(req: dict, resultado: dict, opcion: str = "", template_path: str = None) -> tuple:
    """
    Llena template y convierte a PDF
    opcion: "" para una sola opción, "OPCIÓN 1" o "OPCIÓN 2 - Ajustada a área disponible"
    template_path: Ruta al template PPTX (usa TEMPLATE_PPTX por defecto)
    """
    # Usar template especificado o el por defecto
    template_to_use = template_path if template_path else TEMPLATE_PPTX
    
    if not os.path.isfile(template_to_use):
        raise RuntimeError(f"Template PPTX no encontrado: {template_to_use}")
    
    # Crear archivo temporal con nombre único que incluya timestamp
    timestamp_ms = int(datetime.now().timestamp() * 1000)
    opcion_suffix = "_op1" if opcion and "OPCIÓN 1" in opcion else "_op2" if opcion and "OPCIÓN 2" in opcion else ""
    fd_temp, filled_path = tempfile.mkstemp(
        suffix=f"_{timestamp_ms}{opcion_suffix}.pptx",
        prefix="cotizacion_"
    )
    os.close(fd_temp)
    
    # Copiar template original (FRESCO) al archivo temporal
    shutil.copy(template_to_use, filled_path)
    print(f"\n📄 ===== GENERANDO PDF {opcion if opcion else '(ÚNICA OPCIÓN)'} =====")
    print(f"   Template usado: {os.path.basename(template_to_use)}")
    print(f"   Template copiado a: {os.path.basename(filled_path)}")
    print(f"   Paneles: {resultado['numeroPaneles']}")
    print(f"   Capacidad: {resultado['capacidadInstalada']} kW")
    print(f"   Área requerida: {resultado['areaRequerida']} m²")
    print(f"   Valor total: ${resultado['valorTotalSistema']:,.0f}")

    # Cargar presentación desde la copia fresca
    prs = Presentation(filled_path)
    
    # Reemplazar placeholders (pasando el parámetro opcion)
    mapping = build_placeholders(req, resultado, opcion)
    total_replaced = 0
    print(f"🔄 Iniciando reemplazo de placeholders...")
    print(f"   Total de placeholders definidos: {len(mapping)}")
    print(f"   Opción: {opcion if opcion else '(única)'}") 
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n   📄 Procesando diapositiva {slide_idx + 1}...")
        # Procesar cada shape directamente (sin recursión)
        for shape in slide.shapes:
            count = replace_shape_text(shape, mapping)
            total_replaced += count
    
    print(f"\n✅ Total de reemplazos realizados: {total_replaced}")
    
    # Llenar tabla de ahorros (detecta automáticamente el número de filas)
    fill_ahorros_table_in_ppt(prs, resultado["tablaAhorros"])
    
    prs.save(filled_path)
    
    # Convertir a PDF con LibreOffice
    libreoffice_bin = os.getenv("LIBREOFFICE_PATH", "soffice")
    out_dir = tempfile.mkdtemp()
    cmd = [
        libreoffice_bin,
        "--headless",
        "--norestore",
        "--invisible",
        "--convert-to", "pdf",
        "--outdir", out_dir,
        filled_path
    ]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=90)
    if proc.returncode != 0:
        raise RuntimeError(f"LibreOffice error: {proc.stderr.decode('utf-8')}")
    
    pdf_name = os.path.splitext(os.path.basename(filled_path))[0] + ".pdf"
    pdf_path = os.path.join(out_dir, pdf_name)
    if not os.path.isfile(pdf_path):
        raise RuntimeError("Conversión a PDF fallida")
    
    return filled_path, pdf_path

# ========================================
# 📧 ENVÍO DE EMAIL
# ========================================
def enviar_email(destino: str, pdf_path: str, resultado: dict, pptx_path: Optional[str] = None):
    """Enviar email solo con PDF adjunto (compatibilidad antigua)"""
    return enviar_email_smtp([pdf_path], destino, resultado, 1)

def enviar_email_smtp(pdf_paths: list, destino: str, resultado: dict, num_opciones: int = 1):
    """
    Enviar email usando SMTP con soporte para múltiples PDFs
    pdf_paths: lista de rutas a PDFs (uno o más)
    num_opciones: 1 o 2 opciones de cotización
    """
    SMTP_HOST = os.getenv("SMTP_HOST")
    SMTP_PORT = int(os.getenv("SMTP_PORT", "587"))
    SMTP_USER = os.getenv("SMTP_USER")
    SMTP_PASS = os.getenv("SMTP_PASS")
    EMAIL_FROM = os.getenv("EMAIL_FROM", SMTP_USER)
    EMAIL_NASSA = os.getenv("EMAIL_NASSA", EMAIL_FROM)
    
    print(f"\n📧 CONFIGURACIÓN SMTP:")
    print(f"   SMTP_HOST: {SMTP_HOST}")
    print(f"   SMTP_PORT: {SMTP_PORT}")
    print(f"   EMAIL_FROM: {EMAIL_FROM}")
    print(f"   EMAIL_NASSA: {EMAIL_NASSA}")

    if not all([SMTP_HOST, SMTP_USER, SMTP_PASS, destino]):
        raise RuntimeError("❌ Configuración SMTP incompleta")

    # Mensaje personalizado según número de opciones
    mensaje_opciones = "Le presentamos 2 propuestas para su análisis" if num_opciones == 2 else "Hemos preparado una propuesta personalizada para tu proyecto solar"

    msg = EmailMessage()
    msg["Subject"] = f"Cotización NASSA Solar - {resultado['cotizacionId']}"
    msg["From"] = EMAIL_FROM
    msg["To"] = destino
    if EMAIL_NASSA and EMAIL_NASSA != destino:
        msg["Cc"] = EMAIL_NASSA

    cuerpo = f"""
Estimado/a {resultado.get('nombre', 'cliente')},

{mensaje_opciones}

📊 RESUMEN DE TU INVERSIÓN:
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📋 ID Cotización:      {resultado['cotizacionId']}
💰 Inversión Total:    ${resultado['valorTotalSistema']:,.0f} COP
💡 Ahorro Mensual:     ${resultado['ahorroMensualEnergia']:,.0f} COP
⏱️  Tiempo de Retorno:  {resultado['tiempoRetorno']} años
⚡ Capacidad Instalada: {resultado['capacidadInstalada']} kW
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

🌟 BENEFICIOS DE TU SISTEMA SOLAR:
✓ Ahorro inmediato en tu factura de energía
✓ Protección contra aumentos de tarifas
✓ Valorización de tu propiedad hasta 20%
✓ Contribución ambiental reduciendo CO₂

{'📎 Tus cotizaciones detalladas están adjuntas en formato PDF.' if num_opciones == 2 else '📎 Tu cotización detallada está adjunta en formato PDF.'}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
💬 CONTÁCTANOS:
📞 Teléfono: (057) 313 690 9723
🌐 Web: www.nassasolar.com
📧 Email: comercial@nassasolar.com

💬 WhatsApp: https://wa.me/573136909723?text=Hola,%20me%20interesa%20la%20cotización%20{resultado['cotizacionId']}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

Quedamos atentos a tus consultas.
Esta cotización tiene una validez de 30 días.

Saludos cordiales,
Equipo NASSA SOLAR
Expertos en Energía Solar Fotovoltaica
"""
    msg.set_content(cuerpo.strip())
    
    # Agregar versión HTML con colores anaranjados
    mensaje_adjunto = "Tus cotizaciones detalladas están adjuntas en formato PDF" if num_opciones == 2 else "Tu cotización detallada está adjunta en formato PDF"
    
    cuerpo_html = f"""
<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Cotización NASSA Solar</title>
</head>
<body style="margin: 0; padding: 0; font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 40px 20px;">
    
    <div style="max-width: 650px; margin: 0 auto; background: #ffffff; border-radius: 20px; overflow: hidden; box-shadow: 0 20px 60px rgba(0,0,0,0.3);">
        
        <!-- Header Anaranjado -->
        <div style="background: linear-gradient(135deg, #fbbf24 0%, #f59e0b 100%); padding: 40px 30px; text-align: center; position: relative;">
            <h1 style="margin: 0; color: #ffffff; font-size: 28px; font-weight: 700; text-shadow: 0 2px 8px rgba(0,0,0,0.3);">
                ☀️ NASSA SOLAR
            </h1>
            <p style="margin: 10px 0 0 0; color: #ffffff; font-size: 16px; font-weight: 600; letter-spacing: 2px;">
                ENERGÍA INTELIGENTE
            </p>
        </div>
        
        <!-- Contenido -->
        <div style="padding: 40px 30px;">
            
            <div style="text-align: center; margin-bottom: 30px;">
                <h2 style="color: #1f2937; font-size: 24px; margin: 0 0 10px 0;">
                    ✨ ¡Tu Cotización Está Lista!
                </h2>
                <p style="color: #6b7280; font-size: 16px; margin: 0; line-height: 1.6;">
                    {mensaje_opciones}
                </p>
            </div>
            
            <!-- Tarjeta de Resumen -->
            <div style="background: linear-gradient(135deg, #fef3c7 0%, #fde68a 100%); border-radius: 16px; padding: 25px; margin: 30px 0; border: 3px solid #fbbf24; box-shadow: 0 10px 30px rgba(251, 191, 36, 0.3);">
                <h3 style="color: #92400e; font-size: 20px; margin: 0 0 20px 0; text-align: center; font-weight: 700;">
                    📊 RESUMEN DE TU INVERSIÓN
                </h3>
                
                <table style="width: 100%; border-collapse: collapse;">
                    <tr>
                        <td style="padding: 12px 0; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">📋 ID Cotización</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #92400e; font-size: 16px; font-weight: 700;">{resultado['cotizacionId']}</span>
                        </td>
                    </tr>
                    <tr>
                        <td style="padding: 12px 0; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">💰 Inversión Total</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #dc2626; font-size: 20px; font-weight: 800;">${resultado['valorTotalSistema']:,.0f}</span>
                            <span style="color: #92400e; font-size: 14px; font-weight: 600;"> COP</span>
                        </td>
                    </tr>
                    <tr>
                        <td style="padding: 12px 0; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">💡 Ahorro Mensual</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #16a34a; font-size: 20px; font-weight: 800;">${resultado['ahorroMensualEnergia']:,.0f}</span>
                            <span style="color: #92400e; font-size: 14px; font-weight: 600;"> COP</span>
                        </td>
                    </tr>
                    <tr>
                        <td style="padding: 12px 0; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">⏱️ Tiempo de Retorno</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #2563eb; font-size: 20px; font-weight: 800;">{resultado['tiempoRetorno']}</span>
                            <span style="color: #92400e; font-size: 14px; font-weight: 600;"> años</span>
                        </td>
                    </tr>
                    <tr>
                        <td style="padding: 12px 0;">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">⚡ Capacidad Instalada</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right;">
                            <span style="color: #7c3aed; font-size: 20px; font-weight: 800;">{resultado['capacidadInstalada']}</span>
                            <span style="color: #92400e; font-size: 14px; font-weight: 600;"> kW</span>
                        </td>
                    </tr>
                </table>
            </div>
            
            <!-- Beneficios -->
            <div style="background: #f0fdf4; border-radius: 12px; padding: 20px; margin: 25px 0; border-left: 5px solid #16a34a;">
                <h4 style="color: #166534; margin: 0 0 15px 0; font-size: 18px;">🌟 Beneficios de tu Sistema Solar</h4>
                <ul style="margin: 0; padding-left: 20px; color: #15803d; line-height: 1.8;">
                    <li style="margin-bottom: 8px;"><strong>Ahorro inmediato</strong> en tu factura de energía</li>
                    <li style="margin-bottom: 8px;"><strong>Protección</strong> contra aumentos de tarifas</li>
                    <li style="margin-bottom: 8px;"><strong>Valorización</strong> de tu propiedad hasta 20%</li>
                    <li style="margin-bottom: 8px;"><strong>Contribución ambiental</strong> reduciendo CO₂</li>
                </ul>
            </div>
            
            <!-- Call to Action -->
            <div style="text-align: center; margin: 35px 0;">
                <table border="0" cellpadding="0" cellspacing="0" role="presentation" style="margin: 0 auto;">
                    <tr>
                        <td align="center" style="background: #16a34a; border-radius: 50px; box-shadow: 0 10px 25px rgba(22, 163, 74, 0.4);">
                            <a href="https://wa.me/573136909723?text=Hola,%20me%20interesa%20la%20cotización%20{resultado['cotizacionId']}" 
                               target="_blank"
                               style="background: #16a34a; border: 2px solid #16a34a; color: #ffffff; font-family: 'Segoe UI', Arial, sans-serif; font-size: 18px; font-weight: 700; line-height: 1.5; text-align: center; text-decoration: none; display: block; padding: 18px 45px; border-radius: 50px;">
                                <span style="color: #ffffff; text-decoration: none;">💬 Contáctanos por WhatsApp</span>
                            </a>
                        </td>
                    </tr>
                </table>
                <p style="color: #6b7280; font-size: 14px; margin-top: 15px;">
                    O llámanos al <strong style="color: #ea580c;">(057) 313 690 9723</strong>
                </p>
            </div>
            
            <!-- Archivo Adjunto -->
            <div style="background: #eff6ff; border-radius: 12px; padding: 20px; margin: 25px 0; text-align: center; border: 2px dashed #3b82f6;">
                <p style="margin: 0; color: #1e40af; font-size: 16px; font-weight: 600;">
                    📎 {mensaje_adjunto}
                </p>
            </div>
            
        </div>
        
        <!-- Footer Anaranjado -->
        <div style="background: linear-gradient(135deg, #fbbf24 0%, #f59e0b 100%); padding: 25px 30px; text-align: center; color: white;">
            <p style="margin: 0 0 8px 0; font-size: 18px; color: #ffffff; font-weight: 700;">
                ☀️ NASSA SOLAR
            </p>
            <p style="margin: 8px 0; font-size: 14px; color: #ffffff; font-weight: 600;">
                Expertos en Energía Solar Fotovoltaica
            </p>
            <p style="margin: 5px 0; font-size: 14px; color: #ffffff;">
                📞 (057) 313 690 9723 | 🌐 www.nassasolar.com
            </p>
            <p style="margin: 5px 0; font-size: 14px; color: #ffffff;">
                📧 comercial@nassasolar.com
            </p>
            
            <div style="margin-top: 18px; padding-top: 18px; border-top: 2px solid rgba(255, 255, 255, 0.3);">
                <p style="margin: 0; font-size: 12px; color: #ffffff; line-height: 1.6; opacity: 0.95;">
                    Esta cotización tiene una validez de 30 días.<br>
                    Precios sujetos a disponibilidad y condiciones del mercado.
                </p>
            </div>
        </div>
        
    </div>
    
    <div style="text-align: center; margin-top: 20px; padding: 0 20px;">
        <p style="color: rgba(255,255,255,0.8); font-size: 12px; margin: 0;">
            Has recibido este email porque solicitaste una cotización en nuestro sitio web.
        </p>
    </div>
    
</body>
</html>
"""
    msg.add_alternative(cuerpo_html, subtype='html')
    
    # Adjuntar todos los PDFs
    for i, pdf_path in enumerate(pdf_paths, 1):
        with open(pdf_path, "rb") as f:
            # Nombre del archivo según número de opciones
            if num_opciones == 1:
                filename = f"Cotizacion_{resultado['cotizacionId']}.pdf"
            else:
                filename = f"Cotizacion_{resultado['cotizacionId']}_Opcion{i}.pdf"
            
            msg.add_attachment(f.read(), maintype="application", subtype="pdf", filename=filename)

    # Construir lista de destinatarios
    destinatarios = [destino]
    if EMAIL_NASSA and EMAIL_NASSA != destino:
        destinatarios.append(EMAIL_NASSA)

    # Intentar primero puerto 465 (SSL) y luego 587 (STARTTLS)
    puertos_intentar = [SMTP_PORT, 465, 587] if SMTP_PORT not in [465, 587] else ([SMTP_PORT] if SMTP_PORT == 587 else [465, 587])
    
    print(f"\n📨 Enviando email SMTP...")
    print(f"   From: {EMAIL_FROM}")
    print(f"   To: {destino}")
    if EMAIL_NASSA and EMAIL_NASSA != destino:
        print(f"   CC: {EMAIL_NASSA}")
    print(f"   Attachments: {len(pdf_paths)}")
    print(f"   Puertos a intentar: {puertos_intentar}")
    
    ultimo_error = None
    for puerto in puertos_intentar:
        try:
            print(f"\n⏳ Intentando puerto {puerto}...")
            if puerto == 465:
                # Usar SMTP_SSL para puerto 465
                print(f"   Método: SMTP_SSL (conexión cifrada directa)")
                with smtplib.SMTP_SSL(SMTP_HOST, puerto, timeout=15) as server:
                    print(f"   ✓ Conectado a {SMTP_HOST}:{puerto}")
                    server.login(SMTP_USER, SMTP_PASS)
                    print(f"   ✓ Login exitoso")
                    server.send_message(msg)
                    print(f"   ✓ Mensaje enviado")
                print(f"✅ Email enviado via SMTP puerto {puerto} (SSL) a {destino}")
                return  # Salir si tuvo éxito
            else:
                # Usar SMTP con STARTTLS para puerto 587
                print(f"   Método: SMTP + STARTTLS")
                with smtplib.SMTP(SMTP_HOST, puerto, timeout=15) as server:
                    print(f"   ✓ Conectado a {SMTP_HOST}:{puerto}")
                    server.starttls()
                    print(f"   ✓ STARTTLS activado")
                    server.login(SMTP_USER, SMTP_PASS)
                    print(f"   ✓ Login exitoso")
                    server.send_message(msg)
                    print(f"   ✓ Mensaje enviado")
                print(f"✅ Email enviado via SMTP puerto {puerto} (STARTTLS) a {destino}")
                return  # Salir si tuvo éxito
        except socket.timeout:
            ultimo_error = f"Timeout ({15}s) conectando a {SMTP_HOST}:{puerto}"
            print(f"⚠️ {ultimo_error}")
        except smtplib.SMTPAuthenticationError as e:
            ultimo_error = f"Error de autenticación: {str(e)}"
            print(f"❌ {ultimo_error}")
            # Si falla auth, no intentar otros puertos (credenciales incorrectas)
            raise Exception(f"❌ Credenciales SMTP incorrectas: {ultimo_error}")
        except Exception as e:
            ultimo_error = f"Error en puerto {puerto}: {str(e)}"
            print(f"⚠️ {ultimo_error}")
    
    # Si llegamos aquí, ningún puerto funcionó
    raise Exception(f"❌ Error SMTP: No se pudo enviar email en ningún puerto. Último error: {ultimo_error}")

def enviar_email_original(destino: str, pdf_path: str, resultado: dict, pptx_path: Optional[str] = None):
    """Enviar email solo con PDF adjunto"""
    SMTP_HOST = os.getenv("SMTP_HOST")
    SMTP_PORT = int(os.getenv("SMTP_PORT", "587"))
    SMTP_USER = os.getenv("SMTP_USER")
    SMTP_PASS = os.getenv("SMTP_PASS")
    EMAIL_FROM = os.getenv("EMAIL_FROM", SMTP_USER)
    EMAIL_NASSA = os.getenv("EMAIL_NASSA", EMAIL_FROM)

    if not all([SMTP_HOST, SMTP_USER, SMTP_PASS, destino]):
        raise RuntimeError("SMTP incompleto")

    msg = EmailMessage()
    msg["Subject"] = f"Cotización NASSA Solar - {resultado['cotizacionId']}"
    msg["From"] = EMAIL_FROM
    msg["To"] = destino
    if EMAIL_NASSA and EMAIL_NASSA != destino:
        msg["Cc"] = EMAIL_NASSA

    cuerpo = f"""
Estimado cliente,

Adjuntamos su cotización personalizada:

📊 RESUMEN:
- ID: {resultado['cotizacionId']}
- Inversión: ${resultado['valorTotalSistema']:,.0f} COP
- Ahorro mensual: ${resultado['ahorroMensualEnergia']:,.0f} COP
- Retorno: {resultado['tiempoRetorno']} años
- Capacidad: {resultado['capacidadInstalada']} kW

Quedamos atentos a sus consultas.

Saludos cordiales,
NASSA SOLAR
Tel: (057) 313 690 9723
www.nassasolar.com
"""
    msg.set_content(cuerpo.strip())

    with open(pdf_path, "rb") as f:
        msg.add_attachment(f.read(), maintype="application", subtype="pdf",
                          filename=f"Cotizacion_{resultado['cotizacionId']}.pdf")

    # Construir lista de destinatarios para send_message
    destinatarios = [destino]
    if EMAIL_NASSA and EMAIL_NASSA != destino:
        destinatarios.append(EMAIL_NASSA)

    # Intentar primero puerto 465 (SSL) y luego 587 (STARTTLS)
    puertos_intentar = [465, SMTP_PORT] if SMTP_PORT != 465 else [465]
    
    for puerto in puertos_intentar:
        try:
            if puerto == 465:
                # Usar SMTP_SSL para puerto 465
                with smtplib.SMTP_SSL(SMTP_HOST, puerto, timeout=60) as server:
                    server.login(SMTP_USER, SMTP_PASS)
                    server.send_message(msg)
                print(f"✅ Email enviado via puerto {puerto} (SSL)")
                break
            else:
                # Usar SMTP con STARTTLS para puerto 587
                with smtplib.SMTP(SMTP_HOST, puerto, timeout=60) as server:
                    server.starttls()
                    server.login(SMTP_USER, SMTP_PASS)
                    server.send_message(msg)
                print(f"✅ Email enviado via puerto {puerto} (STARTTLS)")
                break
        except Exception as e:
            print(f"⚠️ Fallo puerto {puerto}: {str(e)}")
            if puerto == puertos_intentar[-1]:
                raise  # Si es el último puerto, propagar error
    
    # Log detallado de destinatarios
    if EMAIL_NASSA and EMAIL_NASSA != destino:
        print(f"   Destinatarios: {destino} | CC: {EMAIL_NASSA}")
    else:
        print(f"   Destinatario: {destino}")

def enviar_email_sendgrid(destino: str, pdf_paths: list, resultado: dict, num_opciones: int = 1, pptx_path: Optional[str] = None):
    """
    Enviar email usando SendGrid API (alternativa a SMTP)
    pdf_paths: lista de rutas a PDFs (uno o más)
    num_opciones: 1 o 2 opciones de cotización
    """
    import base64
    from sendgrid import SendGridAPIClient
    from sendgrid.helpers.mail import Mail, Attachment, FileContent, FileName, FileType, Disposition, ContentId
    
    SENDGRID_API_KEY = os.getenv("SENDGRID_API_KEY")
    EMAIL_FROM = os.getenv("EMAIL_FROM", "nassasolarprecotizacion@gmail.com")
    EMAIL_NASSA = os.getenv("EMAIL_NASSA", EMAIL_FROM)
    
    print(f"\n📧 CONFIGURACIÓN SENDGRID:")
    print(f"   API Key configurada: {'Sí' if SENDGRID_API_KEY else 'NO'}")
    print(f"   EMAIL_FROM: {EMAIL_FROM}")
    print(f"   EMAIL_NASSA: {EMAIL_NASSA}")
    
    if not SENDGRID_API_KEY:
        raise RuntimeError("❌ SENDGRID_API_KEY no configurada en variables de entorno")
    
    # Mensaje personalizado según número de opciones
    mensaje_opciones = "Le presentamos 2 propuestas para su análisis" if num_opciones == 2 else "Hemos preparado una propuesta personalizada para tu proyecto solar"
    mensaje_adjunto = "Sus cotizaciones detalladas están adjuntas en formato PDF" if num_opciones == 2 else "Tu cotización detallada está adjunta en formato PDF"
    
    cuerpo_html = f"""
<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Cotización NASSA Solar</title>
</head>
<body style="margin: 0; padding: 0; font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 40px 20px;">
    
    <!-- Contenedor Principal -->
    <div style="max-width: 650px; margin: 0 auto; background: #ffffff; border-radius: 20px; overflow: hidden; box-shadow: 0 20px 60px rgba(0,0,0,0.3);">
        
        <!-- Header con Gradiente -->
        <div style="background: linear-gradient(135deg, #fbbf24 0%, #f59e0b 100%); padding: 40px 30px; text-align: center; position: relative;">
            <img src="cid:logo_nassa" alt="NASSA Solar Logo" style="max-width: 220px; height: auto; margin-bottom: 15px; display: block; margin-left: auto; margin-right: auto;">
            <p style="margin: 0; color: #ffffff; font-size: 18px; font-weight: 700; letter-spacing: 3px; text-shadow: 0 2px 8px rgba(0,0,0,0.3);">
                ENERGÍA INTELIGENTE
            </p>
        </div>
        
        <!-- Contenido -->
        <div style="padding: 40px 30px;">
            
            <!-- Saludo -->
            <div style="text-align: center; margin-bottom: 30px;">
                <h2 style="color: #1f2937; font-size: 24px; margin: 0 0 10px 0;">
                    ✨ ¡Tu Cotización Está Lista!
                </h2>
                <p style="color: #6b7280; font-size: 16px; margin: 0; line-height: 1.6;">
                    {mensaje_opciones}
                </p>
            </div>
            
            <!-- Tarjeta de Resumen Premium -->
            <div style="background: linear-gradient(135deg, #fef3c7 0%, #fde68a 100%); border-radius: 16px; padding: 25px; margin: 30px 0; border: 3px solid #fbbf24; box-shadow: 0 10px 30px rgba(251, 191, 36, 0.3);">
                <h3 style="color: #92400e; font-size: 20px; margin: 0 0 20px 0; text-align: center; font-weight: 700;">
                    📊 RESUMEN DE TU INVERSIÓN
                </h3>
                
                <table style="width: 100%; border-collapse: collapse;">
                    <tr>
                        <td style="padding: 12px 0; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">📋 ID Cotización</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #92400e; font-size: 16px; font-weight: 700;">{resultado['cotizacionId']}</span>
                        </td>
                    </tr>
                    <tr>
                        <td style="padding: 12px 0; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">💰 Inversión Total</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #dc2626; font-size: 20px; font-weight: 800;">${resultado['valorTotalSistema']:,.0f}</span>
                            <span style="color: #92400e; font-size: 14px; font-weight: 600;"> COP</span>
                        </td>
                    </tr>
                    <tr>
                        <td style="padding: 12px 0; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">💡 Ahorro Mensual</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #16a34a; font-size: 20px; font-weight: 800;">${resultado['ahorroMensualEnergia']:,.0f}</span>
                            <span style="color: #92400e; font-size: 14px; font-weight: 600;"> COP</span>
                        </td>
                    </tr>
                    <tr>
                        <td style="padding: 12px 0; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">⏱️ Tiempo de Retorno</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right; border-bottom: 2px solid rgba(146, 64, 14, 0.2);">
                            <span style="color: #2563eb; font-size: 20px; font-weight: 800;">{resultado['tiempoRetorno']}</span>
                            <span style="color: #92400e; font-size: 14px; font-weight: 600;"> años</span>
                        </td>
                    </tr>
                    <tr>
                        <td style="padding: 12px 0;">
                            <span style="color: #78350f; font-size: 14px; font-weight: 600;">⚡ Capacidad Instalada</span>
                        </td>
                        <td style="padding: 12px 0; text-align: right;">
                            <span style="color: #7c3aed; font-size: 20px; font-weight: 800;">{resultado['capacidadInstalada']}</span>
                            <span style="color: #92400e; font-size: 14px; font-weight: 600;"> kW</span>
                        </td>
                    </tr>
                </table>
            </div>
            
            <!-- Beneficios -->
            <div style="background: #f0fdf4; border-radius: 12px; padding: 20px; margin: 25px 0; border-left: 5px solid #16a34a;">
                <h4 style="color: #166534; margin: 0 0 15px 0; font-size: 18px;">🌟 Beneficios de tu Sistema Solar</h4>
                <ul style="margin: 0; padding-left: 20px; color: #15803d; line-height: 1.8;">
                    <li style="margin-bottom: 8px;"><strong>Ahorro inmediato</strong> en tu factura de energía</li>
                    <li style="margin-bottom: 8px;"><strong>Protección</strong> contra aumentos de tarifas</li>
                    <li style="margin-bottom: 8px;"><strong>Valorización</strong> de tu propiedad hasta 20%</li>
                    <li style="margin-bottom: 8px;"><strong>Contribución ambiental</strong> reduciendo CO₂</li>
                </ul>
            </div>
            
            <!-- Call to Action -->
            <div style="text-align: center; margin: 35px 0;">
                <!-- Botón como tabla para máxima compatibilidad cross-client -->
                <table border="0" cellpadding="0" cellspacing="0" role="presentation" style="margin: 0 auto;">
                    <tr>
                        <td align="center" style="background: #16a34a; border-radius: 50px; box-shadow: 0 10px 25px rgba(22, 163, 74, 0.4);">
                            <a href="https://wa.me/573136909723?text=Hola,%20me%20interesa%20la%20cotización%20{resultado['cotizacionId']}" 
                               target="_blank"
                               style="background: #16a34a; border: 2px solid #16a34a; color: #ffffff; font-family: 'Segoe UI', Arial, sans-serif; font-size: 18px; font-weight: 700; line-height: 1.5; text-align: center; text-decoration: none; display: block; padding: 18px 45px; border-radius: 50px; -webkit-text-size-adjust: none; mso-hide: all;">
                                <span style="color: #ffffff; text-decoration: none;">💬 Contáctanos por WhatsApp</span>
                            </a>
                        </td>
                    </tr>
                </table>
                <p style="color: #6b7280; font-size: 14px; margin-top: 15px;">
                    O llámanos al <strong style="color: #ea580c;">(057) 313 690 9723</strong>
                </p>
            </div>
            
            <!-- Archivo Adjunto -->
            <div style="background: #eff6ff; border-radius: 12px; padding: 20px; margin: 25px 0; text-align: center; border: 2px dashed #3b82f6;">
                <p style="margin: 0; color: #1e40af; font-size: 16px; font-weight: 600;">
                    📎 {mensaje_adjunto}
                </p>
            </div>
            
        </div>
        
        <!-- Footer -->
        <div style="background: linear-gradient(135deg, #fbbf24 0%, #f59e0b 100%); padding: 25px 30px; text-align: center; color: white;">
            <img src="cid:logo_nassa" alt="NASSA Solar" style="max-width: 120px; height: auto; margin: 0 auto 10px auto; display: block; opacity: 0.95;">
            
            <p style="margin: 8px 0; font-size: 14px; color: #ffffff; font-weight: 600;">
                Expertos en Energía Solar Fotovoltaica
            </p>
            <p style="margin: 5px 0; font-size: 14px; color: #ffffff;">
                📞 (057) 313 690 9723 | 🌐 www.nassasolar.com
            </p>
            <p style="margin: 5px 0; font-size: 14px; color: #ffffff;">
                📧 comercial@nassasolar.com
            </p>
            
            <div style="margin-top: 18px; padding-top: 18px; border-top: 2px solid rgba(255, 255, 255, 0.3);">
                <p style="margin: 0; font-size: 12px; color: #ffffff; line-height: 1.6; opacity: 0.95;">
                    Esta cotización tiene una validez de 30 días.<br>
                    Precios sujetos a disponibilidad y condiciones del mercado.
                </p>
            </div>
        </div>
        
    </div>
    
    <!-- Nota de Privacidad -->
    <div style="text-align: center; margin-top: 20px; padding: 0 20px;">
        <p style="color: rgba(255,255,255,0.8); font-size: 12px; margin: 0;">
            Has recibido este email porque solicitaste una cotización en nuestro sitio web.
        </p>
    </div>
    
</body>
</html>
    """
    
    # Leer PDFs y crear attachments
    attachments = []
    for i, pdf_path in enumerate(pdf_paths, 1):
        with open(pdf_path, "rb") as f:
            pdf_data = base64.b64encode(f.read()).decode()
        
        attachment = Attachment()
        attachment.file_content = FileContent(pdf_data)
        
        # Nombre del archivo según número de opciones
        if num_opciones == 1:
            attachment.file_name = FileName(f"Cotizacion_{resultado['cotizacionId']}.pdf")
        else:
            attachment.file_name = FileName(f"Cotizacion_{resultado['cotizacionId']}_Opcion{i}.pdf")
        
        attachment.file_type = FileType("application/pdf")
        attachment.disposition = Disposition("attachment")
        attachments.append(attachment)
    
    # Adjuntar logo de NASSA como inline attachment (CID)
    logo_path = os.path.join(os.path.dirname(__file__), "..", "assets", "images", "loggo-Nassa.png")
    if os.path.exists(logo_path):
        with open(logo_path, "rb") as f:
            logo_data = base64.b64encode(f.read()).decode()
        
        logo_attachment = Attachment()
        logo_attachment.file_content = FileContent(logo_data)
        logo_attachment.file_name = FileName("logo-nassa.png")
        logo_attachment.file_type = FileType("image/png")
        logo_attachment.disposition = Disposition("inline")
        logo_attachment.content_id = ContentId("logo_nassa")
        attachments.append(logo_attachment)
    
    # Crear mensaje
    message = Mail(
        from_email=EMAIL_FROM,
        to_emails=[destino, EMAIL_NASSA],  # Envía a ambos
        subject=f"Cotización NASSA Solar - {resultado['cotizacionId']}",
        html_content=cuerpo_html
    )
    
    # Agregar todos los attachments
    for attachment in attachments:
        message.add_attachment(attachment)
    
    # Enviar
    sg = SendGridAPIClient(SENDGRID_API_KEY)
    
    print(f"\n📨 Enviando email...")
    print(f"   From: {EMAIL_FROM}")
    print(f"   To: {destino}")
    print(f"   CC: {EMAIL_NASSA}")
    print(f"   Subject: Cotización NASSA Solar - {resultado['cotizacionId']}")
    print(f"   Attachments: {len(attachments)}")
    
    response = sg.send(message)
    
    print(f"\n📬 Respuesta SendGrid:")
    print(f"   Status Code: {response.status_code}")
    print(f"   Headers: {response.headers}")
    
    if response.status_code in [200, 201, 202]:
        print(f"✅ Email enviado vía SendGrid a {destino}")
    else:
        print(f"❌ Error SendGrid: {response.status_code} - {response.body}")
        raise RuntimeError(f"SendGrid error: {response.status_code} - {response.body}")

def enviar_email_inteligente(destino: str, pdf_path: str, resultado: dict, pptx_path: Optional[str] = None):
    """Intenta SMTP primero (Gmail funcionando), SendGrid como fallback"""
    # Intentar Gmail SMTP primero (ya está funcionando)
    try:
        enviar_email(destino, pdf_path, resultado, pptx_path)
        return
    except Exception as e:
        print(f"⚠️ Gmail SMTP falló: {e}. Intentando SendGrid...")
        
        # Fallback a SendGrid si está configurado
        if os.getenv("SENDGRID_API_KEY"):
            try:
                enviar_email_sendgrid(destino, pdf_path, resultado, pptx_path)
                return
            except Exception as e2:
                print(f"❌ SendGrid también falló: {e2}")
                raise Exception(f"Error enviando email por ambos métodos: SMTP={e}, SendGrid={e2}")
        else:
            raise e  # Si no hay SendGrid configurado, relanzar error SMTP

# ========================================
# 🌐 ENDPOINTS
# ========================================
@app.get("/", tags=["General"])
def root():
    """Servir página principal"""
    index_path = os.path.join(STATIC_DIR, "index.html")
    if os.path.exists(index_path):
        return FileResponse(index_path, media_type="text/html")
    return {
        "message": "API NASSA Solar - Sistema de Cotización Fotovoltaico",
        "version": "1.0.0",
        "status": "activo"
    }

@app.get("/admin", tags=["General"])
def admin_panel():
    """Servir panel administrativo"""
    admin_path = os.path.join(STATIC_DIR, "admin.html")
    if os.path.exists(admin_path):
        return FileResponse(admin_path, media_type="text/html")
    raise HTTPException(404, "Panel administrativo no encontrado")

@app.get("/crm", tags=["General"])
def crm_panel():
    """Servir panel CRM"""
    crm_path = os.path.join(STATIC_DIR, "crm.html")
    if os.path.exists(crm_path):
        return FileResponse(crm_path, media_type="text/html")
    raise HTTPException(404, "Panel CRM no encontrado")

@app.get("/health", tags=["General"])
def health():
    return {"status": "ok", "timestamp": now_colombia().isoformat()}

@app.post("/api/admin/crear-tabla-cotizaciones", tags=["Admin"], dependencies=[Depends(auth_admin)])
def crear_tabla_cotizaciones_endpoint():
    """
    Crear tabla cotizaciones en PostgreSQL
    Ejecutar después de actualizar models.py con el modelo Cotizacion
    """
    try:
        if not os.getenv("DATABASE_URL"):
            raise HTTPException(500, "DATABASE_URL no configurada")
        
        from models import Cotizacion, create_db_engine
        
        print("🔄 Creando tabla cotizaciones...")
        engine = create_db_engine()
        
        # Crear solo la tabla Cotizacion
        Cotizacion.__table__.create(bind=engine, checkfirst=True)
        
        print("✅ Tabla cotizaciones creada!")
        
        return {
            "status": "success",
            "mensaje": "Tabla cotizaciones creada exitosamente",
            "timestamp": now_colombia().isoformat(),
            "campos": [
                "id", "fecha_creacion",
                "nombre", "email", "telefono", "direccion", "ciudad", "nic",
                "tipo_vivienda", "sistema_electrico", "tipo_sistema_fv",
                "consumo_mensual", "valor_factura", "valor_kwh", "porcentaje_consumo_dia", "hsp_calculado", "area_disponible",
                "panel_id", "panel_nombre", "inversor_id", "inversor_nombre", "bateria_id", "bateria_nombre",
                "num_paneles_op1", "capacidad_instalada_op1", "area_requerida_op1", "valor_total_op1", "ahorro_mensual_op1", "tiempo_retorno_op1",
                "tiene_opcion2", "num_paneles_op2", "capacidad_instalada_op2", "area_requerida_op2", "valor_total_op2", "ahorro_mensual_op2", "tiempo_retorno_op2",
                "datos_completos", "email_enviado", "fecha_envio_email", "num_opciones",
                "legalizacion", "seleccion_manual", "created_at", "updated_at"
            ]
        }
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR: {error_trace}")
        raise HTTPException(500, f"Error creando tabla: {str(e)}")

@app.post("/api/admin/init-database", tags=["Admin"], dependencies=[Depends(auth_admin)])
def init_database_endpoint():
    """
    Inicializa TODAS las tablas en PostgreSQL usando SQLAlchemy
    Incluye la nueva tabla de cotizaciones
    """
    try:
        if not os.getenv("DATABASE_URL"):
            raise HTTPException(500, "DATABASE_URL no configurada")
        
        from models import init_database
        
        print("🔄 Inicializando base de datos...")
        init_database()
        print("✅ Base de datos inicializada correctamente")
        
        return {
            "status": "success",
            "mensaje": "✅ Todas las tablas creadas exitosamente",
            "tablas": [
                "paneles", "inversores", "baterias", "ciudades", 
                "parametros", "consecutivos", "estadisticas", "cotizaciones"
            ]
        }
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR: {error_trace}")
        raise HTTPException(500, f"Error inicializando base de datos: {str(e)}")

@app.post("/api/admin/migrate-to-postgres", tags=["Admin"], dependencies=[Depends(auth_admin)])
def migrate_to_postgres():
    """
    Endpoint ONE-TIME para migrar datos de JSON a PostgreSQL
    ⚠️ EJECUTAR SOLO UNA VEZ después de crear la base de datos
    """
    try:
        # Verificar que DATABASE_URL exista
        if not os.getenv("DATABASE_URL"):
            raise HTTPException(500, "DATABASE_URL no configurada. Agrega PostgreSQL en Railway.")
        
        # Importar el script de migración
        import sys
        from pathlib import Path
        sys.path.insert(0, str(Path(__file__).parent))
        
        from migrate_to_postgres import migrate
        
        # Ejecutar migración
        success = migrate()
        
        if success:
            return {
                "status": "success",
                "mensaje": "¡Migración completada exitosamente!",
                "timestamp": now_colombia().isoformat(),
                "siguiente_paso": "Ahora debes actualizar server.py para usar PostgreSQL en lugar de JSON"
            }
        else:
            raise HTTPException(500, "La migración falló. Revisa los logs de Railway.")
            
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR EN MIGRACIÓN: {error_trace}")
        raise HTTPException(500, f"Error en migración: {str(e)}")

@app.get("/api/admin/verificar-postgres", tags=["Admin"], dependencies=[Depends(auth_admin)])
def verificar_postgres():
    """
    Endpoint para verificar conteos de registros en PostgreSQL
    """
    try:
        # Verificar que DATABASE_URL exista
        if not os.getenv("DATABASE_URL"):
            raise HTTPException(500, "DATABASE_URL no configurada")
        
        from models import get_db_session, Panel, Inversor, Bateria, Ciudad, Parametro, Consecutivo, Cotizacion
        
        session = get_db_session()
        
        try:
            # Contar registros
            paneles_count = session.query(Panel).count()
            inversores_count = session.query(Inversor).count()
            baterias_count = session.query(Bateria).count()
            ciudades_count = session.query(Ciudad).count()
            parametros_count = session.query(Parametro).count()
            consecutivo_count = session.query(Consecutivo).count()
            cotizaciones_count = session.query(Cotizacion).count()
            
            # Muestra de ciudades
            ciudades_muestra = session.query(Ciudad).limit(10).all()
            muestra_data = [
                {
                    "key": c.key,
                    "nombre": c.nombre,
                    "hsp": c.hsp
                }
                for c in ciudades_muestra
            ]
            
            return {
                "status": "success",
                "conteos": {
                    "paneles": paneles_count,
                    "inversores": inversores_count,
                    "baterias": baterias_count,
                    "ciudades": ciudades_count,
                    "parametros": parametros_count,
                    "consecutivos": consecutivo_count,
                    "cotizaciones": cotizaciones_count
                },
                "muestra_ciudades": muestra_data,
                "timestamp": now_colombia().isoformat()
            }
        finally:
            session.close()
            
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR EN VERIFICACIÓN: {error_trace}")
        raise HTTPException(500, f"Error verificando PostgreSQL: {str(e)}")

@app.post("/api/admin/limpiar-parametros-duplicados", tags=["Admin"], dependencies=[Depends(auth_admin)])
def limpiar_parametros_duplicados():
    """
    Endpoint para eliminar parámetros duplicados y dejar solo los últimos valores
    """
    try:
        if not os.getenv("DATABASE_URL"):
            raise HTTPException(500, "DATABASE_URL no configurada")
        
        from models import get_db_session, Parametro
        
        session = get_db_session()
        
        try:
            # Obtener todas las secciones únicas
            secciones = session.query(Parametro.seccion).distinct().all()
            secciones_unicas = [s[0] for s in secciones]
            
            total_antes = session.query(Parametro).count()
            
            # Para cada sección, mantener solo el último registro
            eliminados = 0
            actualizados = {}
            
            for seccion in secciones_unicas:
                # Obtener todos los registros de esta sección ordenados por ID
                registros = session.query(Parametro).filter_by(seccion=seccion).order_by(Parametro.id).all()
                
                if len(registros) > 1:
                    # Mantener el último, eliminar los demás
                    ultimo = registros[-1]
                    actualizados[seccion] = ultimo.data
                    
                    for registro in registros[:-1]:
                        session.delete(registro)
                        eliminados += 1
                else:
                    actualizados[seccion] = registros[0].data if registros else None
            
            session.commit()
            total_despues = session.query(Parametro).count()
            
            return {
                "status": "success",
                "mensaje": "Parámetros duplicados eliminados exitosamente",
                "total_antes": total_antes,
                "total_despues": total_despues,
                "eliminados": eliminados,
                "secciones_unicas": len(secciones_unicas),
                "secciones": list(actualizados.keys()),
                "timestamp": now_colombia().isoformat()
            }
        except Exception as e:
            session.rollback()
            raise e
        finally:
            session.close()
            
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR EN LIMPIEZA: {error_trace}")
        raise HTTPException(500, f"Error limpiando duplicados: {str(e)}")

@app.get("/debug/equipos-file", tags=["Debug"])
def debug_equipos_file():
    """Endpoint temporal para diagnosticar problema con equipos.json"""
    import hashlib
    try:
        with open(EQUIPOS_FILE, 'r', encoding='utf-8') as f:
            content = f.read()
            data = json.loads(content)
        
        # Hash del archivo para verificar versión
        file_hash = hashlib.md5(content.encode()).hexdigest()[:8]
        
        # Contar equipos con campo default
        paneles_con_default = sum(1 for p in data['paneles'] if 'default' in p)
        inversores_con_default = sum(1 for i in data['inversores'] if 'default' in i)
        baterias_con_default = sum(1 for b in data['baterias'] if 'default' in b)
        
        return {
            "file_path": EQUIPOS_FILE,
            "file_exists": os.path.exists(EQUIPOS_FILE),
            "file_hash": file_hash,
            "file_size": os.path.getsize(EQUIPOS_FILE),
            "file_mtime": os.path.getmtime(EQUIPOS_FILE),
            "paneles": {
                "total": len(data['paneles']),
                "con_campo_default": paneles_con_default
            },
            "inversores": {
                "total": len(data['inversores']),
                "con_campo_default": inversores_con_default
            },
            "baterias": {
                "total": len(data['baterias']),
                "con_campo_default": baterias_con_default
            },
            "sample_inversor": data['inversores'][1] if len(data['inversores']) > 1 else None
        }
    except Exception as e:
        return {"error": str(e), "file_path": EQUIPOS_FILE}

# ========================================
# 📊 CRM ENDPOINTS
# ========================================

@app.get("/api/admin/dashboard", tags=["CRM"], dependencies=[Depends(auth_admin)])
def crm_dashboard():
    """
    Dashboard administrativo con métricas principales del CRM
    """
    try:
        from models import get_db_session, Cotizacion
        from sqlalchemy import func, extract
        from datetime import datetime, timedelta
        
        session = get_db_session()
        
        try:
            now = now_colombia()
            mes_actual = now.month
            anio_actual = now.year
            
            # Total de cotizaciones
            total_cotizaciones = session.query(Cotizacion).count()
            
            # Cotizaciones este mes
            cotizaciones_mes = session.query(Cotizacion).filter(
                extract('month', Cotizacion.fecha_creacion) == mes_actual,
                extract('year', Cotizacion.fecha_creacion) == anio_actual
            ).count()
            
            # Cotizaciones este año
            cotizaciones_anio = session.query(Cotizacion).filter(
                extract('year', Cotizacion.fecha_creacion) == anio_actual
            ).count()
            
            # Tasa de conversión (emails enviados)
            emails_enviados = session.query(Cotizacion).filter(
                Cotizacion.email_enviado == True
            ).count()
            tasa_conversion = (emails_enviados / total_cotizaciones * 100) if total_cotizaciones > 0 else 0
            
            # Top 5 ciudades
            top_ciudades = session.query(
                Cotizacion.ciudad,
                func.count(Cotizacion.id).label('count')
            ).group_by(Cotizacion.ciudad).order_by(func.count(Cotizacion.id).desc()).limit(5).all()
            
            # Equipos más populares
            top_paneles = session.query(
                Cotizacion.panel_nombre,
                func.count(Cotizacion.id).label('count')
            ).filter(Cotizacion.panel_nombre.isnot(None)).group_by(
                Cotizacion.panel_nombre
            ).order_by(func.count(Cotizacion.id).desc()).limit(5).all()
            
            top_inversores = session.query(
                Cotizacion.inversor_nombre,
                func.count(Cotizacion.id).label('count')
            ).filter(Cotizacion.inversor_nombre.isnot(None)).group_by(
                Cotizacion.inversor_nombre
            ).order_by(func.count(Cotizacion.id).desc()).limit(5).all()
            
            # Promedios
            promedios = session.query(
                func.avg(Cotizacion.num_paneles_op1).label('avg_paneles'),
                func.avg(Cotizacion.capacidad_instalada_op1).label('avg_capacidad'),
                func.avg(Cotizacion.valor_total_op1).label('avg_valor'),
                func.avg(Cotizacion.tiempo_retorno_op1).label('avg_retorno'),
                func.avg(Cotizacion.ahorro_mensual_op1).label('avg_ahorro')
            ).first()
            
            # Adopción de Opción 2
            cotizaciones_con_op2 = session.query(Cotizacion).filter(
                Cotizacion.tiene_opcion2 == True
            ).count()
            tasa_opcion2 = (cotizaciones_con_op2 / total_cotizaciones * 100) if total_cotizaciones > 0 else 0
            
            # Cotizaciones recientes (últimas 10)
            recientes = session.query(Cotizacion).order_by(
                Cotizacion.fecha_creacion.desc()
            ).limit(10).all()
            
            cotizaciones_recientes = [
                {
                    "id": c.id,
                    "fecha": c.fecha_creacion.isoformat() if c.fecha_creacion else None,
                    "nombre": c.nombre,
                    "email": c.email,
                    "ciudad": c.ciudad,
                    "capacidad": f"{c.capacidad_instalada_op1:.2f} kW" if c.capacidad_instalada_op1 else "N/A",
                    "valor": f"${c.valor_total_op1:,.0f}" if c.valor_total_op1 else "N/A",
                    "email_enviado": c.email_enviado,
                    "num_opciones": c.num_opciones
                }
                for c in recientes
            ]
            
            return {
                "status": "success",
                "timestamp": now.isoformat(),
                "resumen": {
                    "total_cotizaciones": total_cotizaciones,
                    "cotizaciones_mes": cotizaciones_mes,
                    "cotizaciones_anio": cotizaciones_anio,
                    "emails_enviados": emails_enviados,
                    "tasa_conversion": round(tasa_conversion, 2)
                },
                "top_ciudades": [
                    {"ciudad": ciudad, "count": count}
                    for ciudad, count in top_ciudades
                ],
                "equipos_populares": {
                    "paneles": [
                        {"nombre": nombre, "count": count}
                        for nombre, count in top_paneles
                    ],
                    "inversores": [
                        {"nombre": nombre, "count": count}
                        for nombre, count in top_inversores
                    ]
                },
                "promedios": {
                    "paneles": round(promedios.avg_paneles, 1) if promedios.avg_paneles else 0,
                    "capacidad_kw": round(promedios.avg_capacidad, 2) if promedios.avg_capacidad else 0,
                    "valor_total": round(promedios.avg_valor, 0) if promedios.avg_valor else 0,
                    "tiempo_retorno_anos": round(promedios.avg_retorno, 1) if promedios.avg_retorno else 0,
                    "ahorro_mensual": round(promedios.avg_ahorro, 0) if promedios.avg_ahorro else 0
                },
                "tendencias": {
                    "cotizaciones_con_opcion2": cotizaciones_con_op2,
                    "tasa_opcion2_pct": round(tasa_opcion2, 2)
                },
                "cotizaciones_recientes": cotizaciones_recientes
            }
            
        finally:
            session.close()
            
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR EN DASHBOARD: {error_trace}")
        raise HTTPException(500, f"Error generando dashboard: {str(e)}")

@app.get("/api/admin/cotizaciones/buscar", tags=["CRM"], dependencies=[Depends(auth_admin)])
def buscar_cotizaciones(
    nombre: str = None,
    email: str = None,
    telefono: str = None,
    ciudad: str = None,
    fecha_desde: str = None,
    fecha_hasta: str = None,
    email_enviado: bool = None,
    pagina: int = 1,
    por_pagina: int = 50
):
    """
    Buscar cotizaciones con filtros múltiples
    
    Query params:
    - nombre: Búsqueda parcial por nombre de cliente
    - email: Búsqueda parcial por email
    - telefono: Búsqueda parcial por teléfono
    - ciudad: Búsqueda exacta por ciudad
    - fecha_desde: Formato ISO (YYYY-MM-DD)
    - fecha_hasta: Formato ISO (YYYY-MM-DD)
    - email_enviado: true/false
    - pagina: Número de página (default 1)
    - por_pagina: Resultados por página (default 50, max 200)
    """
    try:
        from models import get_db_session, Cotizacion
        from sqlalchemy import and_, or_
        
        session = get_db_session()
        
        try:
            # Construir query base
            query = session.query(Cotizacion)
            
            # Aplicar filtros
            filtros = []
            
            if nombre:
                filtros.append(Cotizacion.nombre.ilike(f"%{nombre}%"))
            
            if email:
                filtros.append(Cotizacion.email.ilike(f"%{email}%"))
            
            if telefono:
                filtros.append(Cotizacion.telefono.ilike(f"%{telefono}%"))
            
            if ciudad:
                filtros.append(Cotizacion.ciudad == ciudad)
            
            if fecha_desde:
                try:
                    from datetime import datetime
                    fecha_desde_dt = datetime.fromisoformat(fecha_desde)
                    filtros.append(Cotizacion.fecha_creacion >= fecha_desde_dt)
                except ValueError:
                    raise HTTPException(400, "fecha_desde debe estar en formato ISO (YYYY-MM-DD)")
            
            if fecha_hasta:
                try:
                    from datetime import datetime, timedelta
                    fecha_hasta_dt = datetime.fromisoformat(fecha_hasta) + timedelta(days=1)
                    filtros.append(Cotizacion.fecha_creacion < fecha_hasta_dt)
                except ValueError:
                    raise HTTPException(400, "fecha_hasta debe estar en formato ISO (YYYY-MM-DD)")
            
            if email_enviado is not None:
                filtros.append(Cotizacion.email_enviado == email_enviado)
            
            # Aplicar todos los filtros
            if filtros:
                query = query.filter(and_(*filtros))
            
            # Contar total
            total = query.count()
            
            # Validar paginación
            por_pagina = min(por_pagina, 200)  # Max 200 por página
            offset = (pagina - 1) * por_pagina
            
            # Obtener resultados paginados
            resultados = query.order_by(
                Cotizacion.fecha_creacion.desc()
            ).offset(offset).limit(por_pagina).all()
            
            cotizaciones = [
                {
                    "id": c.id,
                    "fecha_creacion": c.fecha_creacion.isoformat() if c.fecha_creacion else None,
                    "nombre": c.nombre,
                    "email": c.email,
                    "telefono": c.telefono,
                    "ciudad": c.ciudad,
                    "tipo_sistema_fv": c.tipo_sistema_fv,
                    "consumo_mensual": c.consumo_mensual,
                    "num_paneles": c.num_paneles_op1,
                    "capacidad_instalada": f"{c.capacidad_instalada_op1:.2f} kW" if c.capacidad_instalada_op1 else None,
                    "valor_total": c.valor_total_op1,
                    "tiempo_retorno": c.tiempo_retorno_op1,
                    "email_enviado": c.email_enviado,
                    "fecha_envio_email": c.fecha_envio_email.isoformat() if c.fecha_envio_email else None,
                    "tiene_opcion2": c.tiene_opcion2,
                    "num_opciones": c.num_opciones
                }
                for c in resultados
            ]
            
            total_paginas = (total + por_pagina - 1) // por_pagina
            
            return {
                "status": "success",
                "resultados": cotizaciones,
                "paginacion": {
                    "pagina_actual": pagina,
                    "por_pagina": por_pagina,
                    "total_resultados": total,
                    "total_paginas": total_paginas
                },
                "filtros_aplicados": {
                    "nombre": nombre,
                    "email": email,
                    "telefono": telefono,
                    "ciudad": ciudad,
                    "fecha_desde": fecha_desde,
                    "fecha_hasta": fecha_hasta,
                    "email_enviado": email_enviado
                }
            }
            
        finally:
            session.close()
            
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR EN BÚSQUEDA: {error_trace}")
        raise HTTPException(500, f"Error buscando cotizaciones: {str(e)}")

@app.get("/api/admin/cotizaciones/{cotizacion_id}", tags=["CRM"], dependencies=[Depends(auth_admin)])
def detalle_cotizacion(cotizacion_id: str):
    """
    Obtener detalle completo de una cotización específica
    """
    try:
        from models import get_db_session, Cotizacion
        
        session = get_db_session()
        
        try:
            cotizacion = session.query(Cotizacion).filter(
                Cotizacion.id == cotizacion_id
            ).first()
            
            if not cotizacion:
                raise HTTPException(404, f"Cotización {cotizacion_id} no encontrada")
            
            # Construir respuesta completa
            detalle = {
                "id": cotizacion.id,
                "fecha_creacion": cotizacion.fecha_creacion.isoformat() if cotizacion.fecha_creacion else None,
                
                # Cliente
                "cliente": {
                    "nombre": cotizacion.nombre,
                    "email": cotizacion.email,
                    "telefono": cotizacion.telefono,
                    "direccion": cotizacion.direccion,
                    "ciudad": cotizacion.ciudad,
                    "nic": cotizacion.nic
                },
                
                # Sistema
                "sistema": {
                    "tipo_vivienda": cotizacion.tipo_vivienda,
                    "sistema_electrico": cotizacion.sistema_electrico,
                    "tipo_sistema_fv": cotizacion.tipo_sistema_fv
                },
                
                # Consumo
                "consumo": {
                    "consumo_mensual": cotizacion.consumo_mensual,
                    "valor_factura": cotizacion.valor_factura,
                    "valor_kwh": cotizacion.valor_kwh,
                    "porcentaje_consumo_dia": cotizacion.porcentaje_consumo_dia,
                    "hsp_calculado": cotizacion.hsp_calculado,
                    "area_disponible": cotizacion.area_disponible
                },
                
                # Equipos
                "equipos": {
                    "panel": {
                        "id": cotizacion.panel_id,
                        "nombre": cotizacion.panel_nombre
                    },
                    "inversor": {
                        "id": cotizacion.inversor_id,
                        "nombre": cotizacion.inversor_nombre
                    },
                    "bateria": {
                        "id": cotizacion.bateria_id,
                        "nombre": cotizacion.bateria_nombre
                    } if cotizacion.bateria_id else None
                },
                
                # Opción 1
                "opcion1": {
                    "num_paneles": cotizacion.num_paneles_op1,
                    "capacidad_instalada": cotizacion.capacidad_instalada_op1,
                    "area_requerida": cotizacion.area_requerida_op1,
                    "valor_total": cotizacion.valor_total_op1,
                    "ahorro_mensual": cotizacion.ahorro_mensual_op1,
                    "tiempo_retorno": cotizacion.tiempo_retorno_op1
                },
                
                # Opción 2 (si existe)
                "opcion2": {
                    "tiene_opcion2": cotizacion.tiene_opcion2,
                    "num_paneles": cotizacion.num_paneles_op2,
                    "capacidad_instalada": cotizacion.capacidad_instalada_op2,
                    "area_requerida": cotizacion.area_requerida_op2,
                    "valor_total": cotizacion.valor_total_op2,
                    "ahorro_mensual": cotizacion.ahorro_mensual_op2,
                    "tiempo_retorno": cotizacion.tiempo_retorno_op2
                } if cotizacion.tiene_opcion2 else None,
                
                # Estado
                "estado": {
                    "email_enviado": cotizacion.email_enviado,
                    "fecha_envio_email": cotizacion.fecha_envio_email.isoformat() if cotizacion.fecha_envio_email else None,
                    "num_opciones": cotizacion.num_opciones
                },
                
                # Metadata
                "metadata": {
                    "legalizacion": cotizacion.legalizacion,
                    "seleccion_manual": cotizacion.seleccion_manual,
                    "created_at": cotizacion.created_at.isoformat() if cotizacion.created_at else None,
                    "updated_at": cotizacion.updated_at.isoformat() if cotizacion.updated_at else None
                },
                
                # Datos completos JSON
                "datos_completos": cotizacion.datos_completos
            }
            
            return {
                "status": "success",
                "cotizacion": detalle
            }
            
        finally:
            session.close()
            
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR EN DETALLE: {error_trace}")
        raise HTTPException(500, f"Error obteniendo detalle: {str(e)}")

@app.get("/api/admin/reportes/top-ciudades", tags=["CRM"], dependencies=[Depends(auth_admin)])
def reporte_top_ciudades(limit: int = 10):
    """
    Reporte de ciudades con más cotizaciones
    """
    try:
        from models import get_db_session, Cotizacion
        from sqlalchemy import func
        
        session = get_db_session()
        
        try:
            top_ciudades = session.query(
                Cotizacion.ciudad,
                func.count(Cotizacion.id).label('total_cotizaciones'),
                func.sum(Cotizacion.valor_total_op1).label('valor_total'),
                func.avg(Cotizacion.capacidad_instalada_op1).label('capacidad_promedio'),
                func.sum(Cotizacion.num_paneles_op1).label('total_paneles')
            ).filter(
                Cotizacion.ciudad.isnot(None)
            ).group_by(
                Cotizacion.ciudad
            ).order_by(
                func.count(Cotizacion.id).desc()
            ).limit(limit).all()
            
            resultados = [
                {
                    "ciudad": ciudad,
                    "total_cotizaciones": total,
                    "valor_total_acumulado": round(float(valor), 0) if valor else 0,
                    "capacidad_promedio_kw": round(float(capacidad), 2) if capacidad else 0,
                    "total_paneles": int(paneles) if paneles else 0
                }
                for ciudad, total, valor, capacidad, paneles in top_ciudades
            ]
            
            return {
                "status": "success",
                "reporte": "Top Ciudades",
                "timestamp": now_colombia().isoformat(),
                "resultados": resultados
            }
            
        finally:
            session.close()
            
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR EN REPORTE: {error_trace}")
        raise HTTPException(500, f"Error generando reporte: {str(e)}")

@app.get("/api/admin/reportes/estadisticas", tags=["CRM"], dependencies=[Depends(auth_admin)])
def reporte_estadisticas():
    """
    Reporte de estadísticas generales del sistema
    """
    try:
        from models import get_db_session, Cotizacion
        from sqlalchemy import func
        
        session = get_db_session()
        
        try:
            # Estadísticas generales
            stats = session.query(
                func.count(Cotizacion.id).label('total'),
                func.avg(Cotizacion.num_paneles_op1).label('avg_paneles'),
                func.avg(Cotizacion.capacidad_instalada_op1).label('avg_capacidad'),
                func.avg(Cotizacion.valor_total_op1).label('avg_valor'),
                func.avg(Cotizacion.tiempo_retorno_op1).label('avg_retorno'),
                func.sum(Cotizacion.valor_total_op1).label('valor_total'),
                func.min(Cotizacion.valor_total_op1).label('valor_min'),
                func.max(Cotizacion.valor_total_op1).label('valor_max')
            ).first()
            
            # Adopción de Opción 2
            con_opcion2 = session.query(Cotizacion).filter(
                Cotizacion.tiene_opcion2 == True
            ).count()
            
            # Por tipo de sistema
            por_tipo_sistema = session.query(
                Cotizacion.tipo_sistema_fv,
                func.count(Cotizacion.id).label('count')
            ).group_by(
                Cotizacion.tipo_sistema_fv
            ).all()
            
            # Por tipo de vivienda
            por_tipo_vivienda = session.query(
                Cotizacion.tipo_vivienda,
                func.count(Cotizacion.id).label('count')
            ).group_by(
                Cotizacion.tipo_vivienda
            ).all()
            
            return {
                "status": "success",
                "timestamp": now_colombia().isoformat(),
                "estadisticas_generales": {
                    "total_cotizaciones": stats.total or 0,
                    "promedio_paneles": round(float(stats.avg_paneles), 1) if stats.avg_paneles else 0,
                    "promedio_capacidad_kw": round(float(stats.avg_capacidad), 2) if stats.avg_capacidad else 0,
                    "promedio_valor": round(float(stats.avg_valor), 0) if stats.avg_valor else 0,
                    "promedio_tiempo_retorno_anos": round(float(stats.avg_retorno), 1) if stats.avg_retorno else 0,
                    "valor_total_mercado": round(float(stats.valor_total), 0) if stats.valor_total else 0,
                    "valor_min": round(float(stats.valor_min), 0) if stats.valor_min else 0,
                    "valor_max": round(float(stats.valor_max), 0) if stats.valor_max else 0
                },
                "adopcion_opcion2": {
                    "cotizaciones_con_opcion2": con_opcion2,
                    "porcentaje": round((con_opcion2 / stats.total * 100), 2) if stats.total else 0
                },
                "por_tipo_sistema": [
                    {"tipo": tipo, "cantidad": count}
                    for tipo, count in por_tipo_sistema
                ],
                "por_tipo_vivienda": [
                    {"tipo": tipo, "cantidad": count}
                    for tipo, count in por_tipo_vivienda
                ]
            }
            
        finally:
            session.close()
            
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR EN ESTADÍSTICAS: {error_trace}")
        raise HTTPException(500, f"Error generando estadísticas: {str(e)}")

@app.get("/api/admin/reportes/export", tags=["CRM"], dependencies=[Depends(auth_admin)])
def export_cotizaciones(
    fecha_desde: str = None,
    fecha_hasta: str = None,
    formato: str = "csv"
):
    """
    Exportar cotizaciones a CSV o JSON
    
    Query params:
    - fecha_desde: Formato ISO (YYYY-MM-DD)
    - fecha_hasta: Formato ISO (YYYY-MM-DD)
    - formato: "csv" o "json" (default: csv)
    """
    try:
        from models import get_db_session, Cotizacion
        from sqlalchemy import and_
        import io
        import csv
        
        session = get_db_session()
        
        try:
            # Construir query
            query = session.query(Cotizacion)
            
            filtros = []
            if fecha_desde:
                from datetime import datetime
                fecha_desde_dt = datetime.fromisoformat(fecha_desde)
                filtros.append(Cotizacion.fecha_creacion >= fecha_desde_dt)
            
            if fecha_hasta:
                from datetime import datetime, timedelta
                fecha_hasta_dt = datetime.fromisoformat(fecha_hasta) + timedelta(days=1)
                filtros.append(Cotizacion.fecha_creacion < fecha_hasta_dt)
            
            if filtros:
                query = query.filter(and_(*filtros))
            
            cotizaciones = query.order_by(Cotizacion.fecha_creacion.desc()).all()
            
            if formato == "json":
                # Exportar como JSON
                datos = [
                    {
                        "id": c.id,
                        "fecha_creacion": c.fecha_creacion.isoformat() if c.fecha_creacion else None,
                        "nombre": c.nombre,
                        "email": c.email,
                        "telefono": c.telefono,
                        "ciudad": c.ciudad,
                        "consumo_mensual": c.consumo_mensual,
                        "num_paneles": c.num_paneles_op1,
                        "capacidad_instalada": c.capacidad_instalada_op1,
                        "valor_total": c.valor_total_op1,
                        "ahorro_mensual": c.ahorro_mensual_op1,
                        "tiempo_retorno": c.tiempo_retorno_op1,
                        "email_enviado": c.email_enviado,
                        "tiene_opcion2": c.tiene_opcion2
                    }
                    for c in cotizaciones
                ]
                
                return {
                    "status": "success",
                    "total": len(datos),
                    "datos": datos
                }
            
            else:
                # Exportar como CSV
                output = io.StringIO()
                writer = csv.writer(output)
                
                # Headers
                writer.writerow([
                    "ID", "Fecha", "Nombre", "Email", "Teléfono", "Ciudad",
                    "Consumo (kWh/mes)", "Num Paneles", "Capacidad (kW)",
                    "Valor Total", "Ahorro Mensual", "Tiempo Retorno (años)",
                    "Email Enviado", "Tiene Opción 2"
                ])
                
                # Data
                for c in cotizaciones:
                    writer.writerow([
                        c.id,
                        c.fecha_creacion.strftime("%Y-%m-%d %H:%M") if c.fecha_creacion else "",
                        c.nombre or "",
                        c.email or "",
                        c.telefono or "",
                        c.ciudad or "",
                        c.consumo_mensual or "",
                        c.num_paneles_op1 or "",
                        round(c.capacidad_instalada_op1, 2) if c.capacidad_instalada_op1 else "",
                        round(c.valor_total_op1, 0) if c.valor_total_op1 else "",
                        round(c.ahorro_mensual_op1, 0) if c.ahorro_mensual_op1 else "",
                        round(c.tiempo_retorno_op1, 1) if c.tiempo_retorno_op1 else "",
                        "Sí" if c.email_enviado else "No",
                        "Sí" if c.tiene_opcion2 else "No"
                    ])
                
                from fastapi.responses import StreamingResponse
                output.seek(0)
                
                return StreamingResponse(
                    iter([output.getvalue()]),
                    media_type="text/csv",
                    headers={
                        "Content-Disposition": f"attachment; filename=cotizaciones_{now_colombia().strftime('%Y%m%d')}.csv"
                    }
                )
            
        finally:
            session.close()
            
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ ERROR EN EXPORTACIÓN: {error_trace}")
        raise HTTPException(500, f"Error exportando datos: {str(e)}")

@app.get("/api/equipos", tags=["Equipos"])
def equipos_publicos(sistemaElectrico: str = None):
    """
    Obtiene equipos disponibles (sin precios).
    
    Args:
        sistemaElectrico: Filtra inversores por tipo (monofasico, bifasico, trifasico).
                         Si no se proporciona, devuelve todos los inversores.
    """
    from models import get_db_session, Panel, Inversor, Bateria
    
    session = get_db_session()
    try:
        # Obtener paneles desde PostgreSQL
        paneles_db = session.query(Panel).all()
        paneles = [
            {
                "id": p.id,
                "nombre": p.nombre,
                "capacidad": p.capacidad,
                "descripcion": p.descripcion,
                "eficienciaPanel": p.eficienciaPanel if hasattr(p, 'eficienciaPanel') else 0.90,
                "area": p.area if hasattr(p, 'area') else 2.0,
                "default": p.default if hasattr(p, 'default') else False
            }
            for p in paneles_db
        ]
        
        # Obtener inversores con filtro opcional
        inversores_db = session.query(Inversor).all()
        
        # Filtrar en Python si se especifica sistema eléctrico
        if sistemaElectrico:
            sistema_normalizado = sistemaElectrico.lower().strip()
            inversores_filtrados = []
            for i in inversores_db:
                # sistemaElectrico puede ser string o array JSON
                if isinstance(i.sistemaElectrico, list):
                    # Es un array: ["monofasico", "trifasico"]
                    if sistema_normalizado in [s.lower() for s in i.sistemaElectrico]:
                        inversores_filtrados.append(i)
                elif isinstance(i.sistemaElectrico, str):
                    # Es un string: "monofasico"
                    if sistema_normalizado == i.sistemaElectrico.lower():
                        inversores_filtrados.append(i)
            inversores_db = inversores_filtrados
        
        inversores = [
            {
                "id": i.id,
                "nombre": i.nombre,
                "capacidad": i.capacidad,
                "descripcion": i.descripcion,
                "tipo_sistema": i.sistemaElectrico,
                "eficiencia": i.eficiencia if hasattr(i, 'eficiencia') else 0.96,
                "tipo": i.tipo if hasattr(i, 'tipo') else "STRING",
                "paneles_por_inversor": i.paneles_por_inversor if hasattr(i, 'paneles_por_inversor') else 0,
                "sobredimensionamiento": i.sobredimensionamiento if hasattr(i, 'sobredimensionamiento') else 0.40,
                "default": i.default if hasattr(i, 'default') else False
            }
            for i in inversores_db
        ]
        
        # Obtener baterías
        baterias_db = session.query(Bateria).all()
        baterias = [
            {
                "id": b.id,
                "nombre": b.nombre,
                "capacidad": b.capacidad,
                "descripcion": b.descripcion,
                "default": b.default if hasattr(b, 'default') else False
            }
            for b in baterias_db
        ]
        
        return {
            "paneles": paneles,
            "inversores": inversores,
            "baterias": baterias
        }
    finally:
        session.close()

@app.get("/api/equipos/precios", tags=["Equipos"], dependencies=[Depends(auth_admin)])
def equipos_con_precios():
    """Obtiene todos los equipos CON precios (admin only)"""
    from models import get_db_session, Panel, Inversor, Bateria
    
    session = get_db_session()
    try:
        # Paneles
        paneles_db = session.query(Panel).all()
        paneles = [
            {
                "id": p.id,
                "nombre": p.nombre,
                "capacidad": p.capacidad,
                "precio": p.precio,
                "descripcion": p.descripcion,
                "eficienciaPanel": p.eficienciaPanel,
                "area": p.area,
                "default": p.default
            }
            for p in paneles_db
        ]
        
        # Inversores
        inversores_db = session.query(Inversor).all()
        inversores = [
            {
                "id": i.id,
                "nombre": i.nombre,
                "capacidad": i.capacidad,
                "precio": i.precio,
                "descripcion": i.descripcion,
                "eficiencia": i.eficiencia,
                "sistemaElectrico": i.sistemaElectrico,
                "tipo": i.tipo,
                "paneles_por_inversor": i.paneles_por_inversor,
                "sobredimensionamiento": i.sobredimensionamiento,
                "default": i.default
            }
            for i in inversores_db
        ]
        
        # Baterías
        baterias_db = session.query(Bateria).all()
        baterias = [
            {
                "id": b.id,
                "nombre": b.nombre,
                "capacidad": b.capacidad,
                "precio": b.precio,
                "descripcion": b.descripcion,
                "default": b.default
            }
            for b in baterias_db
        ]
        
        return {
            "paneles": paneles,
            "inversores": inversores,
            "baterias": baterias
        }
    finally:
        session.close()

@app.get("/api/ciudades", tags=["Configuración"])
def ciudades():
    """Obtiene todas las ciudades con sus valores HSP"""
    from models import get_db_session, Ciudad
    
    session = get_db_session()
    try:
        ciudades_db = session.query(Ciudad).order_by(Ciudad.nombre).all()
        
        # Convertir a formato diccionario con key como índice
        ciudades_dict = {}
        for c in ciudades_db:
            ciudades_dict[c.key] = {
                "nombre": c.nombre,
                "hsp": c.hsp
            }
        
        return ciudades_dict
    finally:
        session.close()

@app.get("/api/template/status", tags=["Template"])
def template_status():
    return {"available": os.path.isfile(TEMPLATE_PPTX)}

@app.get("/api/template/download", tags=["Template"], dependencies=[Depends(auth_admin)])
def template_download():
    if not os.path.isfile(TEMPLATE_PPTX):
        raise HTTPException(404, "Template no encontrado")
    return FileResponse(TEMPLATE_PPTX, filename=os.path.basename(TEMPLATE_PPTX))

# ========================================
# 🔧 ENDPOINTS DE ADMINISTRACIÓN
# ========================================

# --- GESTIÓN DE PARÁMETROS DE COSTOS ---
@app.get("/api/admin/parametros", tags=["Admin"], dependencies=[Depends(auth_admin)])
def get_parametros():
    """Obtener todos los parámetros de costos y fiscales"""
    from models import get_db_session, Parametro
    
    session = get_db_session()
    try:
        parametros_db = session.query(Parametro).all()
        
        # Reconstruir el diccionario completo de parámetros
        parametros_dict = {}
        for p in parametros_db:
            parametros_dict[p.seccion] = p.data
        
        return parametros_dict
    finally:
        session.close()

@app.put("/api/admin/parametros", tags=["Admin"], dependencies=[Depends(auth_admin)])
def update_parametros(parametros: dict):
    """Actualizar parámetros de costos y fiscales"""
    from models import get_db_session, Parametro
    
    session = get_db_session()
    try:
        # Actualizar cada sección de parámetros
        for seccion, data in parametros.items():
            # Buscar si existe
            param = session.query(Parametro).filter_by(seccion=seccion).first()
            
            if param:
                # Actualizar existente
                param.data = data
            else:
                # Crear nuevo
                param = Parametro(seccion=seccion, data=data)
                session.add(param)
        
        session.commit()
        return {
            "status": "success",
            "mensaje": "Parámetros actualizados exitosamente en PostgreSQL",
            "secciones_actualizadas": list(parametros.keys())
        }
    except Exception as e:
        session.rollback()
        raise HTTPException(500, f"Error al actualizar parámetros: {e}")
    finally:
        session.close()

# --- GESTIÓN DE PANELES ---
@app.get("/api/admin/paneles", tags=["Admin"], dependencies=[Depends(auth_admin)])
def get_paneles_admin():
    """Obtener todos los paneles con precios (admin)"""
    from models import get_db_session, Panel
    
    session = get_db_session()
    try:
        paneles_db = session.query(Panel).all()
        paneles = [
            {
                "id": p.id,
                "nombre": p.nombre,
                "capacidad": p.capacidad,
                "precio": p.precio,
                "descripcion": p.descripcion,
                "eficienciaPanel": p.eficienciaPanel,
                "area": p.area,
                "default": p.default
            }
            for p in paneles_db
        ]
        return paneles
    finally:
        session.close()

@app.post("/api/admin/paneles", tags=["Admin"], dependencies=[Depends(auth_admin)])
def create_panel(panel: dict):
    """Crear nuevo panel con ID auto-generado"""
    from models import get_db_session, Panel
    
    session = get_db_session()
    try:
        # Validar campos requeridos
        required = ["nombre", "capacidad", "precio", "descripcion"]
        if not all(k in panel for k in required):
            raise HTTPException(400, f"Campos requeridos: {', '.join(required)}")
        
        # Auto-generar ID: encontrar el próximo disponible
        existing_ids = [p.id for p in session.query(Panel.id).all()]
        counter = 1
        while f"panel{counter}" in existing_ids:
            counter += 1
        panel_id = f"panel{counter}"
        
        # Crear panel
        new_panel = Panel(
            id=panel_id,
            nombre=panel["nombre"],
            capacidad=panel["capacidad"],
            precio=panel["precio"],
            descripcion=panel["descripcion"],
            eficienciaPanel=panel.get("eficienciaPanel", 0.90),
            area=panel.get("area", 2.0),
            default=panel.get("default", False)
        )
        
        session.add(new_panel)
        session.commit()
        
        print(f"✅ Panel {panel_id} creado en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Panel {panel_id} creado exitosamente",
            "id": panel_id
        }
    except Exception as e:
        session.rollback()
        print(f"❌ Error al crear panel: {e}")
        raise HTTPException(500, f"Error al crear panel: {e}")
    finally:
        session.close()

@app.put("/api/admin/paneles/{panel_id}", tags=["Admin"], dependencies=[Depends(auth_admin)])
def update_panel(panel_id: str, panel: dict):
    """Actualizar panel existente"""
    from models import get_db_session, Panel
    
    session = get_db_session()
    try:
        # Buscar panel
        panel_db = session.query(Panel).filter_by(id=panel_id).first()
        if not panel_db:
            raise HTTPException(404, f"Panel {panel_id} no encontrado")
        
        # Actualizar campos
        panel_db.nombre = panel.get("nombre", panel_db.nombre)
        panel_db.capacidad = panel.get("capacidad", panel_db.capacidad)
        panel_db.precio = panel.get("precio", panel_db.precio)
        panel_db.descripcion = panel.get("descripcion", panel_db.descripcion)
        panel_db.eficienciaPanel = panel.get("eficienciaPanel", panel_db.eficienciaPanel)
        panel_db.area = panel.get("area", panel_db.area)
        panel_db.default = panel.get("default", panel_db.default)
        
        session.commit()
        
        print(f"✅ Panel {panel_id} actualizado en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Panel {panel_id} actualizado exitosamente"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al actualizar panel: {e}")
        raise HTTPException(500, f"Error al actualizar: {e}")
    finally:
        session.close()

@app.delete("/api/admin/paneles/{panel_id}", tags=["Admin"], dependencies=[Depends(auth_admin)])
def delete_panel(panel_id: str):
    """Eliminar panel"""
    from models import get_db_session, Panel
    
    session = get_db_session()
    try:
        # Buscar panel
        panel_db = session.query(Panel).filter_by(id=panel_id).first()
        if not panel_db:
            raise HTTPException(404, f"Panel {panel_id} no encontrado")
        
        session.delete(panel_db)
        session.commit()
        
        print(f"✅ Panel {panel_id} eliminado de PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Panel {panel_id} eliminado exitosamente"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al eliminar panel: {e}")
        raise HTTPException(500, f"Error al eliminar: {e}")
    finally:
        session.close()

# --- GESTIÓN DE INVERSORES ---
@app.get("/api/admin/inversores", tags=["Admin"], dependencies=[Depends(auth_admin)])
def get_inversores_admin():
    """Obtener todos los inversores con precios (admin)"""
    from models import get_db_session, Inversor
    
    session = get_db_session()
    try:
        inversores_db = session.query(Inversor).all()
        inversores = [
            {
                "id": i.id,
                "nombre": i.nombre,
                "capacidad": i.capacidad,
                "precio": i.precio,
                "descripcion": i.descripcion,
                "eficiencia": i.eficiencia,
                "sistemaElectrico": i.sistemaElectrico,
                "tipo": i.tipo,
                "paneles_por_inversor": i.paneles_por_inversor,
                "sobredimensionamiento": i.sobredimensionamiento,
                "default": i.default
            }
            for i in inversores_db
        ]
        return inversores
    finally:
        session.close()

@app.post("/api/admin/inversores", tags=["Admin"], dependencies=[Depends(auth_admin)])
def create_inversor(inversor: dict):
    """Crear nuevo inversor con ID auto-generado"""
    from models import get_db_session, Inversor
    
    session = get_db_session()
    try:
        # Validar campos requeridos
        required = ["nombre", "capacidad", "precio", "descripcion"]
        if not all(k in inversor for k in required):
            raise HTTPException(400, f"Campos requeridos: {', '.join(required)}")
        
        # Auto-generar ID
        existing_ids = [i.id for i in session.query(Inversor.id).all()]
        counter = 1
        while f"inv{counter}" in existing_ids:
            counter += 1
        inversor_id = f"inv{counter}"
        
        # Crear inversor
        new_inversor = Inversor(
            id=inversor_id,
            nombre=inversor["nombre"],
            capacidad=inversor["capacidad"],
            precio=inversor["precio"],
            descripcion=inversor["descripcion"],
            eficiencia=inversor.get("eficiencia", 0.97),
            sistemaElectrico=inversor.get("sistemaElectrico", "monofasico"),
            tipo=inversor.get("tipo", "STRING"),
            paneles_por_inversor=inversor.get("paneles_por_inversor"),
            sobredimensionamiento=inversor.get("sobredimensionamiento"),
            default=inversor.get("default", False)
        )
        
        session.add(new_inversor)
        session.commit()
        
        print(f"✅ Inversor {inversor_id} creado en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Inversor {inversor_id} creado exitosamente",
            "id": inversor_id
        }
    except Exception as e:
        session.rollback()
        print(f"❌ Error al crear inversor: {e}")
        raise HTTPException(500, f"Error al crear inversor: {e}")
    finally:
        session.close()

@app.put("/api/admin/inversores/{inversor_id}", tags=["Admin"], dependencies=[Depends(auth_admin)])
def update_inversor(inversor_id: str, inversor: dict):
    """Actualizar inversor existente"""
    from models import get_db_session, Inversor
    
    session = get_db_session()
    try:
        # Buscar inversor
        inversor_db = session.query(Inversor).filter_by(id=inversor_id).first()
        if not inversor_db:
            raise HTTPException(404, f"Inversor {inversor_id} no encontrado")
        
        # Actualizar campos
        inversor_db.nombre = inversor.get("nombre", inversor_db.nombre)
        inversor_db.capacidad = inversor.get("capacidad", inversor_db.capacidad)
        inversor_db.precio = inversor.get("precio", inversor_db.precio)
        inversor_db.descripcion = inversor.get("descripcion", inversor_db.descripcion)
        inversor_db.eficiencia = inversor.get("eficiencia", inversor_db.eficiencia)
        inversor_db.sistemaElectrico = inversor.get("sistemaElectrico", inversor_db.sistemaElectrico)
        inversor_db.tipo = inversor.get("tipo", inversor_db.tipo)
        inversor_db.paneles_por_inversor = inversor.get("paneles_por_inversor", inversor_db.paneles_por_inversor)
        inversor_db.sobredimensionamiento = inversor.get("sobredimensionamiento", inversor_db.sobredimensionamiento)
        inversor_db.default = inversor.get("default", inversor_db.default)
        
        session.commit()
        
        print(f"✅ Inversor {inversor_id} actualizado en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Inversor {inversor_id} actualizado exitosamente"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al actualizar inversor: {e}")
        raise HTTPException(500, f"Error al actualizar: {e}")
    finally:
        session.close()

@app.delete("/api/admin/inversores/{inversor_id}", tags=["Admin"], dependencies=[Depends(auth_admin)])
def delete_inversor(inversor_id: str):
    """Eliminar inversor"""
    from models import get_db_session, Inversor
    
    session = get_db_session()
    try:
        # Buscar inversor
        inversor_db = session.query(Inversor).filter_by(id=inversor_id).first()
        if not inversor_db:
            raise HTTPException(404, f"Inversor {inversor_id} no encontrado")
        
        session.delete(inversor_db)
        session.commit()
        
        print(f"✅ Inversor {inversor_id} eliminado de PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Inversor {inversor_id} eliminado exitosamente"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al eliminar inversor: {e}")
        raise HTTPException(500, f"Error al eliminar: {e}")
    finally:
        session.close()

# --- GESTIÓN DE BATERÍAS ---
@app.get("/api/admin/baterias", tags=["Admin"], dependencies=[Depends(auth_admin)])
def get_baterias_admin():
    """Obtener todas las baterías con precios (admin)"""
    from models import get_db_session, Bateria
    
    session = get_db_session()
    try:
        baterias_db = session.query(Bateria).all()
        baterias = [
            {
                "id": b.id,
                "nombre": b.nombre,
                "capacidad": b.capacidad,
                "precio": b.precio,
                "descripcion": b.descripcion,
                "default": b.default
            }
            for b in baterias_db
        ]
        return baterias
    finally:
        session.close()

@app.post("/api/admin/baterias", tags=["Admin"], dependencies=[Depends(auth_admin)])
def create_bateria(bateria: dict):
    """Crear nueva batería con ID auto-generado"""
    from models import get_db_session, Bateria
    
    session = get_db_session()
    try:
        # Validar campos requeridos
        required = ["nombre", "capacidad", "precio", "descripcion"]
        if not all(k in bateria for k in required):
            raise HTTPException(400, f"Campos requeridos: {', '.join(required)}")
        
        # Auto-generar ID
        existing_ids = [b.id for b in session.query(Bateria.id).all()]
        counter = 1
        while f"bat{counter}" in existing_ids:
            counter += 1
        bateria_id = f"bat{counter}"
        
        # Crear batería
        new_bateria = Bateria(
            id=bateria_id,
            nombre=bateria["nombre"],
            capacidad=bateria["capacidad"],
            precio=bateria["precio"],
            descripcion=bateria["descripcion"],
            default=bateria.get("default", False)
        )
        
        session.add(new_bateria)
        session.commit()
        
        print(f"✅ Batería {bateria_id} creada en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Batería {bateria_id} creada exitosamente",
            "id": bateria_id
        }
    except Exception as e:
        session.rollback()
        print(f"❌ Error al crear batería: {e}")
        raise HTTPException(500, f"Error al crear batería: {e}")
    finally:
        session.close()

@app.put("/api/admin/baterias/{bateria_id}", tags=["Admin"], dependencies=[Depends(auth_admin)])
def update_bateria(bateria_id: str, bateria: dict):
    """Actualizar batería existente"""
    from models import get_db_session, Bateria
    
    session = get_db_session()
    try:
        # Buscar batería
        bateria_db = session.query(Bateria).filter_by(id=bateria_id).first()
        if not bateria_db:
            raise HTTPException(404, f"Batería {bateria_id} no encontrada")
        
        # Actualizar campos
        bateria_db.nombre = bateria.get("nombre", bateria_db.nombre)
        bateria_db.capacidad = bateria.get("capacidad", bateria_db.capacidad)
        bateria_db.precio = bateria.get("precio", bateria_db.precio)
        bateria_db.descripcion = bateria.get("descripcion", bateria_db.descripcion)
        bateria_db.default = bateria.get("default", bateria_db.default)
        
        session.commit()
        
        print(f"✅ Batería {bateria_id} actualizada en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Batería {bateria_id} actualizada exitosamente"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al actualizar batería: {e}")
        raise HTTPException(500, f"Error al actualizar: {e}")
    finally:
        session.close()

@app.delete("/api/admin/baterias/{bateria_id}", tags=["Admin"], dependencies=[Depends(auth_admin)])
def delete_bateria(bateria_id: str):
    """Eliminar batería"""
    from models import get_db_session, Bateria
    
    session = get_db_session()
    try:
        # Buscar batería
        bateria_db = session.query(Bateria).filter_by(id=bateria_id).first()
        if not bateria_db:
            raise HTTPException(404, f"Batería {bateria_id} no encontrada")
        
        session.delete(bateria_db)
        session.commit()
        
        print(f"✅ Batería {bateria_id} eliminada de PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Batería {bateria_id} eliminada exitosamente"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al eliminar batería: {e}")
        raise HTTPException(500, f"Error al eliminar: {e}")
    finally:
        session.close()

# --- GESTIÓN DE EQUIPOS DEFAULT ---
@app.put("/api/admin/paneles/{panel_id}/default", tags=["Admin"], dependencies=[Depends(auth_admin)])
def set_panel_default(panel_id: str):
    """Marcar un panel como default (desmarca los demás)"""
    from models import get_db_session, Panel
    
    session = get_db_session()
    try:
        # Verificar que existe el panel
        panel = session.query(Panel).filter_by(id=panel_id).first()
        if not panel:
            raise HTTPException(404, f"Panel {panel_id} no encontrado")
        
        # Desmarcar todos los paneles como default
        session.query(Panel).update({"default": False})
        
        # Marcar el seleccionado como default
        panel.default = True
        
        session.commit()
        
        print(f"✅ Panel {panel_id} marcado como default en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Panel {panel_id} marcado como default"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al marcar panel default: {e}")
        raise HTTPException(500, f"Error al marcar default: {e}")
    finally:
        session.close()

@app.put("/api/admin/inversores/{inversor_id}/default", tags=["Admin"], dependencies=[Depends(auth_admin)])
def set_inversor_default(inversor_id: str):
    """Marcar un inversor como default (desmarca los demás)"""
    from models import get_db_session, Inversor
    
    session = get_db_session()
    try:
        # Verificar que existe el inversor
        inversor = session.query(Inversor).filter_by(id=inversor_id).first()
        if not inversor:
            raise HTTPException(404, f"Inversor {inversor_id} no encontrado")
        
        # Desmarcar todos los inversores como default
        session.query(Inversor).update({"default": False})
        
        # Marcar el seleccionado como default
        inversor.default = True
        
        session.commit()
        
        print(f"✅ Inversor {inversor_id} marcado como default en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Inversor {inversor_id} marcado como default"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al marcar inversor default: {e}")
        raise HTTPException(500, f"Error al marcar default: {e}")
    finally:
        session.close()

@app.put("/api/admin/baterias/{bateria_id}/default", tags=["Admin"], dependencies=[Depends(auth_admin)])
def set_bateria_default(bateria_id: str):
    """Marcar una batería como default (desmarca las demás)"""
    from models import get_db_session, Bateria
    
    session = get_db_session()
    try:
        # Verificar que existe la batería
        bateria = session.query(Bateria).filter_by(id=bateria_id).first()
        if not bateria:
            raise HTTPException(404, f"Batería {bateria_id} no encontrada")
        
        # Desmarcar todas las baterías como default
        session.query(Bateria).update({"default": False})
        
        # Marcar la seleccionada como default
        bateria.default = True
        
        session.commit()
        
        print(f"✅ Batería {bateria_id} marcada como default en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Batería {bateria_id} marcada como default"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al marcar batería default: {e}")
        raise HTTPException(500, f"Error al marcar default: {e}")
    finally:
        session.close()

# --- GESTIÓN DE CIUDADES ---
@app.get("/api/admin/ciudades", tags=["Admin"], dependencies=[Depends(auth_admin)])
def get_ciudades_admin():
    """Obtener todas las ciudades con HSP (admin)"""
    from models import get_db_session, Ciudad
    
    session = get_db_session()
    try:
        ciudades_db = session.query(Ciudad).order_by(Ciudad.nombre).all()
        ciudades_list = []
        for c in ciudades_db:
            # Obtener factorTemperatura con validación
            factorTemp = c.factorTemperatura if hasattr(c, 'factorTemperatura') else 0.90
            
            # VALIDACIÓN: Normalizar si está en formato porcentual
            if factorTemp > 1.0:
                print(f"⚠️ Ciudad {c.nombre}: factorTemperatura = {factorTemp} (corrigiendo a {factorTemp/100})")
                factorTemp = factorTemp / 100
            elif factorTemp < 0.5:
                print(f"⚠️ Ciudad {c.nombre}: factorTemperatura = {factorTemp} (usando default 0.90)")
                factorTemp = 0.90
            
            ciudades_list.append({
                "key": c.key,
                "nombre": c.nombre,
                "hsp": c.hsp,
                "factorTemperatura": factorTemp
            })
        return ciudades_list
    finally:
        session.close()

@app.post("/api/admin/ciudades", tags=["Admin"], dependencies=[Depends(auth_admin)])
def create_ciudad(ciudad: dict):
    """Crear nueva ciudad con HSP"""
    from models import get_db_session, Ciudad
    
    session = get_db_session()
    try:
        # Validar campos requeridos
        if "nombre" not in ciudad or "hsp" not in ciudad:
            raise HTTPException(400, "Campos requeridos: nombre, hsp")
        
        # Normalizar nombre para key (lowercase, underscores, sin acentos)
        ciudad_key = ciudad["nombre"].lower().replace(" ", "_").replace("á", "a").replace("é", "e").replace("í", "i").replace("ó", "o").replace("ú", "u").replace("ñ", "n")
        
        # Verificar que no existe
        existing = session.query(Ciudad).filter_by(key=ciudad_key).first()
        if existing:
            raise HTTPException(400, f"La ciudad {ciudad['nombre']} ya existe")
        
        # Crear ciudad
        new_ciudad = Ciudad(
            key=ciudad_key,
            nombre=ciudad["nombre"],
            hsp=float(ciudad["hsp"]),
            factorTemperatura=float(ciudad.get("factorTemperatura", 0.90))
        )
        
        session.add(new_ciudad)
        session.commit()
        
        print(f"✅ Ciudad {ciudad['nombre']} creada en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Ciudad {ciudad['nombre']} creada exitosamente",
            "key": ciudad_key
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al crear ciudad: {e}")
        raise HTTPException(500, f"Error al crear ciudad: {e}")
    finally:
        session.close()

@app.put("/api/admin/ciudades/{ciudad_key}", tags=["Admin"], dependencies=[Depends(auth_admin)])
def update_ciudad(ciudad_key: str, ciudad: dict):
    """Actualizar ciudad existente"""
    from models import get_db_session, Ciudad
    
    session = get_db_session()
    try:
        # Buscar ciudad
        ciudad_db = session.query(Ciudad).filter_by(key=ciudad_key).first()
        if not ciudad_db:
            raise HTTPException(404, f"Ciudad {ciudad_key} no encontrada")
        
        # Actualizar campos
        if "nombre" in ciudad:
            ciudad_db.nombre = ciudad["nombre"]
        if "hsp" in ciudad:
            ciudad_db.hsp = float(ciudad["hsp"])
        if "factorTemperatura" in ciudad:
            ciudad_db.factorTemperatura = float(ciudad["factorTemperatura"])
        
        session.commit()
        
        print(f"✅ Ciudad {ciudad_key} actualizada en PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Ciudad {ciudad_key} actualizada exitosamente"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al actualizar ciudad: {e}")
        raise HTTPException(500, f"Error al actualizar: {e}")
    finally:
        session.close()

@app.delete("/api/admin/ciudades/{ciudad_key}", tags=["Admin"], dependencies=[Depends(auth_admin)])
def delete_ciudad(ciudad_key: str):
    """Eliminar ciudad"""
    from models import get_db_session, Ciudad
    
    session = get_db_session()
    try:
        # Buscar ciudad
        ciudad_db = session.query(Ciudad).filter_by(key=ciudad_key).first()
        if not ciudad_db:
            raise HTTPException(404, f"Ciudad {ciudad_key} no encontrada")
        
        session.delete(ciudad_db)
        session.commit()
        
        print(f"✅ Ciudad {ciudad_key} eliminada de PostgreSQL")
        return {
            "status": "success",
            "mensaje": f"Ciudad {ciudad_key} eliminada exitosamente"
        }
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        print(f"❌ Error al eliminar ciudad: {e}")
        raise HTTPException(500, f"Error al eliminar: {e}")
    finally:
        session.close()

@app.post("/api/admin/ciudades/fix-temperatura", tags=["Admin"], dependencies=[Depends(auth_admin)])
def fix_factor_temperatura():
    """
    🔧 Endpoint de mantenimiento: Corregir factorTemperatura en formato porcentual
    Convierte valores >1.0 dividiéndolos entre 100
    """
    from models import get_db_session, Ciudad
    
    session = get_db_session()
    try:
        ciudades_db = session.query(Ciudad).all()
        corregidas = []
        sin_cambios = []
        
        for c in ciudades_db:
            if hasattr(c, 'factorTemperatura') and c.factorTemperatura is not None:
                valor_original = c.factorTemperatura
                
                if valor_original > 1.0:
                    # Corregir: dividir entre 100
                    c.factorTemperatura = valor_original / 100
                    corregidas.append({
                        "ciudad": c.nombre,
                        "anterior": valor_original,
                        "nuevo": c.factorTemperatura
                    })
                    print(f"✅ {c.nombre}: {valor_original} → {c.factorTemperatura}")
                elif valor_original < 0.5:
                    # Valor inválido: usar default
                    c.factorTemperatura = 0.90
                    corregidas.append({
                        "ciudad": c.nombre,
                        "anterior": valor_original,
                        "nuevo": 0.90,
                        "nota": "Valor demasiado bajo, usando default"
                    })
                    print(f"✅ {c.nombre}: {valor_original} → 0.90 (default)")
                else:
                    sin_cambios.append(c.nombre)
            else:
                # No tiene factorTemperatura: agregar default
                c.factorTemperatura = 0.90
                corregidas.append({
                    "ciudad": c.nombre,
                    "anterior": "null",
                    "nuevo": 0.90,
                    "nota": "Campo faltante, agregado default"
                })
        
        session.commit()
        
        return {
            "status": "success",
            "mensaje": f"✅ Proceso completado. {len(corregidas)} ciudades corregidas, {len(sin_cambios)} sin cambios",
            "corregidas": corregidas,
            "sin_cambios": sin_cambios
        }
    except Exception as e:
        session.rollback()
        print(f"❌ Error al corregir factorTemperatura: {e}")
        raise HTTPException(500, f"Error: {e}")
    finally:
        session.close()

# --- SISTEMA DE TRACKING Y VALORES POR DEFECTO ---
@app.post("/api/track-seleccion", tags=["Analytics"])
async def track_seleccion(request: Request):
    """Registrar selecciones de una cotización para análisis de frecuencia (sin autenticación)"""
    try:
        data = await request.json()
        
        # Cargar estadísticas existentes
        if os.path.exists(ESTADISTICAS_FILE):
            estadisticas = load_json(ESTADISTICAS_FILE)
        else:
            estadisticas = {"cotizaciones": []}
        
        # Agregar nuevo registro con timestamp
        registro = {
            "timestamp": now_colombia().isoformat(),
            "ciudad": data.get("ciudad", ""),
            "tipoSistemaFV": data.get("tipoSistemaFV", ""),
            "tipo_propiedad": data.get("tipo_propiedad", ""),
            "legalizacion": data.get("legalizacion", ""),
            "seleccionManual": data.get("seleccionManual", ""),
            "sistemaElectrico": data.get("sistemaElectrico", ""),
            "porcentajeConsumodia": data.get("porcentajeConsumodia", 50)
        }
        
        estadisticas["cotizaciones"].append(registro)
        
        # Guardar (limitamos a las últimas 500 cotizaciones para no saturar)
        if len(estadisticas["cotizaciones"]) > 500:
            estadisticas["cotizaciones"] = estadisticas["cotizaciones"][-500:]
        
        with open(ESTADISTICAS_FILE, "w", encoding="utf-8") as f:
            json.dump(estadisticas, f, ensure_ascii=False, indent=2)
        
        return {"status": "success", "mensaje": "Selección registrada"}
    except Exception as e:
        # No fallar si hay error en tracking (no es crítico)
        print(f"Error en tracking: {e}")
        return {"status": "error", "mensaje": str(e)}

@app.get("/api/valores-default", tags=["Analytics"])
def get_valores_default():
    """Obtener valores más frecuentes de los últimos 30 días (público, sin autenticación)"""
    try:
        if not os.path.exists(ESTADISTICAS_FILE):
            # Valores por defecto fijos si no hay estadísticas
            return {
                "ciudad": "santa_marta",
                "tipoSistemaFV": "ongrid",
                "tipo_propiedad": "residencial",
                "legalizacion": "SI",
                "seleccionManual": "NO",
                "sistemaElectrico": "bifasico",
                "porcentajeConsumodia": 50,
                "source": "defaults"
            }
        
        estadisticas = load_json(ESTADISTICAS_FILE)
        cotizaciones = estadisticas.get("cotizaciones", [])
        
        if not cotizaciones:
            # Valores por defecto fijos
            return {
                "ciudad": "santa_marta",
                "tipoSistemaFV": "ongrid",
                "tipo_propiedad": "residencial",
                "legalizacion": "SI",
                "seleccionManual": "NO",
                "sistemaElectrico": "bifasico",
                "porcentajeConsumodia": 50,
                "source": "defaults"
            }
        
        # Filtrar últimos 30 días
        fecha_limite = now_colombia() - timedelta(days=30)
        cotizaciones_recientes = []
        
        for cot in cotizaciones:
            try:
                timestamp = datetime.fromisoformat(cot["timestamp"])
                if timestamp >= fecha_limite:
                    cotizaciones_recientes.append(cot)
            except:
                continue
        
        # Si no hay datos recientes, usar todos los datos disponibles
        if not cotizaciones_recientes:
            cotizaciones_recientes = cotizaciones[-50:]  # Últimas 50
        
        # Calcular frecuencias
        from collections import Counter
        
        ciudades = Counter([c.get("ciudad", "") for c in cotizaciones_recientes if c.get("ciudad")])
        tipos_sistema = Counter([c.get("tipoSistemaFV", "") for c in cotizaciones_recientes if c.get("tipoSistemaFV")])
        tipos_propiedad = Counter([c.get("tipo_propiedad", "") for c in cotizaciones_recientes if c.get("tipo_propiedad")])
        legalizaciones = Counter([c.get("legalizacion", "") for c in cotizaciones_recientes if c.get("legalizacion")])
        selecciones_manual = Counter([c.get("seleccionManual", "") for c in cotizaciones_recientes if c.get("seleccionManual")])
        sistemas_electricos = Counter([c.get("sistemaElectrico", "") for c in cotizaciones_recientes if c.get("sistemaElectrico")])
        porcentajes_consumo = [c.get("porcentajeConsumodia", 50) for c in cotizaciones_recientes if c.get("porcentajeConsumodia") is not None]
        
        # Calcular promedio de porcentaje de consumo día
        porcentaje_promedio = round(sum(porcentajes_consumo) / len(porcentajes_consumo)) if porcentajes_consumo else 50
        
        # Obtener valores más comunes
        resultado = {
            "ciudad": ciudades.most_common(1)[0][0] if ciudades else "santa_marta",
            "tipoSistemaFV": tipos_sistema.most_common(1)[0][0] if tipos_sistema else "ongrid",
            "tipo_propiedad": tipos_propiedad.most_common(1)[0][0] if tipos_propiedad else "residencial",
            "legalizacion": legalizaciones.most_common(1)[0][0] if legalizaciones else "SI",
            "seleccionManual": selecciones_manual.most_common(1)[0][0] if selecciones_manual else "NO",
            "sistemaElectrico": sistemas_electricos.most_common(1)[0][0] if sistemas_electricos else "bifasico",
            "porcentajeConsumodia": porcentaje_promedio,
            "source": "analytics",
            "sample_size": len(cotizaciones_recientes)
        }
        
        return resultado
        
    except Exception as e:
        print(f"Error en valores-default: {e}")
        # En caso de error, retornar defaults fijos
        return {
            "ciudad": "santa_marta",
            "tipoSistemaFV": "ongrid",
            "tipo_propiedad": "residencial",
            "legalizacion": "SI",
            "seleccionManual": "NO",
            "sistemaElectrico": "bifasico",
            "porcentajeConsumodia": 50,
            "source": "defaults"
        }

@app.get("/api/diagnostico-postgres", tags=["Debug"])
def diagnostico_postgres():
    """Endpoint de diagnóstico para verificar carga de datos desde PostgreSQL"""
    resultado = {
        "timestamp": now_colombia().isoformat(),
        "postgres_available": POSTGRES_AVAILABLE,
        "error": None,
        "datos_cargados": {}
    }
    
    try:
        equipos, ciudades, parametros = cargar_datos_desde_postgres()
        
        resultado["datos_cargados"] = {
            "paneles": len(equipos.get("paneles", [])),
            "inversores": len(equipos.get("inversores", [])),
            "baterias": len(equipos.get("baterias", [])),
            "ciudades": len(ciudades),
            "parametros_secciones": list(parametros.keys())
        }
        resultado["status"] = "success"
        
    except Exception as e:
        import traceback
        error_details = {
            "type": type(e).__name__,
            "message": str(e),
            "traceback": traceback.format_exc()
        }
        resultado["error"] = error_details
        resultado["status"] = "error"
    
    return JSONResponse(resultado)


@app.post("/api/admin/migraciones/factor-temperatura", tags=["Admin"])
def migrar_factor_temperatura_endpoint(credentials: HTTPBasicCredentials = Depends(security)):
    """
    Endpoint admin para agregar columna factorTemperatura a la tabla ciudades.
    Requiere autenticación HTTP Basic.
    
    Valores de referencia:
    - Costa Caribe: 0.85 (alta temperatura)
    - Interior/Valles: 0.90 (temperatura moderada)  
    - Alta Montaña: 0.92 (baja temperatura)
    """
    # Verificar credenciales
    auth_admin(credentials)
    
    if not POSTGRES_AVAILABLE:
        raise HTTPException(status_code=500, detail="PostgreSQL no disponible")
    
    # Valores de factorTemperatura por ciudad
    FACTORES_TEMPERATURA = {
        # Costa Caribe - Alta temperatura = menor factor (0.85)
        "santa_marta": 0.85, "barranquilla": 0.85, "cartagena": 0.85,
        "valledupar": 0.86, "riohacha": 0.85, "sincelejo": 0.86,
        "monteria": 0.86, "magangue": 0.85, "cienaga": 0.85,
        "fundacion": 0.85, "aracataca": 0.85, "zona_bananera": 0.85,
        "pueblo_viejo": 0.85, "algarrobo": 0.85, "albania_guajira": 0.84,
        "maicao": 0.84, "uribia": 0.84,
        
        # Interior/Valles - Temperatura moderada (0.88-0.90)
        "medellin": 0.89, "cali": 0.88, "bucaramanga": 0.89,
        "cucuta": 0.88, "pereira": 0.89, "manizales": 0.91,
        "armenia": 0.89, "ibague": 0.88, "neiva": 0.87,
        "villavicencio": 0.88, "yopal": 0.87, "florencia": 0.87,
        
        # Alta Montaña - Baja temperatura = mayor factor (0.92-0.93)
        "bogota": 0.92, "tunja": 0.92, "pasto": 0.93,
        "popayan": 0.91, "duitama": 0.92, "sogamoso": 0.92,
        "zipaquira": 0.92, "chia": 0.92, "facatativa": 0.92,
    }
    
    session = get_db_session()
    try:
        from sqlalchemy import text
        
        # 1. Verificar si la columna ya existe
        result = session.execute(text("""
            SELECT column_name 
            FROM information_schema.columns 
            WHERE table_name='ciudades' AND column_name='factorTemperatura'
        """))
        
        columna_existe = result.fetchone() is not None
        
        if not columna_existe:
            # 2. Agregar la columna con valor por defecto 0.90
            session.execute(text("""
                ALTER TABLE ciudades 
                ADD COLUMN "factorTemperatura" DOUBLE PRECISION DEFAULT 0.90
            """))
            session.commit()
        
        # 3. Actualizar valores específicos por ciudad
        ciudades_actualizadas = 0
        for ciudad_key, factor in FACTORES_TEMPERATURA.items():
            result = session.execute(
                text("""
                    UPDATE ciudades 
                    SET "factorTemperatura" = :factor 
                    WHERE key = :ciudad_key
                """),
                {"factor": factor, "ciudad_key": ciudad_key}
            )
            if result.rowcount > 0:
                ciudades_actualizadas += 1
        
        session.commit()
        
        # 4. Obtener estadísticas
        result = session.execute(text("""
            SELECT COUNT(*) as total,
                   AVG("factorTemperatura") as promedio,
                   MIN("factorTemperatura") as minimo,
                   MAX("factorTemperatura") as maximo
            FROM ciudades
        """))
        stats = result.fetchone()
        
        return {
            "status": "success",
            "columna_ya_existia": columna_existe,
            "total_ciudades": stats[0],
            "ciudades_actualizadas": ciudades_actualizadas,
            "estadisticas": {
                "promedio": round(stats[1], 3),
                "minimo": round(stats[2], 3),
                "maximo": round(stats[3], 3)
            },
            "mensaje": "Migración completada exitosamente. Las ciudades no listadas usan factor default 0.90"
        }
        
    except Exception as e:
        session.rollback()
        import traceback
        raise HTTPException(
            status_code=500,
            detail={
                "error": str(e),
                "type": type(e).__name__,
                "traceback": traceback.format_exc()
            }
        )
    finally:
        session.close()

@app.post("/api/admin/equipos/fix-eficiencias", tags=["Admin"], dependencies=[Depends(auth_admin)])
def fix_eficiencias_equipos():
    """
    🔧 Endpoint de mantenimiento: Corregir eficiencias en formato porcentual
    Convierte valores >1.0 dividiéndolos entre 100 (ej: 100.0 → 1.0, 98.0 → 0.98)
    """
    from models import get_db_session, Panel, Inversor
    
    session = get_db_session()
    try:
        corregidas_paneles = []
        corregidas_inversores = []
        
        # PANELES: Corregir eficienciaPanel
        paneles_db = session.query(Panel).all()
        for p in paneles_db:
            if hasattr(p, 'eficienciaPanel') and p.eficienciaPanel is not None:
                valor_original = p.eficienciaPanel
                
                if valor_original > 1.0:
                    # Corregir: dividir entre 100
                    p.eficienciaPanel = valor_original / 100
                    corregidas_paneles.append({
                        "id": p.id,
                        "nombre": p.nombre,
                        "anterior": valor_original,
                        "nuevo": p.eficienciaPanel
                    })
                    print(f"✅ Panel {p.id}: eficienciaPanel {valor_original} → {p.eficienciaPanel}")
        
        # INVERSORES: Corregir eficiencia
        inversores_db = session.query(Inversor).all()
        for i in inversores_db:
            if hasattr(i, 'eficiencia') and i.eficiencia is not None:
                valor_original = i.eficiencia
                
                if valor_original > 1.0:
                    # Corregir: dividir entre 100
                    i.eficiencia = valor_original / 100
                    corregidas_inversores.append({
                        "id": i.id,
                        "nombre": i.nombre,
                        "anterior": valor_original,
                        "nuevo": i.eficiencia
                    })
                    print(f"✅ Inversor {i.id}: eficiencia {valor_original} → {i.eficiencia}")
        
        session.commit()
        
        total_corregidas = len(corregidas_paneles) + len(corregidas_inversores)
        
        return {
            "status": "success",
            "mensaje": f"✅ Proceso completado. {total_corregidas} equipos corregidos",
            "paneles_corregidos": corregidas_paneles,
            "inversores_corregidos": corregidas_inversores
        }
    except Exception as e:
        session.rollback()
        print(f"❌ Error al corregir eficiencias: {e}")
        raise HTTPException(500, f"Error: {e}")
    finally:
        session.close()


@app.post("/api/cotizar", tags=["Cotización"])
async def cotizar(request: Request, req: CotizarRequest, _: Any = Depends(rate_limit)):
    """Generar cotización completa con 1 o 2 opciones según área disponible"""
    # MIGRACIÓN POSTGRESQL: Cargar desde BD en lugar de archivos JSON
    equipos, ciudades, parametros = cargar_datos_desde_postgres()
    
    # Preparar datos de solicitud
    req_dict = req.dict()
    
    # Si seleccionManual es NO, usar equipos por defecto
    if req.seleccionManual == "NO":
        # FIX #4: Pasar sistema eléctrico para selección inteligente de inversor
        defaults = obtener_equipos_defaults(equipos, req.sistemaElectrico)
        req_dict["panel"] = defaults["panel"]
        req_dict["inversor"] = defaults["inversor"]
        
        # Log para debugging
        print(f"🔧 Selección automática de equipos:")
        print(f"   Sistema eléctrico: {req.sistemaElectrico}")
        print(f"   Panel default: {defaults['panel']}")
        print(f"   Inversor default: {defaults['inversor']}")
        
        # Solo asignar batería default si el sistema la requiere
        if req.tipoSistemaFV in ("offgrid", "hibrido_incluido") and defaults["bateria"]:
            req_dict["bateria"] = defaults["bateria"]
            print(f"   Batería default: {defaults['bateria']}")
    
    
    # Determinar COND_COM según legalización
    if req.legalizacion == "SI":
        cond_com = "Anticipo 70%, 25% en producción, 5% legalizado"
    else:
        cond_com = "Anticipo 70%, 30% contraentrega"
    
    try:
        # Calcular OPCIÓN 1 (ideal sin restricciones)
        resultado_opcion1 = calcular_cotizacion(req_dict, equipos, ciudades, parametros)
        # Agregar COND_COM al resultado
        resultado_opcion1["condicionesComerciales"] = cond_com
    except ValueError as e:
        raise HTTPException(400, str(e))
    except Exception as e:
        raise HTTPException(500, f"Error cálculo: {e}")

    # Verificar si se necesita segunda opción (área disponible < umbral del área requerida)
    areaDisponible = float(req.areaDisponible or 0)
    areaRequerida = resultado_opcion1["areaRequerida"]
    
    # Obtener umbral desde parámetros (default 92%)
    parametros_sistema = parametros.get("parametros_sistema", {})
    umbral_porcentaje = parametros_sistema.get("umbral_segunda_opcion", 0.92)
    umbral = areaRequerida * umbral_porcentaje
    
    # Logs detallados para diagnóstico
    diagnostico = {
        "areaDisponible_raw": req.areaDisponible,
        "areaDisponible_float": areaDisponible,
        "areaRequerida": areaRequerida,
        "umbral_92": round(umbral, 2),
        "condicion1_mayor_cero": areaDisponible > 0,
        "condicion2_menor_umbral": areaDisponible < umbral,
        "generara_segunda_opcion": areaDisponible > 0 and areaDisponible < umbral,
        "error_opcion2": None  # Se llenará si hay error
    }
    
    print(f"\n{'='*80}")
    print(f"🔍 DIAGNÓSTICO COMPLETO - SEGUNDA OPCIÓN")
    print(f"{'='*80}")
    print(f"📥 DATOS RECIBIDOS:")
    print(f"   req.areaDisponible (raw): {req.areaDisponible}")
    print(f"   req.areaDisponible (type): {type(req.areaDisponible)}")
    print(f"   areaDisponible (float convertido): {areaDisponible}")
    print(f"\n📊 CÁLCULOS:")
    print(f"   Área requerida (opción 1): {areaRequerida} m²")
    print(f"   Umbral 92%: {umbral:.2f} m²")
    print(f"\n✅ EVALUACIÓN DE CONDICIONES:")
    print(f"   [1] ¿Área disponible > 0?: {areaDisponible > 0}")
    print(f"   [2] ¿Área disponible < umbral?: {areaDisponible < umbral}")
    print(f"\n📋 DIAGNÓSTICO JSON: {diagnostico}")
    
    necesita_segunda_opcion = areaDisponible > 0 and areaDisponible < umbral
    print(f"\n{'🎯 DECISIÓN FINAL: ' + ('✅ SÍ GENERA 2 OPCIONES' if necesita_segunda_opcion else '❌ NO, SOLO 1 OPCIÓN')}")
    print(f"{'='*80}\n")
    
    # CALCULAR SEGUNDA OPCIÓN si es necesario (sin generar PDFs aún)
    resultado_opcion2 = None
    num_opciones = 1
    
    if necesita_segunda_opcion:
        print(f"\n🚀 CALCULANDO SEGUNDA OPCIÓN")
        print(f"   Área disponible: {areaDisponible} m²")
        print(f"   Área requerida original: {areaRequerida} m²")
        
        try:
            print(f"   📊 Calculando segunda opción...")
            # Extraer ID base sin sufijo para mantener consistencia
            cotizacion_id_base = resultado_opcion1["cotizacionId"]
            resultado_opcion2 = calcular_segunda_opcion(req_dict, equipos, ciudades, areaDisponible, cotizacion_id_base, parametros)
            # Agregar COND_COM también a la opción 2
            resultado_opcion2["condicionesComerciales"] = cond_com
            num_opciones = 2
            print(f"   ✅ Cálculo completado:")
            print(f"      - Paneles: {resultado_opcion2['numeroPaneles']} (vs {resultado_opcion1['numeroPaneles']} original)")
            print(f"      - Capacidad: {resultado_opcion2['capacidadInstalada']} kW")
            print(f"      - Valor: ${resultado_opcion2['valorTotalSistema']:,.0f}")
        except Exception as e:
            error_msg = f"{type(e).__name__}: {str(e)}"
            diagnostico["error_opcion2"] = error_msg
            print(f"   ❌ ERROR CALCULANDO SEGUNDA OPCIÓN: {error_msg}")
            import traceback
            traceback.print_exc()
            # No lanzar excepción, continuar con una sola opción
    else:
        print(f"\n❌ NO SE CALCULA SEGUNDA OPCIÓN (condición no cumplida)")

    # Combinar datos de la solicitud con el resultado para el frontend
    resumen_completo = {
        **resultado_opcion1,
        "nombre": req.nombre,
        "email": req.email,
        "telefono": req.telefono,
        "ciudad": req.ciudad,
        "direccion": req.direccion,
        "tipoVivienda": req.tipoVivienda,
        "sistemaElectrico": req.sistemaElectrico,
        "tipoSistemaFV": req.tipoSistemaFV,
        "consumoMensual": req.consumoMensual,
        "valorFactura": req.valorFactura,
        "valorKwh": req.valorKwh,
        "porcentajeConsumodia": req.porcentajeConsumodia,
        "hspCalculado": req.hspCalculado,
        "panelSeleccionado": resultado_opcion1["panel"],
        "inversorSeleccionado": resultado_opcion1["inversor"],
        "bateriaSeleccionada": resultado_opcion1["bateria"],
        "hora": now_colombia().strftime("%H:%M:%S"),
        "numOpciones": num_opciones,
        # Opción 2: Incluir TODOS los campos, no solo resumen
        "opcion2": {
            **resultado_opcion2,  # Spread completo de resultado_opcion2
            # Agregar datos del cliente para mostrar en modal
            "nombre": req.nombre,
            "email": req.email,
            "telefono": req.telefono,
            "ciudad": req.ciudad,
            "direccion": req.direccion,
            "tipoVivienda": req.tipoVivienda,
            "sistemaElectrico": req.sistemaElectrico,
            "tipoSistemaFV": req.tipoSistemaFV,
            "consumoMensual": req.consumoMensual,
            "valorFactura": req.valorFactura,
            "valorKwh": req.valorKwh,
            "porcentajeConsumodia": req.porcentajeConsumodia,
            "hspCalculado": req.hspCalculado,
            "areaDisponible": req.areaDisponible,  # ← AGREGAR PARA MENSAJE EN MODAL
            "panelSeleccionado": resultado_opcion2["panel"],
            "inversorSeleccionado": resultado_opcion2["inversor"],
            "bateriaSeleccionada": resultado_opcion2["bateria"],
            "hora": now_colombia().strftime("%H:%M:%S"),
            "cotizacionId": resultado_opcion1["cotizacionId"]  # Mismo ID para ambas opciones
        } if resultado_opcion2 else None
    }
    
    # TRACKING AUTOMÁTICO: Registrar selección para estadísticas (sin bloquear si falla)
    try:
        if os.path.exists(ESTADISTICAS_FILE):
            estadisticas = load_json(ESTADISTICAS_FILE)
        else:
            estadisticas = {"cotizaciones": []}
        
        registro = {
            "timestamp": now_colombia().isoformat(),
            "ciudad": req.ciudad,
            "tipoSistemaFV": req.tipoSistemaFV,
            "tipo_propiedad": req.tipoVivienda,
            "legalizacion": req.legalizacion,
            "seleccionManual": req.seleccionManual,
            "sistemaElectrico": req.sistemaElectrico,
            "porcentajeConsumodia": req.porcentajeConsumodia
        }
        
        estadisticas["cotizaciones"].append(registro)
        
        # Limitar a 500 registros
        if len(estadisticas["cotizaciones"]) > 500:
            estadisticas["cotizaciones"] = estadisticas["cotizaciones"][-500:]
        
        with open(ESTADISTICAS_FILE, "w", encoding="utf-8") as f:
            json.dump(estadisticas, f, ensure_ascii=False, indent=2)
        
        print(f"📊 Tracking: Selección registrada exitosamente")
    except Exception as e:
        print(f"⚠️ Error en tracking (no crítico): {e}")
    
    # GUARDAR COTIZACIÓN COMPLETA EN POSTGRESQL (para CRM y trazabilidad)
    cotizacion_id_guardado = None
    try:
        from models import get_db_session, Cotizacion
        
        session = get_db_session()
        cotizacion_id_guardado = resultado_opcion1["cotizacionId"]
        
        # Extraer datos de opción 2 si existe
        op2_data = {}
        if resultado_opcion2:
            op2_data = {
                "tiene_opcion2": True,
                "num_paneles_op2": resultado_opcion2["numeroPaneles"],
                "capacidad_instalada_op2": resultado_opcion2["capacidadInstalada"],
                "area_requerida_op2": resultado_opcion2["areaRequerida"],
                "valor_total_op2": resultado_opcion2["valorTotalSistema"],
                "ahorro_mensual_op2": resultado_opcion2["ahorroMensualEnergia"],
                "tiempo_retorno_op2": resultado_opcion2["tiempoRetorno"]
            }
        else:
            op2_data = {"tiene_opcion2": False}
        
        # Crear registro completo
        cotizacion = Cotizacion(
            id=cotizacion_id_guardado,
            fecha_creacion=now_colombia(),
            # Datos del cliente
            nombre=req.nombre,
            email=req.email,
            telefono=req.telefono,
            direccion=req.direccion,
            ciudad=req.ciudad,
            nic=req.nic if hasattr(req, 'nic') else None,
            # Datos del sistema
            tipo_vivienda=req.tipoVivienda,
            sistema_electrico=req.sistemaElectrico,
            tipo_sistema_fv=req.tipoSistemaFV,
            # Datos de consumo
            consumo_mensual=req.consumoMensual,
            valor_factura=req.valorFactura,
            valor_kwh=req.valorKwh,
            porcentaje_consumo_dia=req.porcentajeConsumodia,
            hsp_calculado=req.hspCalculado,
            area_disponible=req.areaDisponible if hasattr(req, 'areaDisponible') and req.areaDisponible else None,
            # Equipos
            panel_id=resultado_opcion1["panel"]["id"],
            panel_nombre=resultado_opcion1["panel"]["nombre"],
            inversor_id=resultado_opcion1["inversor"]["id"],
            inversor_nombre=resultado_opcion1["inversor"]["nombre"],
            bateria_id=resultado_opcion1["bateria"]["id"] if resultado_opcion1["bateria"] else None,
            bateria_nombre=resultado_opcion1["bateria"]["nombre"] if resultado_opcion1["bateria"] else None,
            # Resultados opción 1
            num_paneles_op1=resultado_opcion1["numeroPaneles"],
            capacidad_instalada_op1=resultado_opcion1["capacidadInstalada"],
            area_requerida_op1=resultado_opcion1["areaRequerida"],
            valor_total_op1=resultado_opcion1["valorTotalSistema"],
            ahorro_mensual_op1=resultado_opcion1["ahorroMensualEnergia"],
            tiempo_retorno_op1=resultado_opcion1["tiempoRetorno"],
            # Resultados opción 2 (si existe)
            **op2_data,
            # JSON completo
            datos_completos=resumen_completo,
            # Estado
            email_enviado=False,
            num_opciones=num_opciones,
            # Metadata
            legalizacion=req.legalizacion if hasattr(req, 'legalizacion') else "no",
            seleccion_manual=req.seleccionManual if hasattr(req, 'seleccionManual') else "no"
        )
        
        session.merge(cotizacion)  # merge = insert or update
        session.commit()
        session.close()
        
        print(f"💾 Cotización guardada en PostgreSQL: {cotizacion_id_guardado}")
        print(f"   📊 Cliente: {req.nombre} ({req.email})")
        print(f"   🏠 Ciudad: {req.ciudad}")
        print(f"   ⚡ Paneles: {resultado_opcion1['numeroPaneles']} x {resultado_opcion1['panel']['nombre']}")
        print(f"   💰 Valor: ${resultado_opcion1['valorTotalSistema']:,.0f}")
        if resultado_opcion2:
            print(f"   📋 Opción 2: {resultado_opcion2['numeroPaneles']} paneles, ${resultado_opcion2['valorTotalSistema']:,.0f}")

    except Exception as e:
        print(f"⚠️ Error guardando cotización en PostgreSQL (no crítico): {e}")
        import traceback
        traceback.print_exc()
    
    return JSONResponse({
        "status": "success",
        "mensaje": f"✅ Cotización calculada exitosamente. {num_opciones} opción(es) disponible(s).",
        "numOpciones": num_opciones,
        "cotizacionId": cotizacion_id_guardado,  # ID para envío posterior
        "resumen": resumen_completo,
        "diagnostico": diagnostico  # Agregar diagnóstico para debug
    })


@app.post("/api/enviar-cotizacion", tags=["Cotización"])
async def enviar_cotizacion(request: Request, data: dict, _: Any = Depends(rate_limit)):
    """
    Genera PDFs y envía cotización por email.
    Solo requiere cotizacionId - el backend ya tiene todos los datos.
    """
    try:
        # Validar campo requerido
        if "cotizacionId" not in data:
            raise HTTPException(400, "Campo requerido: cotizacionId")
        
        cotizacion_id = data["cotizacionId"]
        
        # Recuperar cotización desde PostgreSQL usando SQLAlchemy
        print(f"\n📧 RECUPERANDO COTIZACIÓN PARA ENVÍO")
        print(f"   ID: {cotizacion_id}")
        
        from models import get_db_session, Cotizacion
        
        session = get_db_session()
        cotizacion = session.query(Cotizacion).filter(Cotizacion.id == cotizacion_id).first()
        
        if not cotizacion:
            session.close()
            raise HTTPException(404, f"Cotización {cotizacion_id} no encontrada")
        
        # Extraer datos completos
        datos_completos = cotizacion.datos_completos
        num_opciones = cotizacion.num_opciones
        email_cliente = cotizacion.email
        
        print(f"   Email destino: {email_cliente}")
        print(f"   Número de opciones: {num_opciones}")
        
        session.close()
        
        # Extraer datos necesarios - INCLUIR TODOS LOS CAMPOS
        datos_cliente = {
            "nombre": datos_completos["nombre"],
            "email": datos_completos["email"],
            "telefono": datos_completos["telefono"],
            "ciudad": datos_completos["ciudad"],
            "direccion": datos_completos["direccion"],
            "nic": datos_completos.get("nic", ""),
            "tipoVivienda": datos_completos["tipoVivienda"],
            "sistemaElectrico": datos_completos["sistemaElectrico"],
            "tipoSistemaFV": datos_completos["tipoSistemaFV"],
            "consumoMensual": datos_completos["consumoMensual"],
            "valorFactura": datos_completos["valorFactura"],
            "valorKwh": datos_completos["valorKwh"],
            "porcentajeConsumodia": datos_completos["porcentajeConsumodia"],
            "hspCalculado": datos_completos["hspCalculado"],
            # AGREGAR CAMPOS FALTANTES PARA PLACEHOLDERS
            "areaDisponible": datos_completos.get("areaDisponible", 0),
            "numeroPisos": datos_completos.get("numeroPisos", "1"),
            "legalizacion": datos_completos.get("legalizacion", "NO"),
            "seleccionManual": datos_completos.get("seleccionManual", "NO")
        }
        
        resumen_opcion1 = datos_completos  # Ya tiene todos los campos necesarios
        opcion2 = datos_completos.get("opcion2")
        
        pdf_paths = []
        pptx_paths = []
        
        # GENERAR OPCIÓN 1
        print(f"\n🔄 Generando PDF Opción 1...")
        pptx_path1, pdf_path1 = fill_template_and_convert(
            datos_cliente,
            resumen_opcion1,
            opcion="" if num_opciones == 1 else "OPCIÓN 1 DE 2"
        )
        pdf_paths.append(pdf_path1)
        pptx_paths.append(pptx_path1)
        print(f"✅ Opción 1: {os.path.basename(pdf_path1)}")
        
        # GENERAR OPCIÓN 2 si existe
        if opcion2 and num_opciones == 2:
            print(f"\n🔄 Generando PDF Opción 2...")
            # Preparar resumen completo para opción 2
            resumen_opcion2 = {**datos_cliente, **opcion2}
            
            # Verificar que existe Template-PreCotizacion2.pptx
            if os.path.isfile(TEMPLATE_PPTX_OP2):
                pptx_path2, pdf_path2 = fill_template_and_convert(
                    datos_cliente,
                    resumen_opcion2,
                    opcion="OPCIÓN 2 - Ajustada a área disponible",
                    template_path=TEMPLATE_PPTX_OP2
                )
                print(f"✅ Opción 2 (template 2): {os.path.basename(pdf_path2)}")
            else:
                pptx_path2, pdf_path2 = fill_template_and_convert(
                    datos_cliente,
                    resumen_opcion2,
                    opcion="OPCIÓN 2 - Ajustada a área disponible"
                )
                print(f"✅ Opción 2 (template principal): {os.path.basename(pdf_path2)}")
            
            pdf_paths.append(pdf_path2)
            pptx_paths.append(pptx_path2)
        
        # ENVIAR EMAIL (NO CRÍTICO - Sistema continúa si falla)
        email_enviado = False
        email_error = None
        
        print(f"\n📧 ENVIANDO EMAIL")
        print(f"   Total PDFs: {len(pdf_paths)}")
        
        try:
            # Intentar SendGrid primero (configurado para Railway)
            SENDGRID_API_KEY = os.getenv("SENDGRID_API_KEY")
            if SENDGRID_API_KEY:
                print(f"   Método: SendGrid API")
                enviar_email_sendgrid(email_cliente, pdf_paths, resumen_opcion1, num_opciones)
                email_enviado = True
                print(f"✅ Email enviado exitosamente via SendGrid a {email_cliente}")
            else:
                # Fallback a SMTP (local/desarrollo)
                print(f"   Método: SMTP (PrivateEmail/Gmail)")
                print(f"   ⚠️ SendGrid no configurado, usando SMTP como fallback")
                enviar_email_smtp(pdf_paths, email_cliente, resumen_opcion1, num_opciones)
                email_enviado = True
                print(f"✅ Email enviado exitosamente via SMTP a {email_cliente}")
        except Exception as e:
            email_error = str(e)
            print(f"⚠️ Warning - No se pudo enviar email: {e}")
            print(f"⚠️ Cotización generada correctamente. Email pendiente por configuración.")
            # NO lanzar excepción - permitir que la cotización continúe
        finally:
            # Limpiar archivos temporales
            print("\n🧹 Limpiando archivos temporales...")
            for pdf_path in pdf_paths:
                if pdf_path and os.path.exists(pdf_path):
                    try:
                        os.remove(pdf_path)
                        print(f"   🗑️  {os.path.basename(pdf_path)}")
                    except Exception as e:
                        print(f"   ⚠️ Error eliminando: {e}")
            
            for pptx_path in pptx_paths:
                if pptx_path and os.path.exists(pptx_path):
                    try:
                        os.remove(pptx_path)
                        print(f"   🗑️  {os.path.basename(pptx_path)}")
                    except Exception as e:
                        print(f"   ⚠️ Error eliminando: {e}")
        
        # Determinar mensaje de respuesta según estado del email
        if email_enviado:
            mensaje_respuesta = f"✅ Cotización enviada exitosamente a {email_cliente}"
            
            # Actualizar estado en PostgreSQL
            try:
                from models import get_db_session, Cotizacion
                session = get_db_session()
                cotizacion = session.query(Cotizacion).filter(Cotizacion.id == cotizacion_id).first()
                if cotizacion:
                    cotizacion.email_enviado = True
                    cotizacion.fecha_envio_email = now_colombia()
                    session.commit()
                    print(f"💾 Estado actualizado: email_enviado = True")
                session.close()
            except Exception as e:
                print(f"⚠️ Error actualizando estado (no crítico): {e}")
        else:
            mensaje_respuesta = f"✅ Cotización generada correctamente. Email pendiente por configuración de servidor SMTP."
        
        return JSONResponse({
            "status": "success",
            "mensaje": mensaje_respuesta,
            "emailEnviado": email_enviado,
            "emailError": email_error if not email_enviado else None,
            "numOpciones": num_opciones
        })
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Error en enviar_cotizacion: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"Error procesando envío: {e}")


# ========================================
# 🚀 EJECUCIÓN
# ========================================
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(
        app,
        host="0.0.0.0",
        port=int(os.getenv("PORT", "5000")),
        log_level="info"
    )# Trigger Railway redeploy - Thu Dec  4 10:42:41 -05 2025
