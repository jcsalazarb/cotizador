#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Compara placeholders definidos vs encontrados en el template
"""
import re
from pptx import Presentation

TEMPLATE_PATH = "../Template/Template-PreCotizacion.pptx"

# Placeholders definidos en build_placeholders()
DEFINIDOS = [
    "{{COT_ID}}", "{{FECHA}}",
    "{{NOMBRE}}", "{{EMAIL}}", "{{TELEFONO}}", "{{CIUDAD}}", "{{DIRECC}}", "{{NIC}}",
    "{{CONSUMO}}", "{{FACTURA}}", "{{VAL_KWH}}",
    "{{VIVIENDA}}", "{{SIS_ELEC}}", "{{TIPO_FV}}",
    "{{N_PANEL}}", "{{M_PANEL}}", "{{N_INVER}}", "{{M_INVER}}", "{{N_BATER}}", "{{M_BATER}}",
    "{{CAP_KW}}", "{{GEN_MES}}",
    "{{INVER}}", "{{SUBTOT}}", "{{AHO_MES}}", "{{RETORNO}}", "{{PORC_PR}}",
    "{{ACUM_GEN}}", "{{ACUM_DED}}", "{{ACUM_DEP}}", "{{TOT_ACUM}}",
    "{{NPISOS}}", "{{HSPC}}", "{{AREA}}", "{{PCTDIA}}"
]

def find_placeholders_in_template():
    """Encuentra todos los placeholders en el template usando el método mejorado"""
    prs = Presentation(TEMPLATE_PATH)
    found = set()
    
    for slide in prs.slides:
        for shape in slide.shapes:
            # Método 1: text_frame
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        matches = re.findall(r'\{\{[^}]+\}\}', run.text)
                        found.update(matches)
            
            # Método 2: tablas
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        matches = re.findall(r'\{\{[^}]+\}\}', cell.text)
                        found.update(matches)
            
            # Método 3: Buscar en XML (SmartArt y objetos complejos)
            try:
                namespaces = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                for text_elem in shape.element.findall('.//a:t', namespaces):
                    if text_elem.text:
                        matches = re.findall(r'\{\{[^}]+\}\}', text_elem.text)
                        found.update(matches)
            except:
                pass
    
    return found

def main():
    print("=" * 70)
    print("📋 COMPARACIÓN DE PLACEHOLDERS")
    print("=" * 70)
    
    encontrados = find_placeholders_in_template()
    definidos_set = set(DEFINIDOS)
    
    # Placeholders definidos pero NO en template
    faltantes = definidos_set - encontrados
    # Placeholders en template pero NO definidos
    extras = encontrados - definidos_set
    # Placeholders que coinciden
    coinciden = definidos_set & encontrados
    
    print(f"\n✅ PLACEHOLDERS QUE COINCIDEN ({len(coinciden)}):")
    for ph in sorted(coinciden):
        print(f"   {ph}")
    
    print(f"\n⚠️  PLACEHOLDERS DEFINIDOS PERO NO EN TEMPLATE ({len(faltantes)}):")
    print("   (Estos no se reemplazarán porque no están en el PPTX)")
    for ph in sorted(faltantes):
        print(f"   {ph}")
    
    if extras:
        print(f"\n⚠️  PLACEHOLDERS EN TEMPLATE PERO NO DEFINIDOS ({len(extras)}):")
        print("   (Estos no se reemplazarán porque no están en el código)")
        for ph in sorted(extras):
            print(f"   {ph}")
    
    print("\n" + "=" * 70)
    print(f"📊 RESUMEN:")
    print(f"   Definidos en código: {len(definidos_set)}")
    print(f"   Encontrados en template: {len(encontrados)}")
    print(f"   Coinciden: {len(coinciden)}")
    print(f"   Faltantes en template: {len(faltantes)}")
    print(f"   Extras en template: {len(extras)}")
    print("=" * 70)

if __name__ == "__main__":
    main()
