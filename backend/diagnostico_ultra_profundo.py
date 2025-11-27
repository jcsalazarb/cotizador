#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Diagnóstico ULTRA PROFUNDO - Busca en TODO el PPTX incluyendo masters, layouts, notes
"""
import os
import re
from pptx import Presentation

TEMPLATE_PATH = "../Template/Template-PreCotizacion.pptx"

def find_all_placeholders_in_xml(element):
    """Busca todos los placeholders en cualquier elemento XML"""
    placeholders = []
    try:
        # Buscar en TODO el árbol XML
        xml_str = str(element.xml) if hasattr(element, 'xml') else str(element)
        matches = re.findall(r'\{\{[^}]+\}\}', xml_str)
        return list(set(matches))
    except:
        return []

def analyze_ultra_deep():
    """Análisis ultra profundo del template"""
    if not os.path.exists(TEMPLATE_PATH):
        print(f"❌ Template no encontrado: {TEMPLATE_PATH}")
        return
    
    print("=" * 70)
    print("🔬 DIAGNÓSTICO ULTRA PROFUNDO DEL TEMPLATE PPTX")
    print("=" * 70)
    
    prs = Presentation(TEMPLATE_PATH)
    all_placeholders = set()
    locations = {}
    
    # 1. Buscar en todas las diapositivas
    print(f"\n📄 DIAPOSITIVAS ({len(prs.slides)}):")
    for slide_idx, slide in enumerate(prs.slides):
        slide_placeholders = find_all_placeholders_in_xml(slide)
        all_placeholders.update(slide_placeholders)
        if slide_placeholders:
            print(f"   Diapositiva {slide_idx + 1}: {len(slide_placeholders)} placeholders")
            for ph in slide_placeholders:
                if ph not in locations:
                    locations[ph] = []
                locations[ph].append(f"Diapositiva {slide_idx + 1}")
    
    # 2. Buscar en el Slide Master
    try:
        print(f"\n🎨 SLIDE MASTER:")
        if hasattr(prs, 'slide_master') or hasattr(prs, 'slide_masters'):
            masters = prs.slide_masters if hasattr(prs, 'slide_masters') else [prs.slide_master]
            for master_idx, master in enumerate(masters):
                master_placeholders = find_all_placeholders_in_xml(master)
                all_placeholders.update(master_placeholders)
                if master_placeholders:
                    print(f"   Master {master_idx + 1}: {len(master_placeholders)} placeholders")
                    for ph in master_placeholders:
                        if ph not in locations:
                            locations[ph] = []
                        locations[ph].append(f"Master {master_idx + 1}")
    except Exception as e:
        print(f"   (No se pudo acceder a Slide Master: {e})")
    
    # 3. Buscar en Layouts
    try:
        print(f"\n📐 SLIDE LAYOUTS:")
        layout_count = 0
        for slide_idx, slide in enumerate(prs.slides):
            if hasattr(slide, 'slide_layout'):
                layout_placeholders = find_all_placeholders_in_xml(slide.slide_layout)
                all_placeholders.update(layout_placeholders)
                if layout_placeholders:
                    layout_count += 1
                    for ph in layout_placeholders:
                        if ph not in locations:
                            locations[ph] = []
                        locations[ph].append(f"Layout de Diapositiva {slide_idx + 1}")
        print(f"   Encontrados placeholders en {layout_count} layouts")
    except Exception as e:
        print(f"   (No se pudo acceder a Layouts: {e})")
    
    # 4. Buscar en Notes
    try:
        print(f"\n📝 NOTES:")
        notes_count = 0
        for slide_idx, slide in enumerate(prs.slides):
            if hasattr(slide, 'notes_slide'):
                notes_placeholders = find_all_placeholders_in_xml(slide.notes_slide)
                all_placeholders.update(notes_placeholders)
                if notes_placeholders:
                    notes_count += 1
                    for ph in notes_placeholders:
                        if ph not in locations:
                            locations[ph] = []
                        locations[ph].append(f"Notes de Diapositiva {slide_idx + 1}")
        print(f"   Encontrados placeholders en {notes_count} notes")
    except Exception as e:
        print(f"   (No se pudo acceder a Notes: {e})")
    
    # 5. Resumen
    print(f"\n" + "=" * 70)
    print(f"📊 RESUMEN:")
    print(f"   Total de placeholders únicos encontrados: {len(all_placeholders)}")
    print("=" * 70)
    
    if all_placeholders:
        print(f"\n🏷️  LISTA COMPLETA DE PLACEHOLDERS ({len(all_placeholders)}):\n")
        for ph in sorted(all_placeholders):
            print(f"   {ph}")
            if ph in locations:
                for loc in locations[ph][:3]:  # Mostrar máximo 3 ubicaciones
                    print(f"      └─ {loc}")
                if len(locations[ph]) > 3:
                    print(f"      └─ ... y {len(locations[ph]) - 3} más")
    else:
        print("\n   ⚠️  No se encontraron placeholders")
    
    print("\n" + "=" * 70)
    
    # Comparar con los definidos
    DEFINIDOS = [
        "{{COT_ID}}", "{{FECHA}}",
        "{{NOMBRE}}", "{{EMAIL}}", "{{TELEFONO}}", "{{CIUDAD}}", "{{DIRECC}}", "{{NIC}}",
        "{{CONSUMO}}", "{{FACTURA}}", "{{VAL_KWH}}",
        "{{VIVIENDA}}", "{{SIS_ELEC}}", "{{TIPO_FV}}",
        "{{N_PANEL}}", "{{M_PANEL}}", "{{N_INVER}}", "{{M_INVER}}", "{{N_BATER}}", "{{M_BATER}}",
        "{{CAP_KW}}", "{{GEN_MES}}",
        "{{INVER}}", "{{SUBTOT}}", "{{AHO_MES}}", "{{RETORNO}}", "{{PORC_PR}}",
        "{{ACUM_GEN}}", "{{ACUM_DED}}", "{{ACUM_DEP}}", "{{TOT_ACUM}}",
        "{{NPISOS}}", "{{HSPC}}", "{{AREA}}", "{{PCTDIA}}"
    ]
    
    definidos_set = set(DEFINIDOS)
    encontrados_set = all_placeholders
    
    faltantes = definidos_set - encontrados_set
    extras = encontrados_set - definidos_set
    
    if faltantes:
        print(f"\n⚠️  PLACEHOLDERS DEFINIDOS PERO NO ENCONTRADOS ({len(faltantes)}):")
        for ph in sorted(faltantes):
            print(f"   {ph}")
    
    if extras:
        print(f"\n⚠️  PLACEHOLDERS ENCONTRADOS PERO NO DEFINIDOS ({len(extras)}):")
        for ph in sorted(extras):
            print(f"   {ph}")
    
    print("\n" + "=" * 70)

if __name__ == "__main__":
    analyze_ultra_deep()
